"""Generate a self-explanatory supervisor report PPT for test08.

The slides intentionally do not show source code. Each step explains:
input, function, core logic, output, and the corresponding debug image.
"""

from __future__ import annotations

import csv
import json
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
TEST_DIR = ROOT / "debug_output" / "test08"
REGION_DIR = TEST_DIR / "debug_region"
TILE_DIR = TEST_DIR / "debug_vis" / "tile_37_13_4992_14208"
OUT_PPT = ROOT / "docs" / "reports" / "CD34_test08_supervisor_report_pipeline_explained.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(248, 249, 250)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(91, 101, 116)
LINE = RGBColor(221, 227, 235)
DARK = RGBColor(17, 24, 39)
BLUE = RGBColor(40, 99, 175)
BLUE_LIGHT = RGBColor(231, 240, 255)
GREEN = RGBColor(42, 135, 92)
GREEN_LIGHT = RGBColor(229, 246, 238)
RED = RGBColor(160, 54, 45)
RED_LIGHT = RGBColor(252, 236, 232)
GOLD = RGBColor(178, 116, 34)
GOLD_LIGHT = RGBColor(255, 244, 226)
PURPLE = RGBColor(106, 81, 165)
PURPLE_LIGHT = RGBColor(241, 237, 250)

FONT = "Microsoft YaHei"


def load_json(path: Path) -> dict:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def load_area_summary() -> dict:
    summary = {}
    with (TEST_DIR / "DC2200155 A3 CD34_summary.csv").open("r", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            summary[row["metric"]] = float(row["value"])
    areas = []
    with (TEST_DIR / "DC2200155 A3 CD34_statistics.csv").open("r", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            areas.append(float(row["area_px2"]))
    areas.sort()
    return {
        "count": int(summary["count"]),
        "mean": summary["area_mean"],
        "std": summary["area_std"],
        "min": summary["area_min"],
        "max": summary["area_max"],
        "median": areas[len(areas) // 2] if areas else 0,
        "lt20": sum(a < 20 for a in areas),
        "lt100": sum(a < 100 for a in areas),
        "lt500": sum(a < 500 for a in areas),
    }


def prepare_image(path: Path, tempdir: Path, max_px: int = 1700) -> Path:
    image = Image.open(path).convert("RGB")
    image.thumbnail((max_px, max_px), Image.Resampling.LANCZOS)
    out = tempdir / f"{path.stem}.jpg"
    image.save(out, "JPEG", quality=88, optimize=True)
    return out


def add_image_fit(slide, path: Path, x, y, w, h, tempdir: Path):
    img_path = prepare_image(path, tempdir)
    iw, ih = Image.open(img_path).size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        pic_w = w
        pic_h = w / img_ratio
        pic_x = x
        pic_y = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * img_ratio
        pic_x = x + (w - pic_w) / 2
        pic_y = y
    slide.shapes.add_picture(str(img_path), pic_x, pic_y, width=pic_w, height=pic_h)


def set_font(paragraph, size=12, bold=False, color=TEXT):
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()


def textbox(slide, x, y, w, h, text="", size=12, bold=False,
            color=TEXT, align=None, valign=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    if valign is not None:
        frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    set_font(paragraph, size, bold, color)
    if align is not None:
        paragraph.alignment = align
    return box


def rounded_box(slide, x, y, w, h, fill=WHITE, line=LINE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def title(slide, main, sub=None):
    textbox(slide, Inches(0.55), Inches(0.25), Inches(12.15), Inches(0.45),
            main, 23, True, DARK)
    if sub:
        textbox(slide, Inches(0.58), Inches(0.73), Inches(12.0), Inches(0.32),
                sub, 10.5, False, MUTED)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.11),
        Inches(12.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def bullets(slide, x, y, w, h, items, size=11.2, color=TEXT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    for index, item in enumerate(items):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.space_after = Pt(4)
        paragraph.line_spacing = 1.06
        set_font(paragraph, size, False, color)
    return tb


def info_card(slide, x, y, w, h, header, items, accent=BLUE, fill=WHITE):
    rounded_box(slide, x, y, w, h, fill, LINE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.07), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    textbox(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.35),
            Inches(0.28), header, 12.5, True, accent)
    bullets(slide, x + Inches(0.2), y + Inches(0.52), w - Inches(0.35),
            h - Inches(0.62), items, 10.4)


def image_panel(slide, path, x, y, w, h, label, tempdir):
    rounded_box(slide, x, y, w, h + Inches(0.34), WHITE, LINE)
    add_image_fit(slide, path, x + Inches(0.08), y + Inches(0.08),
                  w - Inches(0.16), h - Inches(0.07), tempdir)
    textbox(slide, x + Inches(0.08), y + h + Inches(0.09),
            w - Inches(0.16), Inches(0.21), label, 9.2, False, MUTED,
            PP_ALIGN.CENTER)


def metric(slide, x, y, w, h, value, label, accent=BLUE, fill=WHITE):
    rounded_box(slide, x, y, w, h, fill, LINE)
    textbox(slide, x + Inches(0.08), y + Inches(0.12), w - Inches(0.16),
            Inches(0.34), value, 18, True, accent, PP_ALIGN.CENTER)
    textbox(slide, x + Inches(0.08), y + Inches(0.50), w - Inches(0.16),
            Inches(0.22), label, 9.2, False, MUTED, PP_ALIGN.CENTER)


def step_io(slide, y, input_text, logic_text, output_text):
    x0 = Inches(0.66)
    widths = [Inches(3.45), Inches(4.62), Inches(3.45)]
    labels = [("输入", input_text, BLUE, BLUE_LIGHT),
              ("处理逻辑", logic_text, GOLD, GOLD_LIGHT),
              ("输出", output_text, GREEN, GREEN_LIGHT)]
    for idx, (label, text, accent, fill) in enumerate(labels):
        x = x0 + sum(widths[:idx]) + Inches(0.25 * idx)
        rounded_box(slide, x, y, widths[idx], Inches(0.88), fill, LINE)
        textbox(slide, x + Inches(0.15), y + Inches(0.13), Inches(0.75),
                Inches(0.24), label, 10.8, True, accent)
        textbox(slide, x + Inches(0.9), y + Inches(0.12), widths[idx] - Inches(1.05),
                Inches(0.55), text, 9.7, False, TEXT)
        if idx < 2:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + widths[idx] + Inches(0.04),
                y + Inches(0.28),
                Inches(0.18),
                Inches(0.24),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()


def process_note(slide, x, y, w, h, title_text, body, accent=RED):
    rounded_box(slide, x, y, w, h, WHITE, LINE)
    textbox(slide, x + Inches(0.18), y + Inches(0.13), w - Inches(0.35),
            Inches(0.25), title_text, 12.4, True, accent)
    textbox(slide, x + Inches(0.18), y + Inches(0.47), w - Inches(0.35),
            h - Inches(0.55), body, 10.4, False, TEXT)


def footer(slide, text="样例：debug_output/test08 / tile_37_13_4992_14208"):
    textbox(slide, Inches(0.65), Inches(7.12), Inches(12.0), Inches(0.2),
            text, 8.4, False, MUTED, PP_ALIGN.RIGHT)


def flow_node(slide, x, y, w, h, header, input_text, output_text, accent):
    rounded_box(slide, x, y, w, h, WHITE, accent)
    textbox(slide, x + Inches(0.1), y + Inches(0.12), w - Inches(0.2),
            Inches(0.22), header, 10.6, True, accent, PP_ALIGN.CENTER)
    textbox(slide, x + Inches(0.14), y + Inches(0.44), w - Inches(0.28),
            Inches(0.28), f"入：{input_text}", 8.1, False, TEXT, PP_ALIGN.CENTER)
    textbox(slide, x + Inches(0.14), y + Inches(0.75), w - Inches(0.28),
            Inches(0.28), f"出：{output_text}", 8.1, False, TEXT, PP_ALIGN.CENTER)


def add_arrow(slide, x, y, w=0.28, h=0.22):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MUTED
    arrow.line.fill.background()


def build():
    metadata = load_json(REGION_DIR / "metadata.json")
    region = load_json(REGION_DIR / "region_summary.json")
    stitched = load_json(REGION_DIR / "stitched_deepliif_metadata.json")
    prompt = load_json(TILE_DIR / "step3_34_weighted_prompt_summary.json")
    sam = load_json(TILE_DIR / "step4_weighted_summary.json")
    merge = load_json(TILE_DIR / "step5_06_merge_filter_summary.json")
    summary = load_area_summary()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)

        # 1 Cover
        s = prs.slides.add_slide(blank)
        add_bg(s)
        textbox(s, Inches(0.7), Inches(0.72), Inches(9.2), Inches(0.62),
                "CD34 微血管识别项目当前进展", 30, True, DARK)
        textbox(s, Inches(0.73), Inches(1.45), Inches(8.6), Inches(0.36),
                "基于 test08 的完整流程说明：每一步的输入、功能、逻辑和输出", 15, False, MUTED)
        image_panel(s, REGION_DIR / "05_region_geojson_overlay.png",
                    Inches(7.3), Inches(0.72), Inches(5.1), Inches(5.72),
                    "当前最终效果：自动生成的微血管 GeoJSON 叠加到原图", tmp)
        metric(s, Inches(0.78), Inches(2.45), Inches(1.65), Inches(0.85),
               str(region["tile_count"]), "处理 tile", BLUE, BLUE_LIGHT)
        metric(s, Inches(2.68), Inches(2.45), Inches(1.65), Inches(0.85),
               str(summary["count"]), "最终区域", RED, RED_LIGHT)
        metric(s, Inches(4.58), Inches(2.45), Inches(1.65), Inches(0.85),
               str(region["accepted_stitch_matches_drawn"]), "跨 tile 合并", GREEN, GREEN_LIGHT)
        info_card(s, Inches(0.78), Inches(3.75), Inches(5.65), Inches(1.75),
                  "本次汇报读法", [
                      "先看最终结果，再用一个具体 tile 逐步解释处理过程。",
                      "每一步都说明输入是什么、为什么要做、怎么判断、输出是什么。",
                      "最后回到整个 ROI，说明多个 tile 如何拼接成最终 GeoJSON。",
                  ], BLUE)
        textbox(s, Inches(0.8), Inches(6.35), Inches(5.55), Inches(0.3),
                "样例 tile：tile_37_13_4992_14208", 13, True, RED,
                PP_ALIGN.CENTER)
        footer(s, "CD34 test08 supervisor report")

        # 2 final result
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "当前结果：pipeline 已能输出可复查的 CD34 微血管实例",
              "红色区域是自动识别出的微血管实例，结果已经导出为 GeoJSON，可以在 QuPath 中复核")
        image_panel(s, REGION_DIR / "01_region_original_mosaic.png",
                    Inches(0.65), Inches(1.45), Inches(5.2), Inches(4.85),
                    "输入 ROI 原图：含 tile 网格和目标区域", tmp)
        image_panel(s, REGION_DIR / "05_region_geojson_overlay.png",
                    Inches(6.2), Inches(1.45), Inches(5.2), Inches(4.85),
                    "输出结果：GeoJSON 实例轮廓叠加", tmp)
        metric(s, Inches(11.65), Inches(1.85), Inches(0.85), Inches(0.72),
               str(summary["count"]), "实例", RED, RED_LIGHT)
        metric(s, Inches(11.65), Inches(2.78), Inches(0.85), Inches(0.72),
               str(region["tile_merged_instances"]), "tile 内", GOLD, GOLD_LIGHT)
        metric(s, Inches(11.65), Inches(3.71), Inches(0.85), Inches(0.72),
               str(region["accepted_stitch_matches_drawn"]), "合并", GREEN, GREEN_LIGHT)
        process_note(s, Inches(0.85), Inches(6.55), Inches(11.25), Inches(0.45),
                     "一句话结论",
                     "目前已经跑通“WSI 区域 -> tile 处理 -> SAM2 分割 -> 跨 tile 拼接 -> GeoJSON 输出”的完整链路。", RED)
        footer(s)

        # 3 how to read
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "这张结果图怎么看", "先把图中的元素解释清楚，后面每一步都围绕这些元素展开")
        image_panel(s, REGION_DIR / "05_region_geojson_overlay.png",
                    Inches(0.75), Inches(1.45), Inches(6.2), Inches(5.1),
                    "最终 overlay 图", tmp)
        info_card(s, Inches(7.35), Inches(1.48), Inches(4.55), Inches(1.05),
                  "红色半透明区域", [
                      "最终导出的微血管实例。",
                      "红色边界是实例轮廓，内部是实例覆盖范围。",
                  ], RED, RED_LIGHT)
        info_card(s, Inches(7.35), Inches(2.78), Inches(4.55), Inches(1.05),
                  "蓝/黄色网格", [
                      "每个格子对应一个 tile 的位置。",
                      "用来检查 tile 边界附近是否断裂或重复。",
                  ], BLUE, BLUE_LIGHT)
        info_card(s, Inches(7.35), Inches(4.08), Inches(4.55), Inches(1.05),
                  "中间框线", [
                      "表示本次 debug region / ROI 的核心区域。",
                      "neighbor tile 只提供边界上下文，不一定全部作为最终统计区域。",
                  ], GREEN, GREEN_LIGHT)
        info_card(s, Inches(7.35), Inches(5.38), Inches(4.55), Inches(1.05),
                  "当前需要继续优化的地方", [
                      "有些小碎片也被保留下来。",
                      "下一步需要用面积、形态和置信度规则进一步过滤。",
                  ], GOLD, GOLD_LIGHT)
        footer(s)

        # 4 overall
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "整体流程：每一步都有明确输入和输出",
              "后面用同一个 tile 逐步展开，避免只讲抽象流程")
        nodes = [
            ("1 选 tile", "WSI + ROI", "待处理 tile"),
            ("2 DeepLIIF", "原始 tile", "Seg/Marker/DAPI"),
            ("3 构建 prompt", "中间通道 + 原图", "mask_input + points"),
            ("4 SAM2", "prompt + 原图", "raw mask"),
            ("5 后处理", "raw mask", "tile 实例"),
            ("6 拼接导出", "所有 tile 实例", "GeoJSON + 统计"),
        ]
        colors = [BLUE, GREEN, GOLD, RED, PURPLE, DARK]
        for i, (h, inp, out) in enumerate(nodes):
            x = Inches(0.66 + i * 2.08)
            flow_node(s, x, Inches(1.55), Inches(1.72), Inches(1.16),
                      h, inp, out, colors[i])
            if i < len(nodes) - 1:
                add_arrow(s, x + Inches(1.76), Inches(2.02), 0.25, 0.2)
        image_panel(s, REGION_DIR / "01_region_original_mosaic.png",
                    Inches(0.7), Inches(3.25), Inches(2.8), Inches(2.85),
                    "输入：ROI tiles", tmp)
        image_panel(s, REGION_DIR / "08_stitched_deepliif_seg.png",
                    Inches(3.85), Inches(3.25), Inches(2.8), Inches(2.85),
                    "中间：DeepLIIF Seg", tmp)
        image_panel(s, TILE_DIR / "step3_22_weighted_mask_input_256.png",
                    Inches(7.0), Inches(3.25), Inches(2.2), Inches(2.85),
                    "中间：SAM2 prompt", tmp)
        image_panel(s, REGION_DIR / "05_region_geojson_overlay.png",
                    Inches(9.55), Inches(3.25), Inches(2.8), Inches(2.85),
                    "输出：GeoJSON", tmp)
        footer(s)

        # 5 sample tile
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "后续用这个具体 tile 讲完整过程",
              "选择 tile_37_13_4992_14208，因为它完整保存了 step1 到 step7 的调试图")
        image_panel(s, TILE_DIR / "step1_original.png",
                    Inches(0.9), Inches(1.55), Inches(4.65), Inches(4.65),
                    "样例 tile 原图", tmp)
        info_card(s, Inches(6.1), Inches(1.65), Inches(5.6), Inches(1.2),
                  "为什么用单个 tile 展开讲", [
                      "导师可以看到一张原图如何逐步变成 prompt、SAM2 mask 和最终实例。",
                      "避免只看最终 ROI 图时，不知道每个中间环节做了什么。",
                  ], BLUE)
        info_card(s, Inches(6.1), Inches(3.12), Inches(5.6), Inches(1.2),
                  "这个 tile 的关键数字", [
                      f"Marker 自动阈值：{prompt['marker_thresh']}。",
                      f"SAM2 positive points：{sam['point_count']} 个。",
                      f"SAM2 最佳候选分数：{sam['best_score']:.4f}。",
                      f"后处理后 tile 内实例：{merge['final_regions']} 个。",
                  ], GREEN)
        process_note(s, Inches(6.1), Inches(4.75), Inches(5.6), Inches(1.1),
                     "讲解主线",
                     "原始 tile -> DeepLIIF 三个中间通道 -> weighted prompt -> SAM2 候选 mask -> 连通域实例 -> 全区域拼接。", RED)
        footer(s)

        # 6 step 1 input
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 1：输入是一张 512x512 的原始 CD34 IHC tile",
              "这一步的任务不是识别血管，而是把 WSI 中目标位置的图块取出来，作为后续模型的共同输入")
        step_io(s, Inches(1.24),
                "WSI 中某个 ROI 位置；包含 row/col 和全局坐标",
                "按 tile_size 和 overlap 从 WSI 读取图块；保留坐标信息",
                "RGB tile 图像；后续 DeepLIIF 和 SAM2 都使用它")
        image_panel(s, TILE_DIR / "step1_original.png",
                    Inches(0.9), Inches(2.45), Inches(4.35), Inches(4.0),
                    "原始 CD34 IHC tile", tmp)
        info_card(s, Inches(5.75), Inches(2.55), Inches(3.05), Inches(1.35),
                  "功能", [
                      "提供真实染色图像。",
                      "作为 DeepLIIF 的输入。",
                      "作为 SAM2 后续精细分割的原图背景。",
                  ], BLUE)
        info_card(s, Inches(9.05), Inches(2.55), Inches(3.05), Inches(1.35),
                  "为什么保留坐标", [
                      "单 tile 分割只得到局部坐标。",
                      "最终需要回到 WSI 全局坐标导出 GeoJSON。",
                      "跨 tile 拼接也依赖 row/col 信息。",
                  ], GREEN)
        process_note(s, Inches(5.75), Inches(4.45), Inches(6.35), Inches(1.35),
                     "这一页要传达的信息",
                     "pipeline 不是直接在整张 WSI 上做一次分割，而是先把 ROI 切成有重叠的 tile。每个 tile 独立推理，最后再拼回全局结果。", RED)
        footer(s)

        # 7 DeepLIIF
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 2：DeepLIIF 把原图转换成三个可用线索",
              "DeepLIIF 不直接输出最终微血管实例，它输出的是后续构建 SAM2 prompt 的中间证据")
        step_io(s, Inches(1.18),
                "原始 tile 图像",
                "用 DeepLIIF 生成结构、染色和核/腔隙相关通道",
                "Seg、Marker、DAPI 三张中间图")
        image_panel(s, TILE_DIR / "step1_original.png",
                    Inches(0.65), Inches(2.25), Inches(2.5), Inches(2.65),
                    "输入：原图", tmp)
        image_panel(s, TILE_DIR / "step2_01_deepliif_Seg.png",
                    Inches(3.45), Inches(2.25), Inches(2.5), Inches(2.65),
                    "输出 1：Seg", tmp)
        image_panel(s, TILE_DIR / "step2_02_deepliif_Marker.png",
                    Inches(6.25), Inches(2.25), Inches(2.5), Inches(2.65),
                    "输出 2：Marker", tmp)
        image_panel(s, TILE_DIR / "step2_03_deepliif_DAPI.png",
                    Inches(9.05), Inches(2.25), Inches(2.5), Inches(2.65),
                    "输出 3：DAPI", tmp)
        info_card(s, Inches(0.85), Inches(5.55), Inches(3.55), Inches(0.95),
                  "Seg 的作用", [
                      "提供组织结构和边界方向。",
                      "帮助判断哪些区域更像血管壁结构。",
                  ], BLUE_LIGHT, WHITE)
        info_card(s, Inches(4.65), Inches(5.55), Inches(3.55), Inches(0.95),
                  "Marker 的作用", [
                      "提供 CD34 阳性染色强度。",
                      "强阳性区域更应该成为 SAM2 的正提示。",
                  ], RED, WHITE)
        info_card(s, Inches(8.45), Inches(5.55), Inches(3.55), Inches(0.95),
                  "DAPI 的作用", [
                      "辅助识别无核的管腔区域。",
                      "避免管腔内部被错误当成背景断开。",
                  ], GREEN, WHITE)
        footer(s)

        # 8 Seg logic
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 3：Seg 通道提供“结构是否像血管”的线索",
              "核心思想：不是简单二值化，而是按颜色关系和强度给每个像素一个支持程度")
        step_io(s, Inches(1.16),
                "DeepLIIF Seg RGB 图像",
                "根据 R/B/G 关系判断阳性方向、阴性方向和强弱等级",
                "seg_logits：每个像素得到 -5 到 5 的支持分数")
        image_panel(s, TILE_DIR / "step2_01_deepliif_Seg.png",
                    Inches(0.75), Inches(2.25), Inches(3.0), Inches(3.0),
                    "Seg 图像", tmp)
        image_panel(s, TILE_DIR / "step2_03_seg_positive_r_intensity_curve.png",
                    Inches(4.1), Inches(2.25), Inches(4.0), Inches(2.2),
                    "Seg 红色强度分布", tmp)
        info_card(s, Inches(8.45), Inches(2.25), Inches(3.75), Inches(1.55),
                  "判断逻辑", [
                      "红色成分强，说明更接近阳性结构。",
                      "蓝色成分更强，说明偏阴性或非目标。",
                      "绿色过高的区域容易是背景或不可靠区域。",
                  ], BLUE)
        info_card(s, Inches(8.45), Inches(4.1), Inches(3.75), Inches(1.25),
                  "输出含义", [
                      "-5 表示强烈不支持。",
                      "0 表示不确定区域。",
                      "1 到 5 表示越来越强的阳性支持。",
                  ], GREEN)
        process_note(s, Inches(0.85), Inches(5.75), Inches(11.35), Inches(0.65),
                     "为什么这样做",
                     "SAM2 需要一个形状先验。Seg 通道负责告诉 SAM2：哪些位置在结构上更可能属于血管壁，哪些位置应该被压低。", RED)
        footer(s)

        # 9 Marker logic
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 4：Marker 通道提供“CD34 阳性强度”的线索",
              "核心思想：Marker 越强，越应该成为 SAM2 的正向提示；弱染色区域需要谨慎处理")
        step_io(s, Inches(1.16),
                "DeepLIIF Marker 图像",
                "自动估计染色阈值，并按强度映射为不同等级",
                "marker_logits：CD34 阳性强度支持图")
        image_panel(s, TILE_DIR / "step2_02_deepliif_Marker.png",
                    Inches(0.75), Inches(2.25), Inches(3.0), Inches(3.0),
                    "Marker 图像", tmp)
        image_panel(s, TILE_DIR / "step2_04_marker_nonzero_intensity_curve.png",
                    Inches(4.1), Inches(2.25), Inches(4.0), Inches(2.2),
                    "Marker 非零强度分布", tmp)
        info_card(s, Inches(8.45), Inches(2.25), Inches(3.75), Inches(1.45),
                  "本 tile 的阈值", [
                      f"自动阈值：{prompt['marker_thresh']}。",
                      "阈值来源：two-stage Multi-Otsu。",
                      f"Marker 阳性像素：{prompt['marker_positive_px']:,} px。",
                  ], RED)
        info_card(s, Inches(8.45), Inches(4.0), Inches(3.75), Inches(1.35),
                  "功能", [
                      "补充 Seg 中不明显但染色阳性的区域。",
                      "让强 CD34 染色在后续 prompt 中占更高权重。",
                  ], GREEN)
        process_note(s, Inches(0.85), Inches(5.75), Inches(11.35), Inches(0.65),
                     "为什么要自动阈值",
                     "不同 tile 的染色强度不完全一致，固定阈值容易过严或过松；自动阈值可以让每个 tile 根据自身强度分布调整。", BLUE)
        footer(s)

        # 10 fusion
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 5：Seg 和 Marker 融合成初始 weighted prompt",
              "核心思想：结构线索和染色线索互补，取二者中更强的支持作为初始提示")
        step_io(s, Inches(1.16),
                "seg_logits + marker_logits",
                "逐像素融合：如果任一通道强烈支持，就提高该位置的 prompt 权重",
                "pre-DAB 初始热图")
        image_panel(s, TILE_DIR / "step3_01_weighted_raw_heatmap.png",
                    Inches(0.8), Inches(2.2), Inches(3.35), Inches(3.05),
                    "融合后的初始热图", tmp)
        image_panel(s, TILE_DIR / "step3_05_weighted_dab_filtered_heatmap.png",
                    Inches(4.55), Inches(2.2), Inches(3.35), Inches(3.05),
                    "DAB 过滤后的热图", tmp)
        info_card(s, Inches(8.35), Inches(2.25), Inches(3.85), Inches(1.4),
                  "这一步的功能", [
                      "把两个 DeepLIIF 通道变成统一的 prompt 空间。",
                      "弱支持、不确定、强支持都保留下来，而不是直接二值化。",
                      "给下一步 DAB 过滤提供基础。",
                  ], GOLD)
        info_card(s, Inches(8.35), Inches(3.95), Inches(3.85), Inches(1.25),
                  "输出特点", [
                      "不是最终 mask。",
                      "是给 SAM2 使用的形状先验草稿。",
                      "还需要用原图真实 DAB 强度继续清洗。",
                  ], BLUE)
        process_note(s, Inches(0.9), Inches(5.8), Inches(11.3), Inches(0.58),
                     "重点",
                     "这个项目的关键不是让 DeepLIIF 直接分割血管，而是用 DeepLIIF 生成 SAM2 更容易理解的 weighted prompt。", RED)
        footer(s)

        # 11 DAB filter
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 6：用原图 DAB 强度过滤低可信 prompt",
              "核心思想：真正的 CD34 阳性应当在原始染色图中有 DAB 支持，低 DAB 的 prompt 更可能是噪声或伪影")
        step_io(s, Inches(1.16),
                "初始 prompt + 原始 RGB tile",
                "从原图估计 DAB 强度，移除 DAB 不足的 prompt 像素",
                "filtered prompt：更可信的阳性提示")
        image_panel(s, TILE_DIR / "step3_02_weighted_dab_intensity.png",
                    Inches(0.7), Inches(2.18), Inches(2.7), Inches(2.65),
                    "原图估计的 DAB 强度", tmp)
        image_panel(s, TILE_DIR / "step3_03_weighted_dab_intensity_keep_mask.png",
                    Inches(3.72), Inches(2.18), Inches(2.45), Inches(2.65),
                    "DAB 保留区域", tmp)
        image_panel(s, TILE_DIR / "step3_06_weighted_dab_filter_overlay.png",
                    Inches(6.5), Inches(2.18), Inches(2.7), Inches(2.65),
                    "DAB 过滤 overlay", tmp)
        info_card(s, Inches(9.55), Inches(2.25), Inches(2.75), Inches(1.25),
                  "本 tile 数字", [
                      f"DAB 阈值：{prompt['dab_min_intensity']}。",
                      f"候选 prompt：{prompt['dab_prompt_candidate_px']:,} px。",
                      f"移除：{prompt['dab_prompt_removed_px']:,} px。",
                  ], RED)
        info_card(s, Inches(9.55), Inches(3.8), Inches(2.75), Inches(1.05),
                  "这一步解决什么问题", [
                      "减少弱染色背景被误当作血管。",
                      "减少 DeepLIIF 中间图带来的假阳性。",
                  ], GREEN)
        process_note(s, Inches(0.85), Inches(5.55), Inches(11.45), Inches(0.85),
                     "注意",
                     "DAB 过滤不是简单删除所有低强度区域。对于可能的管腔和被结构包围的区域，后续还会通过 lumen 逻辑保护或补回。", GOLD)
        footer(s)

        # 12 strong DAB lumen
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 7：强 DAB 和管腔逻辑把真实血管补回来",
              "核心思想：既要过滤噪声，也不能把真实血管壁或管腔结构误删")
        step_io(s, Inches(1.16),
                "filtered prompt + DAB 强度 + DAPI/管腔候选",
                "强 DAB 区域升级权重；被血管壁包围的管腔区域增加 positive point 或弱填充",
                "最终 weighted prompt + positive points")
        image_panel(s, TILE_DIR / "step3_11_weighted_dab_support_heatmap.png",
                    Inches(0.65), Inches(2.15), Inches(2.55), Inches(2.55),
                    "强 DAB 支持热图", tmp)
        image_panel(s, TILE_DIR / "step3_20_weighted_final_heatmap.png",
                    Inches(3.55), Inches(2.15), Inches(2.55), Inches(2.55),
                    "最终 prompt 热图", tmp)
        image_panel(s, TILE_DIR / "step3_21_weighted_final_overlay.png",
                    Inches(6.45), Inches(2.15), Inches(2.55), Inches(2.55),
                    "最终 prompt 叠加原图", tmp)
        image_panel(s, TILE_DIR / "step3_33_weighted_positive_points.png",
                    Inches(9.35), Inches(2.15), Inches(2.55), Inches(2.55),
                    "positive points", tmp)
        info_card(s, Inches(0.85), Inches(5.35), Inches(3.45), Inches(1.0),
                  "补强", [
                      f"强 DAB 新增 prompt：{prompt['dab_prompt_added_px']:,} px。",
                      f"强 DAB 升级 prompt：{prompt['dab_prompt_upgraded_px']:,} px。",
                  ], GREEN)
        info_card(s, Inches(4.6), Inches(5.35), Inches(3.45), Inches(1.0),
                  "管腔", [
                      f"DAB lumen 接受数：{prompt['dab_lumen_accepted_count']}。",
                      f"DAPI lumen 接受数：{prompt['dapi_lumen_accepted_count']}。",
                      "这些点帮助 SAM2 保持管腔结构连续。",
                  ], BLUE)
        info_card(s, Inches(8.35), Inches(5.35), Inches(3.45), Inches(1.0),
                  "最终结果", [
                      f"最终非负 prompt 像素：{prompt['final_nonnegative_px']:,} px。",
                      f"总 positive points：{prompt['total_positive_point_count']} 个。",
                  ], RED)
        footer(s)

        # 13 SAM2 inputs
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 8：真正输入给 SAM2 的不是二值 mask",
              "SAM2 接收的是一张低分辨率 weighted mask_input，加上一组阳性点")
        step_io(s, Inches(1.16),
                "最终 prompt logits + positive points",
                "把 512x512 prompt 池化到 256x256，作为 SAM2 的 mask_input；阳性点作为额外约束",
                "SAM2 可理解的 mask_input 和 point prompts")
        image_panel(s, TILE_DIR / "step3_22_weighted_mask_input_256.png",
                    Inches(1.0), Inches(2.2), Inches(3.55), Inches(3.55),
                    "mask_input：1x256x256 weighted logits", tmp)
        image_panel(s, TILE_DIR / "step3_33_weighted_positive_points.png",
                    Inches(5.0), Inches(2.2), Inches(3.55), Inches(3.55),
                    "point prompts：强阳性点和管腔点", tmp)
        info_card(s, Inches(9.0), Inches(2.35), Inches(3.0), Inches(1.35),
                  "mask_input 含义", [
                      "-5 表示强烈排除。",
                      "0 表示不确定。",
                      "1 到 5 表示越来越强的支持。",
                  ], GOLD)
        info_card(s, Inches(9.0), Inches(4.0), Inches(3.0), Inches(1.25),
                  "point prompts 含义", [
                      "点的位置来自强阳性连通区域或管腔候选。",
                      f"本 tile 输入 SAM2 的点数：{sam['point_count']} 个。",
                  ], BLUE)
        process_note(s, Inches(1.0), Inches(6.35), Inches(11.0), Inches(0.45),
                     "本页重点",
                     "weighted prompt 的目的，是把病理染色线索转换成 SAM2 能使用的形状先验，而不是直接替代 SAM2。", RED)
        footer(s)

        # 14 SAM2 candidates
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 9：SAM2 生成多个候选 mask，选择分数最高的一个",
              "SAM2 根据原图、mask_input 和点提示生成候选结果；当前选择 score 最大的候选作为 raw mask")
        step_io(s, Inches(1.16),
                "原图 tile + mask_input + positive points",
                "SAM2 生成多个候选 mask，并为每个候选给出 score",
                "best raw mask")
        image_panel(s, TILE_DIR / "step4_weighted_candidate_0_score_0.9157.png",
                    Inches(0.75), Inches(2.2), Inches(3.15), Inches(3.15),
                    f"候选 0：score={sam['scores'][0]:.4f}", tmp)
        image_panel(s, TILE_DIR / "step4_weighted_candidate_1_score_0.6561.png",
                    Inches(4.25), Inches(2.2), Inches(3.15), Inches(3.15),
                    f"候选 1：score={sam['scores'][1]:.4f}", tmp)
        image_panel(s, TILE_DIR / "step4_weighted_candidate_2_score_0.9497.png",
                    Inches(7.75), Inches(2.2), Inches(3.15), Inches(3.15),
                    f"候选 2：score={sam['scores'][2]:.4f}，最终选择", tmp)
        metric(s, Inches(11.3), Inches(2.55), Inches(1.0), Inches(0.76),
               str(sam["best_idx"]), "best", RED, RED_LIGHT)
        metric(s, Inches(11.3), Inches(3.55), Inches(1.0), Inches(0.76),
               f"{sam['best_score']:.3f}", "score", GREEN, GREEN_LIGHT)
        process_note(s, Inches(0.95), Inches(5.85), Inches(11.15), Inches(0.7),
                     "这一页要传达的信息",
                     "SAM2 负责把前面的提示转成更贴近图像边界的 mask。weighted prompt 负责告诉 SAM2 大概哪里应该分割，SAM2 负责细化边界。", BLUE)
        footer(s)

        # 15 postprocess
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 10：SAM2 raw mask 需要拆成多个实例",
              "一个 SAM2 mask 可能包含多个不相连的血管段，因此需要连通域拆分，才能得到实例级结果")
        step_io(s, Inches(1.16),
                "SAM2 best raw mask",
                "做连通域拆分；每个独立连通区域分配一个实例编号；可按面积等规则过滤",
                "tile 内实例 mask")
        image_panel(s, TILE_DIR / "step5_02_merge_filter_sam2_raw_overlay.png",
                    Inches(0.75), Inches(2.18), Inches(3.1), Inches(3.05),
                    "SAM2 raw mask overlay", tmp)
        image_panel(s, TILE_DIR / "step5_05_merge_filter_merged_overlay.png",
                    Inches(4.25), Inches(2.18), Inches(3.1), Inches(3.05),
                    "拆分后的实例 overlay", tmp)
        image_panel(s, TILE_DIR / "step7_sam2_merge_diff.png",
                    Inches(7.75), Inches(2.18), Inches(3.1), Inches(3.05),
                    "raw vs merged diff", tmp)
        metric(s, Inches(11.25), Inches(2.55), Inches(1.05), Inches(0.72),
               str(merge["total_connected_components"]), "连通域", RED, RED_LIGHT)
        metric(s, Inches(11.25), Inches(3.45), Inches(1.05), Inches(0.72),
               str(merge["final_regions"]), "保留", GREEN, GREEN_LIGHT)
        info_card(s, Inches(0.9), Inches(5.75), Inches(3.65), Inches(0.9),
                  "为什么要拆分", [
                      "SAM2 可能把多个血管段放进同一个 mask。",
                      "统计时需要每个不相连区域作为独立实例。",
                  ], RED)
        info_card(s, Inches(4.85), Inches(5.75), Inches(3.65), Inches(0.9),
                  "当前状态", [
                      "本 tile raw mask 拆成 56 个区域。",
                      "当前小面积过滤较宽松，所以小碎片也会保留。",
                  ], GOLD)
        info_card(s, Inches(8.8), Inches(5.75), Inches(3.25), Inches(0.9),
                  "后续优化", [
                      "增加面积、形态和置信度联合过滤。",
                      "减少微小碎片进入最终统计。",
                  ], GREEN)
        footer(s)

        # 16 stitching
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 11：多个 tile 拼接时要避免 overlap 重复计数",
              "center-valid 只导出每个 tile 的中心有效区；overlap 区域用于判断跨 tile 是否属于同一血管")
        step_io(s, Inches(1.16),
                "每个 tile 的实例 mask + tile 坐标",
                "中心有效区导出；重叠区域做身份匹配；跨 tile 相同血管合并",
                "全 ROI 的实例集合")
        image_panel(s, REGION_DIR / "03_region_tile_merged_mosaic.png",
                    Inches(0.65), Inches(2.05), Inches(3.35), Inches(3.25),
                    "tile 内实例拼到 ROI", tmp)
        image_panel(s, REGION_DIR / "07_region_overlap_matches.png",
                    Inches(4.35), Inches(2.05), Inches(3.35), Inches(3.25),
                    "跨 tile overlap 匹配", tmp)
        image_panel(s, REGION_DIR / "06_region_tile_vs_geojson_diff.png",
                    Inches(8.05), Inches(2.05), Inches(3.35), Inches(3.25),
                    "tile 结果与 GeoJSON 差异", tmp)
        metric(s, Inches(0.95), Inches(5.75), Inches(1.65), Inches(0.68),
               str(region["tile_merged_instances"]), "tile 内实例", GOLD, GOLD_LIGHT)
        metric(s, Inches(3.0), Inches(5.75), Inches(1.65), Inches(0.68),
               str(region["final_geojson_features"]), "最终实例", RED, RED_LIGHT)
        metric(s, Inches(5.05), Inches(5.75), Inches(1.65), Inches(0.68),
               str(region["accepted_stitch_matches_drawn"]), "合并匹配", GREEN, GREEN_LIGHT)
        info_card(s, Inches(7.25), Inches(5.55), Inches(4.7), Inches(1.0),
                  "拼接逻辑的作用", [
                      "减少 tile overlap 带来的重复导出。",
                      "保留 overlap 证据判断跨边界血管是否应合并。",
                      "把局部坐标统一到全局 ROI/WSI 坐标。",
                  ], BLUE)
        footer(s)

        # 17 GeoJSON
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "Step 12：最终输出 GeoJSON 和统计表",
              "最后一步把所有实例轮廓转成 WSI 全局坐标，形成可复查、可统计的结果文件")
        step_io(s, Inches(1.16),
                "全 ROI 实例集合 + tile 全局坐标 + WSI 分辨率",
                "轮廓提取、跨 tile 合并、多边形导出、面积统计",
                "GeoJSON、summary.csv、statistics.csv、面积分布图")
        image_panel(s, REGION_DIR / "05_region_geojson_overlay.png",
                    Inches(0.7), Inches(2.1), Inches(4.25), Inches(3.3),
                    "最终 GeoJSON overlay", tmp)
        image_panel(s, TEST_DIR / "DC2200155 A3 CD34_area_histogram.png",
                    Inches(5.25), Inches(2.1), Inches(6.55), Inches(2.75),
                    "实例面积分布", tmp)
        metric(s, Inches(5.45), Inches(5.25), Inches(1.3), Inches(0.68),
               str(summary["count"]), "count", RED, RED_LIGHT)
        metric(s, Inches(7.0), Inches(5.25), Inches(1.3), Inches(0.68),
               f"{summary['mean']:.0f}", "mean", GOLD, GOLD_LIGHT)
        metric(s, Inches(8.55), Inches(5.25), Inches(1.3), Inches(0.68),
               f"{summary['median']:.0f}", "median", BLUE, BLUE_LIGHT)
        metric(s, Inches(10.1), Inches(5.25), Inches(1.3), Inches(0.68),
               f"{summary['max']:.0f}", "max", GREEN, GREEN_LIGHT)
        process_note(s, Inches(0.95), Inches(6.25), Inches(10.9), Inches(0.58),
                     "当前结果的一个重要观察",
                     f"面积分布显示小碎片偏多：<20 px2 的区域有 {summary['lt20']} 个，<100 px2 的区域有 {summary['lt100']} 个。后续需要更严格的过滤规则。", RED)
        footer(s)

        # 18 summary
        s = prs.slides.add_slide(blank)
        add_bg(s)
        title(s, "阶段总结：主流程已跑通，下一步重点是结果质量优化",
              "这页可以作为汇报结束页，直接回答“现在做到哪了、问题是什么、下一步做什么”")
        info_card(s, Inches(0.85), Inches(1.55), Inches(3.55), Inches(4.35),
                  "已经完成", [
                      "ROI tile 读取和坐标记录。",
                      "DeepLIIF 生成 Seg / Marker / DAPI 中间通道。",
                      "基于 Seg、Marker、DAB、DAPI 的 weighted prompt。",
                      "SAM2 分割和候选选择。",
                      "连通域拆分、center-valid 拼接和 GeoJSON 输出。",
                  ], GREEN, GREEN_LIGHT)
        info_card(s, Inches(4.9), Inches(1.55), Inches(3.55), Inches(4.35),
                  "当前问题", [
                      "部分小碎片会进入最终统计。",
                      "有些区域可能存在边界过宽或粘连。",
                      "目前是局部 ROI 验证，还需要扩展到更多 ROI 和 WSI。",
                      "缺少与人工标注的定量对照。",
                  ], RED, RED_LIGHT)
        info_card(s, Inches(8.95), Inches(1.55), Inches(3.55), Inches(4.35),
                  "下一步计划", [
                      "加入面积、形态、边界支持度和 score 的联合过滤。",
                      "整理人工复核样本，用召回、误检和边界质量评估。",
                      "在多 ROI / 多 WSI 上复测稳定性。",
                      "把当前 debug 报告流程固定下来，方便后续参数迭代。",
                  ], BLUE, BLUE_LIGHT)
        process_note(s, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.48),
                     "汇报结论",
                     "当前项目已经从“模型能否跑通”推进到“输出质量是否足够稳定”的阶段，核心优化点是小碎片过滤和多样本验证。", RED)
        footer(s, "End")

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPT)
    return OUT_PPT


if __name__ == "__main__":
    print(build())
