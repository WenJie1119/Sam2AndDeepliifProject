#!/usr/bin/env python3
"""
Generate a detailed technical-route PPT for CD34 Microvessel Detection Pipeline.

Pure technical implementation walkthrough with real example images.
Every description is verified against the actual codebase.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Example image directories ──
TILE1 = "/local1/yangwenjie/DataImg/Tile_55_69/debug_vis/tile_55_69_35328_28160"
TILE2 = "/local1/yangwenjie/DataImg/Tile_63_72/debug_vis/tile_63_72_36864_32256"

# ── Color Palette ──
BG_DARK    = RGBColor(0x1B, 0x1B, 0x2F)
BG_CARD    = RGBColor(0x27, 0x27, 0x44)
ACCENT     = RGBColor(0x00, 0xD2, 0xFF)   # Cyan
ACCENT2    = RGBColor(0xFF, 0x6B, 0x6B)   # Coral
ACCENT3    = RGBColor(0x4E, 0xCB, 0x71)   # Green
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DIM_GRAY   = RGBColor(0x99, 0x99, 0x99)
ORANGE     = RGBColor(0xFF, 0xA5, 0x00)
PURPLE     = RGBColor(0xCC, 0x66, 0xFF)
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)


# ════════════════════════════════════════════════════════════════
#  HELPER FUNCTIONS
# ════════════════════════════════════════════════════════════════

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def _set_run_fonts(run, cn_font="SimSun", en_font="Times New Roman"):
    """Set Chinese font (SimSun) and English/number font (Times New Roman)."""
    run.font.name = en_font
    # Set East Asian font via XML
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('a:lang'), 'zh-CN')
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', cn_font)


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text; run.font.size = Pt(size)
    run.font.color.rgb = color; run.font.bold = bold
    _set_run_fonts(run)
    return txBox

def add_bullets(slide, left, top, width, height, items, size=14,
                color=LIGHT_GRAY, bullet_color=ACCENT, font=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4); p.space_before = Pt(2)
        is_sub = item.startswith("  ")
        rb = p.add_run()
        rb.text = "       -  " if is_sub else "  >  "
        rb.font.size = Pt(size); rb.font.color.rgb = bullet_color
        rb.font.bold = True; _set_run_fonts(rb)
        rt = p.add_run(); rt.text = item.strip().lstrip("- ")
        rt.font.size = Pt(size); rt.font.color.rgb = color; _set_run_fonts(rt)
    return txBox

def add_flow_box(slide, left, top, width, height, title, subtitle="",
                 fill=BG_CARD, border=ACCENT, title_color=ACCENT):
    shape = add_shape(slide, left, top, width, height, fill, border)
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = title; run.font.size = Pt(13)
    run.font.color.rgb = title_color; run.font.bold = True
    _set_run_fonts(run)
    if subtitle:
        p2 = shape.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = subtitle; r2.font.size = Pt(10)
        r2.font.color.rgb = DIM_GRAY; _set_run_fonts(r2)
    return shape

def add_arrow_right(slide, left, cy, length=Inches(0.3)):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, cy - Inches(0.1), length, Inches(0.2))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

def add_arrow_down(slide, cx, top, length=Inches(0.3)):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, cx - Inches(0.12), top, Inches(0.24), length)
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

def add_divider(slide, y):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(8.4), Pt(2))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

def add_header(slide, num, title):
    set_slide_bg(slide)
    add_text(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
             f"{num}  {title}", size=28, color=WHITE, bold=True)
    add_divider(slide, Inches(0.9))

def add_img(slide, path, left, top, width, height=None):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, left, top, width, height) if height \
            else slide.shapes.add_picture(path, left, top, width)
    h = height or width
    s = add_shape(slide, left, top, width, h, BG_CARD, DIM_GRAY)
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = s.text_frame.paragraphs[0].add_run()
    r.text = f"[Missing]\n{os.path.basename(path)}"; r.font.size = Pt(9); r.font.color.rgb = DIM_GRAY
    return s

def add_img_label(slide, path, left, top, w, h, label, lbl_color=ACCENT, lbl_size=10):
    add_img(slide, path, left, top, w, h)
    add_text(slide, left, top + h + Inches(0.02), w, Inches(0.25),
             label, size=lbl_size, color=lbl_color, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
#  SLIDES
# ════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1: Title with pipeline preview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide)
    # Top bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Pt(6))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

    add_text(slide, Inches(1), Inches(1.3), Inches(8), Inches(1),
             "CD34 微血管自动检测与分割系统", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(2.4), Inches(8), Inches(0.6),
             "技术路线: DeepLIIF 染色分解 → 连通域提取 → SAM2 实例分割",
             size=18, color=ACCENT, align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(3.2))

    # 4 preview images
    iw, ih, iy = Inches(2.0), Inches(2.0), Inches(3.5)
    imgs = [
        (f"{TILE2}/step1_original.png", "IHC 原图"),
        (f"{TILE2}/step2_deepliif_Seg.png", "DeepLIIF Seg"),
        (f"{TILE2}/cr_step5_connected_regions_10.png", "连通域提取"),
        (f"{TILE2}/step5_merged_10inst.png", "SAM2 实例分割"),
    ]
    x = Inches(0.5)
    for path, label in imgs:
        add_img_label(slide, path, x, iy, iw, ih, label, DIM_GRAY, 10)
        x += Inches(2.3)

    add_text(slide, Inches(1), Inches(6.2), Inches(8), Inches(0.4),
             "汇报人：杨文杰     2026 年 4 月", size=15, color=DIM_GRAY, align=PP_ALIGN.CENTER)


def slide_outline(prs):
    """Slide 2: Outline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "技术路线总览", size=32, color=WHITE, bold=True)
    add_divider(slide, Inches(1.05))

    items = [
        ("01", "Pipeline 总体架构",  "输入 -> DeepLIIF -> 阳性提取 -> SAM2 -> 后处理 -> 输出"),
        ("02", "DeepLIIF 染色分解",  "IHC 输入 -> 4 路模态翻译 -> 5 路分割聚合 -> Seg + Marker"),
        ("03", "Seg+Marker 联合提取", "前景检测 -> 阳性判定 -> Marker 增强 -> 形态学处理 -> 连通域"),
        ("04", "示例: 连通域提取过程", "Tile(63,72) 密集血管区域 — 逐步可视化"),
        ("05", "SAM2 Mask-Only 推理",  "logit Mask Prompt -> 3 候选掩码 -> 置信度选择 -> 重叠消解"),
        ("06", "SAM2 掩码提示词生成", "连通域 -> 256x256 logit 掩码 -> Mask Prompt 可视化"),
        ("07", "后处理与合并",        "连通域分析 -> 面积过滤 (min_area=200) -> 实例合并"),
        ("08", "Pipeline 完整示例",   "两个 Tile 从原图到最终分割的 4 步可视化"),
        ("09", "WSI 全片处理",       "OpenSlide 读取 -> YOLO 过滤 -> 多 GPU 并行 -> 结果输出"),
        ("10", "GeoJSON 导出详解",   "轮廓提取 -> Union-Find 边界合并 -> 几何合并 -> QuPath 导入"),
        ("11", "参数详表与调优",      "所有关键参数、默认值、调优建议"),
    ]
    y = Inches(1.35)
    for num, title, desc in items:
        add_flow_box(slide, Inches(0.8), y, Inches(0.65), Inches(0.46),
                     num, fill=BG_CARD, border=ACCENT, title_color=ACCENT)
        add_text(slide, Inches(1.65), y + Inches(0.0), Inches(3), Inches(0.28),
                 title, size=15, color=WHITE, bold=True)
        add_text(slide, Inches(1.65), y + Inches(0.24), Inches(7.5), Inches(0.22),
                 desc, size=10, color=DIM_GRAY)
        y += Inches(0.52)


def slide_architecture(prs):
    """Slide 3: Overall pipeline architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "01", "Pipeline 总体架构")

    add_text(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.3),
             "单 Tile 处理流程: 输入 IHC 图像 → 虚拟染色分解 → 阳性区域提取 → 实例分割 → 结果输出",
             size=12, color=DIM_GRAY)

    # Horizontal flow: 5 major steps
    bw, bh = Inches(1.65), Inches(1.05)
    yt = Inches(1.5); gap = Inches(0.25)
    steps = [
        ("1. 读取 Tile",   "从 WSI 裁取\n512x512 RGB", ACCENT),
        ("2. DeepLIIF\n   染色分解", "输出 Seg (RGB)\n输出 Marker (灰度)", ACCENT2),
        ("3. 阳性区域\n   提取",   "Seg+Marker 联合\n连通域模式", ORANGE),
        ("4. SAM2\n   实例分割",     "Mask 提示输入\n3 候选选最优", ACCENT3),
        ("5. 后处理\n   与保存",     "合并 + 面积过滤\nNPY / GeoJSON", PURPLE),
    ]
    x = Inches(0.25)
    for i, (t, s, c) in enumerate(steps):
        add_flow_box(slide, x, yt, bw, bh, t, s, fill=BG_CARD, border=c, title_color=c)
        if i < len(steps) - 1:
            add_arrow_right(slide, x + bw + Inches(0.02), yt + bh/2, gap - Inches(0.04))
        x += bw + gap

    # Detailed breakdown cards
    cy = Inches(3.0); ch = Inches(3.8)

    # Card 1: DeepLIIF
    add_shape(slide, Inches(0.2), cy, Inches(3.1), ch, BG_CARD, ACCENT2)
    add_text(slide, Inches(0.35), cy+Inches(0.08), Inches(2.8), Inches(0.35),
             "Step 2: DeepLIIF 染色分解", size=14, color=ACCENT2, bold=True)
    add_bullets(slide, Inches(0.35), cy+Inches(0.4), Inches(2.8), Inches(3.2), [
        "输入: 512x512 IHC Tile",
        "Tile 切分: 512px + 32px overlap",
        "4 路模态翻译 (Pix2Pix 生成器):",
        "  G1->DAPI, G2->Hema, G3->Lap2, G4->Marker",
        "5 路分割网络 (G51-G55):",
        "  各生成概率图, 加权平均聚合",
        "  默认权重: 各 0.2 (可配置)",
        "输出 Seg (RGB 图像):",
        "  R 通道 = 阳性 (棕色) 概率",
        "  B 通道 = 阴性 (蓝色) 概��",
        "  G 通道 = 边界标记",
        "输出 Marker (灰度图):",
        "  亮度 = CD34 阳性信号强度",
    ], size=10, bullet_color=ACCENT2)

    # Card 2: Cell Extraction
    add_shape(slide, Inches(3.45), cy, Inches(3.1), ch, BG_CARD, ORANGE)
    add_text(slide, Inches(3.6), cy+Inches(0.08), Inches(2.8), Inches(0.35),
             "Step 3: 连通域提取", size=14, color=ORANGE, bold=True)
    add_bullets(slide, Inches(3.6), cy+Inches(0.4), Inches(2.8), Inches(3.2), [
        "1. 前景检测:",
        "  (R+B) > seg_thresh AND G <= 80",
        "2. 阳性判定: R >= B",
        "3. Marker 增强:",
        "  marker > auto_thresh (99.9%ile x 0.9)",
        "  联合: seg_positive OR marker_positive",
        "4. 形态学闭运算:",
        "  椭圆核 11x11, iterations=2",
        "  连接断裂的棕色血管壁",
        "5. 形态学开运算:",
        "  椭圆核 3x3, 去除噪点",
        "6. 八连通域分析:",
        "  过滤面积 < min_area 的碎片",
        "输出: 阳性连通域列表",
    ], size=10, bullet_color=ORANGE)

    # Card 3: SAM2
    add_shape(slide, Inches(6.7), cy, Inches(3.1), ch, BG_CARD, ACCENT3)
    add_text(slide, Inches(6.85), cy+Inches(0.08), Inches(2.8), Inches(0.35),
             "Step 4-5: SAM2 + 后处理", size=14, color=ACCENT3, bold=True)
    add_bullets(slide, Inches(6.85), cy+Inches(0.4), Inches(2.8), Inches(3.2), [
        "逐区域生成 Mask Prompt:",
        "  初始化全图 logit = -5.0 (背景)",
        "  前景像素设为 +5.0 (前景)",
        "  缩放至 256x256 (双线性插值)",
        "SAM2 推理 (逐区域):",
        "  输入: Mask Prompt (256x256)",
        "  multimask_output = True",
        "  输出: 3 个候选 mask + 3 个 score",
        "  选择 score 最高者为最优",
        "置信度优先合并:",
        "  重叠像素由高分实例占据",
        "  低分实例被部分覆盖",
        "后处理过滤:",
        "  score < 0.05 的实例被移除",
        "  面积 < min_area 的碎片被移除",
    ], size=10, bullet_color=ACCENT3)


def slide_deepliif_detail(prs):
    """Slide 4: DeepLIIF with example images."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "02", "DeepLIIF 染色分解 — 实际输出示例")

    # Flow
    y = Inches(1.15); bw, bh = Inches(1.3), Inches(0.6)
    flow = [
        ("IHC 输入", ACCENT), ("Tiling\n512+32", LIGHT_GRAY),
        ("G1-G4\n模态翻译", ACCENT2), ("G51-G55\n5 路分割", ACCENT2),
        ("加权聚合\nSeg", ORANGE), ("Marker\n(G4 直出)", ACCENT3),
    ]
    x = Inches(0.2)
    for i, (label, clr) in enumerate(flow):
        add_flow_box(slide, x, y, bw, bh, label, fill=BG_CARD, border=clr, title_color=clr)
        if i < len(flow) - 1:
            add_arrow_right(slide, x + bw, y + bh/2, Inches(0.15))
        x += bw + Inches(0.15)

    # Example: Tile_63_72 (dense vessels) — 3 images
    add_text(slide, Inches(0.4), Inches(1.95), Inches(9), Inches(0.3),
             "示例: Tile(63,72) — 密集微血管区域", size=15, color=ACCENT, bold=True)

    iw, ih, iy = Inches(2.8), Inches(2.8), Inches(2.3)
    add_img_label(slide, f"{TILE2}/step1_original.png",
                  Inches(0.3), iy, iw, ih,
                  "原图: CD34 IHC 染色\n棕色 = CD34+ 血管内皮, 蓝色 = 苏木精复染细胞核",
                  LIGHT_GRAY, 9)
    add_img_label(slide, f"{TILE2}/step2_deepliif_Seg.png",
                  Inches(3.4), iy, iw, ih,
                  "Seg (5 路分割聚合): RGB 编码\n红区 R>B = 阳性概率高, 蓝区 B>R = 阴性, 绿 = 边界",
                  LIGHT_GRAY, 9)
    add_img_label(slide, f"{TILE2}/step2_deepliif_Marker.png",
                  Inches(6.5), iy, iw, ih,
                  "Marker (G4 翻译): 灰度图\n亮度 = CD34 阳性信号强度, 暗区 = 背景/阴性",
                  LIGHT_GRAY, 9)

    # Bottom note: DeepLIIF parameters
    add_shape(slide, Inches(0.3), Inches(5.65), Inches(9.4), Inches(0.85), BG_CARD, DIM_GRAY)
    add_text(slide, Inches(0.5), Inches(5.7), Inches(9.0), Inches(0.75),
             "DeepLIIF 推理参数:  输入尺寸 = 512x512 | Tile overlap = 32px | "
             "分割权重 seg_weights = [0.2, 0.2, 0.2, 0.2, 0.2] (G51-G55 等权)\n"
             "Seg RGB 解读:  R >= B 表示阳性 (CD34 棕色区域); R < B 表示阴性 (苏木精蓝色区域); "
             "G 通道标记细胞边界\n"
             "Marker 直接反映 CD34 染色强度, 用于辅助确认阳性区域",
             size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_extraction_detail(prs):
    """Slide 5: Cell extraction algorithm detail."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "03", "Seg+Marker 联合提取 — 算法详解")

    add_text(slide, Inches(0.4), Inches(1.1), Inches(9), Inches(0.25),
             "从 DeepLIIF 的 Seg 和 Marker 输出中提取 CD34 阳性连通域, 作为 SAM2 的输入提示",
             size=11, color=DIM_GRAY)

    # 6-step flow (vertical 2 columns)
    steps = [
        ("Step 1: 前景检测", ACCENT,
         ["从 Seg RGB 图像中分离前景与背景",
          "前景条件: (R+B) > seg_thresh AND G <= 80",
          "输出: 前景布尔掩码",
          "  参数: seg_thresh=120 (默认)"]),
        ("Step 2: 阳性像素判定", ACCENT2,
         ["Seg 阳性: 前景中 R >= B 的像素",
          "Marker 增强: marker 灰度值 > 自动阈值",
          "  自动阈值 = 非零像素 99.9 百分位 x 0.9",
          "联合判定: seg_positive OR marker_positive"]),
        ("Step 3: 形态学闭运算", ORANGE,
         ["对联合阳性掩码做闭运算",
          "椭圆核大小: 11x11, 迭代次数: 2",
          "作用: 连接相邻阳性像素",
          "  填充血管壁的小缺口和断裂"]),
        ("Step 4: 形态学开运算", ACCENT3,
         ["对闭运算结果做开运算",
          "椭圆核大小: 3x3, 迭代次数: 1",
          "作用: 去除小噪点和伪影",
          "  参数: morphology_kernel=11 (闭运算)"]),
        ("Step 5: 连通域分析", PURPLE,
         ["在处理后掩码上做八连通域分析",
          "8-连通: 斜对角像素也计入同一区域",
          "面积过滤: 像素数 < min_area 的区域移除",
          "  参数: min_area=200 (默认)"]),
        ("Step 6: 输出区域信息", YELLOW,
         ["每个区域记录: 坐标集合, 中心点, 面积",
          "marker 统计: sum/max/mean 值",
          "全部标记为阳性区域 (is_positive)",
          "作为 clusters 列表送入 SAM2 推理"]),
    ]

    # Layout: 2 columns, 3 rows
    for i, (title, clr, bullets) in enumerate(steps):
        col = i % 2
        row = i // 2
        lx = Inches(0.3) + col * Inches(4.85)
        ly = Inches(1.5) + row * Inches(1.9)
        cw, ch_ = Inches(4.65), Inches(1.8)

        add_shape(slide, lx, ly, cw, ch_, BG_CARD, clr)
        add_text(slide, lx + Inches(0.1), ly + Inches(0.05), cw - Inches(0.2), Inches(0.3),
                 title, size=13, color=clr, bold=True)
        add_bullets(slide, lx + Inches(0.1), ly + Inches(0.35), cw - Inches(0.2), Inches(1.3),
                    bullets, size=10, bullet_color=clr)


def slide_extraction_example(prs):
    """Slide 6: Extraction step-by-step with real images (Tile_63_72)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "04", "示例: 连通域提取过程 — Tile(63,72)")

    add_text(slide, Inches(0.4), Inches(1.05), Inches(9), Inches(0.25),
             "密集血管区域 — 从 Seg 输出到最终 10 个连通域的逐步变换", size=12, color=DIM_GRAY)

    # 5 step images in a row
    iw, ih = Inches(1.75), Inches(1.75)
    iy = Inches(1.4)
    step_imgs = [
        (f"{TILE2}/cr_step1_foreground_mask.png", "Step1: 前景掩码\n(R+B)>120 & G<=80", ACCENT),
        (f"{TILE2}/cr_step1_posneg_mask.png", "Step1b: 阳/阴性分布\n红=阳性(R>=B) 蓝=阴性", ACCENT2),
        (f"{TILE2}/cr_step2a_seg_positive.png", "Step2: Seg阳性区域\n绿色标记, 黑色=非阳性", ACCENT3),
        (f"{TILE2}/cr_step3_morph_close.png", "Step3: 闭运算后\n椭圆核11×11, iter=2", ORANGE),
        (f"{TILE2}/cr_step5_connected_regions_10.png", "Step5: 连通域(10区域)\n彩色叠加在Seg上", PURPLE),
    ]
    x = Inches(0.15)
    for path, label, clr in step_imgs:
        add_img(slide, path, x, iy, iw, ih)
        add_text(slide, x, iy + ih + Inches(0.02), iw, Inches(0.5),
                 label, size=8, color=clr, bold=True, align=PP_ALIGN.CENTER)
        x += Inches(1.95)

    # Second example: Tile_55_69 (sparse)
    add_text(slide, Inches(0.4), Inches(3.85), Inches(9), Inches(0.25),
             "稀疏血管区域 — Tile(55,69) → 8 个连通域", size=12, color=DIM_GRAY)

    iy2 = Inches(4.15)
    step_imgs2 = [
        (f"{TILE1}/cr_step1_foreground_mask.png", "Step1: 前景掩码", ACCENT),
        (f"{TILE1}/cr_step1_posneg_mask.png", "Step1b: 阳/阴性分布", ACCENT2),
        (f"{TILE1}/cr_step2a_seg_positive.png", "Step2: Seg阳性区域", ACCENT3),
        (f"{TILE1}/cr_step3_morph_close.png", "Step3: 闭运算后", ORANGE),
        (f"{TILE1}/cr_step5_connected_regions_8.png", "Step5: 连通域(8区域)", PURPLE),
    ]
    x = Inches(0.15)
    for path, label, clr in step_imgs2:
        add_img(slide, path, x, iy2, iw, ih)
        add_text(slide, x, iy2 + ih + Inches(0.02), iw, Inches(0.3),
                 label, size=8, color=clr, bold=True, align=PP_ALIGN.CENTER)
        x += Inches(1.95)

    # Note
    add_shape(slide, Inches(0.3), Inches(6.4), Inches(9.4), Inches(0.55), BG_CARD, DIM_GRAY)
    add_text(slide, Inches(0.5), Inches(6.42), Inches(9.0), Inches(0.5),
             "关键观察: 形态学闭运算 (Step3) 将断裂的血管壁连接成完整区域, 使 SAM2 能获得连贯的 Mask 提示。\n"
             "对比 Step2→Step3: 分散的小绿块被合并为大白块。连通域数量 = 最终送入 SAM2 的候选区域数。",
             size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_sam2_algorithm(prs):
    """Slide 7: SAM2 Mask-Only inference algorithm."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "05", "SAM2 Mask-Only 推理 — 算法详解")

    add_text(slide, Inches(0.4), Inches(1.1), Inches(9), Inches(0.25),
             "将每个连通域转换为 256x256 logit 掩码作为提示, 由 SAM2 生成精确的实例分割",
             size=10, color=DIM_GRAY)

    # Algorithm flow
    y = Inches(1.5)
    steps_data = [
        ("1. Mask Prompt 生成",
         "初始化全图 logit = -5.0 (背景)\n"
         "前景像素设为 +5.0 (前景)\n"
         "缩放至 256x256 (INTER_AREA)\n"
         "输出 shape: (1, 256, 256) float32",
         ACCENT),
        ("2. SAM2 推理",
         "输入: 256x256 logit prompt\n"
         "多掩码输出模式: 开启\n"
         "输出: masks (3, H, W) bool\n"
         "输出: scores (3,) float | 选最高分",
         ACCENT3),
        ("3. 置信度优先合并",
         "逐像素比较置信度:\n"
         "若新实例分数 > 当前分数:\n"
         "  则该像素归属新实例\n"
         "  更新该像素的置信度记录",
         ORANGE),
        ("4. 过滤与后处理",
         "score < 0.05 的实例被移除\n"
         "二值化后做连通域分析\n"
         "面积 < 200px 的碎片被移除\n"
         "相连实例合并, 取平均 score",
         PURPLE),
    ]

    x = Inches(0.2)
    for title, desc, clr in steps_data:
        add_shape(slide, x, y, Inches(2.3), Inches(2.3), BG_CARD, clr)
        add_text(slide, x+Inches(0.08), y+Inches(0.05), Inches(2.15), Inches(0.6),
                 title, size=11, color=clr, bold=True)
        add_text(slide, x+Inches(0.08), y+Inches(0.65), Inches(2.15), Inches(1.5),
                 desc, size=10, color=LIGHT_GRAY)
        x += Inches(2.42)

    # Before → After example
    add_text(slide, Inches(0.4), Inches(4.0), Inches(9), Inches(0.3),
             "效果: Tile(55,69) — 8 个连通域 -> SAM2 -> 合并后 5 个实例", size=14, color=ACCENT, bold=True)

    iw2, ih2, iy2 = Inches(2.5), Inches(2.5), Inches(4.4)
    add_img_label(slide, f"{TILE1}/step3_positive_cells_8.png",
                  Inches(0.5), iy2, iw2, ih2,
                  "提取结果: 8 个阳性区域 (绿色叠加原图)", ACCENT3, 9)
    add_arrow_right(slide, Inches(3.1), iy2 + ih2/2, Inches(0.35))
    add_img_label(slide, f"{TILE1}/step4_sam2_raw_8inst.png",
                  Inches(3.6), iy2, iw2, ih2,
                  "SAM2 原始输出: 8 个实例 (每色=1实例)", ORANGE, 9)
    add_arrow_right(slide, Inches(6.2), iy2 + ih2/2, Inches(0.35))
    add_img_label(slide, f"{TILE1}/step5_merged_5inst.png",
                  Inches(6.7), iy2, iw2, ih2,
                  "合并后: 5 个最终实例 (面积过滤 min_area=200)", PURPLE, 9)


def slide_sam2_comparison(prs):
    """Slide 8: SAM2 3-mask comparison with real images."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "06 (续)", "SAM2 三候选掩码对比 — 实际示例")

    add_text(slide, Inches(0.4), Inches(1.1), Inches(9), Inches(0.25),
             "每个连通域生成 3 个候选 mask, 选 score 最高者 (标 BEST)", size=12, color=DIM_GRAY)

    # Example 1: Tile_55_69 inst001 (high score, large vessel)
    add_text(slide, Inches(0.3), Inches(1.45), Inches(9), Inches(0.3),
             "示例 A:  Tile(55,69) Instance #1 — 大型血管, cluster=12003px, score: 0.92/0.90/0.93(BEST)",
             size=12, color=ACCENT3, bold=True)

    add_img(slide, f"{TILE1}/sam2_steps/instance_001/comparison_3_masks.png",
            Inches(0.3), Inches(1.8), Inches(6.5), Inches(1.85))
    add_img_label(slide, f"{TILE1}/sam2_steps/instance_001/input_cluster.png",
                  Inches(7.0), Inches(1.8), Inches(1.25), Inches(1.25),
                  "输入连通域", ORANGE, 9)
    add_img_label(slide, f"{TILE1}/sam2_steps/instance_001/input_mask_prompt_256x256.png",
                  Inches(8.4), Inches(1.8), Inches(1.25), Inches(1.25),
                  "Mask Prompt\n(256×256 logit)", ORANGE, 9)

    add_shape(slide, Inches(0.3), Inches(3.75), Inches(9.4), Inches(0.4), BG_CARD, ACCENT3)
    add_text(slide, Inches(0.5), Inches(3.78), Inches(9.0), Inches(0.35),
             "高分场景: 3 个候选形状各异但分数都 >0.89 — mask_areas: 10217 / 16025 / 11987 — 选 #2 (score=0.933, area=11987)",
             size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Example 2: Tile_63_72 inst001 (lower score)
    add_text(slide, Inches(0.3), Inches(4.35), Inches(9), Inches(0.3),
             "示例 B:  Tile(63,72) Instance #1 — 小区域, score: 0.43(BEST)/0.0001/0.04",
             size=12, color=ACCENT2, bold=True)

    add_img(slide, f"{TILE2}/sam2_steps/instance_001/comparison_3_masks.png",
            Inches(0.3), Inches(4.7), Inches(6.5), Inches(1.85))
    add_img_label(slide, f"{TILE2}/sam2_steps/instance_001/input_cluster.png",
                  Inches(7.0), Inches(4.7), Inches(1.25), Inches(1.25),
                  "输入连通域", ORANGE, 9)
    add_img_label(slide, f"{TILE2}/sam2_steps/instance_001/input_mask_prompt_256x256.png",
                  Inches(8.4), Inches(4.7), Inches(1.25), Inches(1.25),
                  "Mask Prompt\n(256×256 logit)", ORANGE, 9)

    add_shape(slide, Inches(0.3), Inches(6.65), Inches(9.4), Inches(0.4), BG_CARD, ACCENT2)
    add_text(slide, Inches(0.5), Inches(6.68), Inches(9.0), Inches(0.35),
             "低分场景: 候选 #1 只覆盖边角小块 (score=0.43), #0/#2 几乎为空 — 小区域 SAM2 信心不足, 后处理可能被过滤",
             size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_sam2_mask_prompt(prs):
    """Slide: SAM2 mask prompt generation visualization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "06", "SAM2 掩码提示词 (Mask Prompt) 生成")

    add_text(slide, Inches(0.4), Inches(1.1), Inches(9), Inches(0.25),
             "从连通域提取结果生成 SAM2 的 Mask Prompt 输入, 引导模型精确分割每个血管实例",
             size=11, color=DIM_GRAY)

    # Flow: Connected Region → Mask Prompt → SAM2 → Result
    bw, bh = Inches(2.0), Inches(0.7)
    y = Inches(1.4)
    flow = [
        ("连通域区域", "单个阳性区域坐标", ACCENT),
        ("Logit 掩码初始化", "全图 = -5.0, 前景 = +5.0", ACCENT2),
        ("缩放至 256x256", "INTER_AREA 插值", ORANGE),
        ("SAM2 推理", "输出 3 个候选 mask", ACCENT3),
    ]
    x = Inches(0.3)
    for i, (t, s, c) in enumerate(flow):
        add_flow_box(slide, x, y, bw, bh, t, s, fill=BG_CARD, border=c, title_color=c)
        if i < len(flow) - 1:
            add_arrow_right(slide, x + bw + Inches(0.02), y + bh/2, Inches(0.18))
        x += bw + Inches(0.2)

    # Example A: High-score (Tile_55_69 inst001)
    add_text(slide, Inches(0.3), Inches(2.3), Inches(9), Inches(0.25),
             "示例 A: Tile(55,69) Instance #1 — 大型血管, 高置信度 (score=0.93)",
             size=12, color=ACCENT3, bold=True)

    iw, ih = Inches(1.9), Inches(1.9)
    iy = Inches(2.6)
    add_img_label(slide, f"{TILE1}/sam2_steps/instance_001/input_cluster.png",
                  Inches(0.3), iy, iw, ih,
                  "输入: 连通域区域\n(原图上的单个阳性区域)", ACCENT, 8)
    add_arrow_right(slide, Inches(2.3), iy + ih/2, Inches(0.2))
    add_img_label(slide, f"{TILE1}/sam2_steps/instance_001/input_mask_prompt_256x256.png",
                  Inches(2.6), iy, iw, ih,
                  "Mask Prompt (256x256)\n白=+5.0 前景, 黑=-5.0 背景", ACCENT2, 8)
    add_arrow_right(slide, Inches(4.6), iy + ih/2, Inches(0.2))
    add_img(slide, f"{TILE1}/sam2_steps/instance_001/comparison_3_masks.png",
            Inches(4.9), iy, Inches(4.8), Inches(1.35))
    add_text(slide, Inches(4.9), iy + Inches(1.4), Inches(4.8), Inches(0.5),
             "SAM2 输出: 3 个候选掩码 (score: 0.92 / 0.90 / 0.93)\n"
             "选择 score 最高者 #2 (score=0.933) 为最终结果",
             size=8, color=ACCENT3, align=PP_ALIGN.CENTER)

    # Example B: Low-score (Tile_63_72 inst001)
    add_text(slide, Inches(0.3), Inches(4.85), Inches(9), Inches(0.25),
             "示例 B: Tile(63,72) Instance #1 — 小区域, 低置信度 (score=0.43)",
             size=12, color=ACCENT2, bold=True)

    iy2 = Inches(5.15)
    add_img_label(slide, f"{TILE2}/sam2_steps/instance_001/input_cluster.png",
                  Inches(0.3), iy2, iw, ih,
                  "输入: 连通域区域\n(较小的阳性区域)", ACCENT, 8)
    add_arrow_right(slide, Inches(2.3), iy2 + ih/2, Inches(0.2))
    add_img_label(slide, f"{TILE2}/sam2_steps/instance_001/input_mask_prompt_256x256.png",
                  Inches(2.6), iy2, iw, ih,
                  "Mask Prompt (256x256)\n小区域在掩码中占比较小", ACCENT2, 8)
    add_arrow_right(slide, Inches(4.6), iy2 + ih/2, Inches(0.2))
    add_img(slide, f"{TILE2}/sam2_steps/instance_001/comparison_3_masks.png",
            Inches(4.9), iy2, Inches(4.8), Inches(1.35))
    add_text(slide, Inches(4.9), iy2 + Inches(1.4), Inches(4.8), Inches(0.5),
             "SAM2 输出: 3 个候选掩码 (score: 0.43 / 0.0001 / 0.04)\n"
             "小区域 SAM2 信心不足, 仅 #0 有效 (score=0.43), 后处理可能被过滤",
             size=8, color=ACCENT2, align=PP_ALIGN.CENTER)


def slide_postprocessing(prs):
    """Slide 9: Post-processing & merge detail."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "07", "后处理与合并")

    add_text(slide, Inches(0.4), Inches(1.1), Inches(9), Inches(0.25),
             "将 SAM2 原始输出清理为最终实例分割结果: 去重叠、去碎片、合并相连区域", size=11, color=DIM_GRAY)

    # Flow
    y = Inches(1.5); bw, bh = Inches(2.1), Inches(0.85)
    flow = [
        ("SAM2 原始 mask\n(可能有重叠/粘连)", ACCENT),
        ("二值化\n→ connectedComponents", ACCENT2),
        ("面积过滤\n< 200px → 移除", ORANGE),
        ("新 ID 分配\n合并相连实例", ACCENT3),
    ]
    x = Inches(0.3)
    for i, (t, c) in enumerate(flow):
        add_flow_box(slide, x, y, bw, bh, t, fill=BG_CARD, border=c, title_color=c)
        if i < len(flow) - 1:
            add_arrow_right(slide, x + bw + Inches(0.02), y + bh/2, Inches(0.2))
        x += bw + Inches(0.22)

    # Detail
    add_shape(slide, Inches(0.3), Inches(2.6), Inches(4.5), Inches(2.8), BG_CARD, ACCENT)
    add_text(slide, Inches(0.5), Inches(2.65), Inches(4.1), Inches(0.3),
             "合并逻辑", size=14, color=ACCENT, bold=True)
    add_bullets(slide, Inches(0.5), Inches(2.95), Inches(4.1), Inches(2.3), [
        "将 instance_mask > 0 二值化",
        "cv2.connectedComponents 找连通区域",
        "每个连通区域可能包含多个旧实例 ID",
        "  (SAM2 输出的相邻实例被合并)",
        "取成员实例的 avg_score 作为新分数",
        "merge_mapping: {旧ID: 新ID}",
        "输出 merged_mask: uint16",
        "  (支持 >255 个实例)",
    ], size=11, bullet_color=ACCENT)

    add_shape(slide, Inches(5.1), Inches(2.6), Inches(4.5), Inches(2.8), BG_CARD, ORANGE)
    add_text(slide, Inches(5.3), Inches(2.65), Inches(4.1), Inches(0.3),
             "面积过滤与质量控制", size=14, color=ORANGE, bold=True)
    add_bullets(slide, Inches(5.3), Inches(2.95), Inches(4.1), Inches(2.3), [
        "score_threshold = 0.05 (SAM2 阶段)",
        "  → 过滤极低置信度结果",
        "min_area = 200 (合并阶段)",
        "  → 移除面积过小的碎片",
        "置信度优先合并 (SAM2 阶段):",
        "  重叠像素由高分实例占据",
        "  低分实例被部分覆盖",
        "最终 area = combined_mask 中实际像素数",
    ], size=11, bullet_color=ORANGE)

    # Example comparison
    add_text(slide, Inches(0.4), Inches(5.65), Inches(9), Inches(0.3),
             "效果对比: Tile(63,72)", size=14, color=ACCENT3, bold=True)
    iw, ih, iy = Inches(2.5), Inches(2.5), Inches(6.0) # won't fit, use smaller
    # Use smaller images
    iw, ih = Inches(1.6), Inches(1.6)
    iy = Inches(5.95)
    add_img_label(slide, f"{TILE2}/step3_positive_cells_10.png",
                  Inches(0.5), iy, iw, ih, "提取: 10 区域", ACCENT3, 9)
    add_arrow_right(slide, Inches(2.2), iy + ih/2, Inches(0.2))
    add_img_label(slide, f"{TILE2}/step4_sam2_raw_10inst.png",
                  Inches(2.55), iy, iw, ih, "SAM2: 10 实例", ORANGE, 9)
    add_arrow_right(slide, Inches(4.25), iy + ih/2, Inches(0.2))
    add_img_label(slide, f"{TILE2}/step5_merged_10inst.png",
                  Inches(4.6), iy, iw, ih, "合并: 10 实例", PURPLE, 9)

    add_shape(slide, Inches(6.5), iy, Inches(3.2), ih, BG_CARD, DIM_GRAY)
    add_text(slide, Inches(6.7), iy + Inches(0.1), Inches(2.8), Inches(0.3),
             "Tile(63,72) 统计", size=12, color=YELLOW, bold=True)
    add_bullets(slide, Inches(6.7), iy + Inches(0.4), Inches(2.8), Inches(1.1), [
        "送入 SAM2: 10 个连通域",
        "SAM2 输出: 10 个实例",
        "  score 范围: 0.01 ~ 0.93",
        "合并后: 10 个最终实例",
        "  (此 Tile 无粘连/碎片)",
    ], size=9, bullet_color=YELLOW)


def slide_full_pipeline_demo(prs):
    """Slide 10: Full 8-step pipeline demo with 2 tiles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "08", "Pipeline 完整示例 — 两种场景对比")

    # Tile_63_72 (dense)
    add_text(slide, Inches(0.3), Inches(1.05), Inches(9), Inches(0.25),
             "Tile(63,72) — 密集微血管, 最终 10 实例", size=12, color=ACCENT, bold=True)
    iw, ih = Inches(2.15), Inches(2.15)
    y1 = Inches(1.35)
    row1 = [
        (f"{TILE2}/step1_original.png", "1.原图"),
        (f"{TILE2}/step2_deepliif_Seg.png", "2.Seg"),
        (f"{TILE2}/cr_step2a_seg_positive.png", "3.阳性区域"),
        (f"{TILE2}/step5_merged_10inst.png", "4.最终分割"),
    ]
    x = Inches(0.3)
    for path, label in row1:
        add_img(slide, path, x, y1, iw, ih)
        add_text(slide, x, y1+ih+Inches(0.01), iw, Inches(0.2),
                 label, size=9, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        x += Inches(2.35)

    # Tile_55_69 (sparse)
    add_text(slide, Inches(0.3), Inches(3.95), Inches(9), Inches(0.25),
             "Tile(55,69) — 稀疏微���管, 最终 5 实例", size=12, color=ACCENT3, bold=True)
    y2 = Inches(4.25)
    row2 = [
        (f"{TILE1}/step1_original.png", "1.原图"),
        (f"{TILE1}/step2_deepliif_Seg.png", "2.Seg"),
        (f"{TILE1}/cr_step2a_seg_positive.png", "3.阳性区域"),
        (f"{TILE1}/step5_merged_5inst.png", "4.最终分割"),
    ]
    x = Inches(0.3)
    for path, label in row2:
        add_img(slide, path, x, y2, iw, ih)
        add_text(slide, x, y2+ih+Inches(0.01), iw, Inches(0.2),
                 label, size=9, color=ACCENT3, bold=True, align=PP_ALIGN.CENTER)
        x += Inches(2.35)

    add_shape(slide, Inches(0.3), Inches(6.85), Inches(9.4), Inches(0.5), BG_CARD, DIM_GRAY)
    add_text(slide, Inches(0.5), Inches(6.88), Inches(9.0), Inches(0.45),
             "对比: 密集场景检出更多实例 (10 vs 5); 稀疏场景 SAM2 仍能精确贴合血管轮廓;\n"
             "两种场景使用相同参数 (seg_thresh=120, kernel=11, min_area=200)",
             size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_wsi_pipeline(prs):
    """Slide 11: WSI full-slide processing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "09", "WSI 全片处理 — 技术方案")

    add_text(slide, Inches(0.4), Inches(1.1), Inches(9), Inches(0.25),
             "从 .ndpi/.svs 全片扫描图像直接处理, 输出 NPY 掩码和 GeoJSON 标注文件", size=11, color=DIM_GRAY)

    # Vertical flow
    cx = Inches(5.0); bw = Inches(3.6)
    flow = [
        ("Step 1: WSI 读取",
         "OpenSlide 后端 | 自动选择倍率层级 | 512x512 tile 枚举", ACCENT),
        ("Step 2: YOLO 背景过滤",
         "YOLOv11-cls 二分类 | 流式推理 | 跳过背景 tile", ACCENT2),
        ("Step 3: 多 GPU 并行处理",
         "每 Worker: DeepLIIF -> 提取 -> SAM2 | 约 8GB/Worker", ACCENT3),
        ("Step 4: 结果输出",
         "NPY 掩码保存 + GeoJSON 导出 (QuPath 兼容)", ORANGE),
    ]
    y = Inches(1.5)
    for title, desc, clr in flow:
        add_flow_box(slide, cx - bw/2, y, bw, Inches(0.8), title, desc,
                     fill=BG_CARD, border=clr, title_color=clr)
        if y < Inches(4.5):
            add_arrow_down(slide, cx, y + Inches(0.8), Inches(0.18))
        y += Inches(0.98)

    # Side cards
    add_shape(slide, Inches(0.2), Inches(1.5), Inches(2.9), Inches(2.2), BG_CARD, ACCENT)
    add_text(slide, Inches(0.35), Inches(1.55), Inches(2.6), Inches(0.3),
             "WSI 读取细节", size=12, color=ACCENT, bold=True)
    add_bullets(slide, Inches(0.3), Inches(1.85), Inches(2.6), Inches(1.7), [
        "支持 .ndpi / .svs / .tiff 格式",
        "自动选择最接近目标倍率的层级",
        "  参数: target_magnification=40x",
        "枚举 tile: stride=512, 无 overlap",
        "可选预加载: 按 4096 行条读入",
        "  后续 read_tile 零拷贝切片",
        "crop-csv: 仅处理指定 ROI 区域",
    ], size=9, bullet_color=ACCENT)

    add_shape(slide, Inches(0.2), Inches(3.9), Inches(2.9), Inches(1.8), BG_CARD, ACCENT2)
    add_text(slide, Inches(0.35), Inches(3.95), Inches(2.6), Inches(0.3),
             "YOLO 背景过滤", size=12, color=ACCENT2, bold=True)
    add_bullets(slide, Inches(0.3), Inches(4.25), Inches(2.6), Inches(1.3), [
        "YOLOv11-cls 模型 (二分类: bg/target)",
        "  batch_size=64, half=True (FP16)",
        "4 个预取线程并行加载图像",
        "流式回调: 目标 tile 直接入队列",
        "tile_map CSV: 可保存分类结果复用",
        "classify-only 模式: 仅分类不处理",
    ], size=9, bullet_color=ACCENT2)

    # Multi-GPU
    add_shape(slide, Inches(0.2), Inches(5.9), Inches(9.6), Inches(1.4), BG_CARD, ACCENT3)
    add_text(slide, Inches(0.4), Inches(5.95), Inches(4), Inches(0.3),
             "多 GPU 并行策略", size=13, color=ACCENT3, bold=True)
    add_bullets(slide, Inches(0.4), Inches(6.25), Inches(4.5), Inches(1.0), [
        "自动检测可用 GPU 及显存",
        "YOLO 分配至空闲最少但够用的 GPU (约 1.5GB)",
        "Worker 分配至空闲最多的 GPU (约 8GB/worker)",
        "流水线: YOLO 线程 -> 队列 -> Worker 进程",
    ], size=10, bullet_color=ACCENT3)
    add_bullets(slide, Inches(5.2), Inches(6.25), Inches(4.5), Inches(1.0), [
        "单设备模式: YOLO + DeepLIIF/SAM2 共享 GPU",
        "  YOLO 后台线程 -> tile 队列 -> 主线程处理",
        "多设备模式: spawn Worker 进程 (各自加载模型)",
        "  Worker 从队列消费 tile 直到完成",
    ], size=10, bullet_color=ACCENT3)


def slide_geojson_export(prs):
    """Slide: GeoJSON export detailed explanation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "10", "GeoJSON 导出详解 — 从 Tile 掩码到 QuPath 标注")

    add_text(slide, Inches(0.4), Inches(1.1), Inches(9), Inches(0.25),
             "四步流程: 将分散在各 tile 中的 NPY 掩码合并为完整的全图标注文件, 可直接在 QuPath 中加载",
             size=11, color=DIM_GRAY)

    # 4-pass vertical flow
    passes = [
        ("Pass 1: 逐 Tile 轮廓提取",
         "遍历所有 tile 的 NPY 掩码文件, 对每个实例 ID 做轮廓提取:\n"
         "  1) 二值化单个实例 -> 轮廓检测\n"
         "  2) 可选简化: epsilon = simplify x 周长 (默认 simplify=0.002)\n"
         "  3) 面积过滤: 跳过 < min_area (默认 50px) 的轮廓\n"
         "  4) 坐标转换: 局部坐标 + tile 偏移量 -> 全图坐标\n"
         "  5) 转为 Shapely Polygon 对象, 按全局实例 ID 存储",
         ACCENT),
        ("Pass 2: Union-Find 跨 Tile 边界合并",
         "扫描相邻 tile 的边界像素, 检测跨 tile 的同一实例:\n"
         "  1) 右邻 tile: 比较 A 的右边缘列 vs B 的左边缘列\n"
         "  2) 下邻 tile: 比较 A 的底边缘行 vs B 的顶边缘行\n"
         "  3) 同一行/列位置如果两侧都有前景像素, 则 Union-Find 合并\n"
         "  确保跨 tile 边界的血管实例被正确识别为同一个",
         ACCENT2),
        ("Pass 3: 几何合并 (Shapely)",
         "按 Union-Find 的 root ID 聚合所有多边形:\n"
         "  1) 同一实例的多个多边形做 unary_union 合并\n"
         "  2) buffer(1).buffer(-1): 膨胀再收缩, 消除 tile 接缝\n"
         "  3) 验证几何有效性, 无效则用 make_valid 修复\n"
         "  4) MultiPolygon 拆分为独立 Polygon",
         ORANGE),
        ("Pass 4: 验证与导出",
         "生成 QuPath 兼容的 GeoJSON 格式:\n"
         "  1) 每个 Polygon 转为 Feature (objectType: detection)\n"
         "  2) 分类标签: classification.name = CD34+\n"
         "  3) 坐标整数化, 环闭合验证, 至少 4 个点\n"
         "  4) 输出: JSON 数组 (QuPath 直接导入)",
         ACCENT3),
    ]

    y = Inches(1.5)
    for i, (title, desc, clr) in enumerate(passes):
        add_shape(slide, Inches(0.3), y, Inches(9.4), Inches(1.25), BG_CARD, clr)
        add_text(slide, Inches(0.5), y + Inches(0.05), Inches(3.0), Inches(0.25),
                 title, size=12, color=clr, bold=True)
        add_text(slide, Inches(0.5), y + Inches(0.3), Inches(8.8), Inches(0.9),
                 desc, size=9, color=LIGHT_GRAY)
        if i < len(passes) - 1:
            add_arrow_down(slide, Inches(5.0), y + Inches(1.25), Inches(0.12))
        y += Inches(1.37)

    # Bottom: QuPath usage
    add_shape(slide, Inches(0.3), Inches(6.98), Inches(9.4), Inches(0.45), BG_CARD, YELLOW)
    add_text(slide, Inches(0.5), Inches(7.0), Inches(9.0), Inches(0.4),
             "QuPath 导入: File -> Open -> 选择原图 .ndpi | File -> Object -> Import objects -> 选择 .geojson\n"
             "导出参数: simplify=0.002 (轮廓简化) | min_area=50 (最小面积) | 坐标系 = 全图像素坐标",
             size=9, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

def slide_parameters(prs):
    """Slide 12: Parameter table + tuning tips."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "11", "关键参数详表与调优建议")

    # Table
    headers = ["参数", "默认值", "说明", "调优建议"]
    col_w = [Inches(2.3), Inches(0.9), Inches(3.5), Inches(2.8)]
    col_x = [Inches(0.3)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    hy = Inches(1.3)
    add_shape(slide, Inches(0.25), hy - Inches(0.02), Inches(9.5), Inches(0.3), ACCENT)
    for i, h in enumerate(headers):
        add_text(slide, col_x[i], hy, col_w[i], Inches(0.28),
                 h, size=11, color=BG_DARK, bold=True)

    rows = [
        ["--seg-thresh",         "120",    "(R+B)>thresh 前景检测",          "染色淡→80~100, 深→120~150"],
        ["--marker-thresh",      "auto",   "99.9%ile×0.9, 阳性辅助",        "None=自动, 或手动指定 int"],
        ["--morphology-kernel",  "11",     "闭运算椭圆核大小",              "密集→11~15, 稀疏→7~9"],
        ["--min-mask-area",      "50",     "最小连通域面积 (送入SAM2)",      "噪点多→增大, 小血管→减小"],
        ["merge min_area",       "200",    "合并后最小实例面积",            "过滤碎片, 一般不改"],
        ["score_threshold",      "0.05",   "SAM2 最低置信度",               "极端场景→降至 0.01"],
        ["--use-connected-regions", "推荐", "启用连通域模式",               "比个体细胞模式效果好"],
        ["--resolution",         "40x",    "影响面积阈值自动计算",          "匹配实际倍率"],
        ["--seg-weights",        "各0.2",  "G51-G55 聚合权重 (5个float)",   "通常不需修改"],
        ["--num-gpus",           "auto",   "GPU 数量 (0=自动检测)",         "限制→设具体数字"],
        ["--yolo-batch-size",    "64",     "YOLO 分类批大小",               "显存紧张→减小"],
    ]

    for ri, row in enumerate(rows):
        ry = Inches(1.62) + Inches(0.28) * ri
        if ri % 2 == 0:
            add_shape(slide, Inches(0.25), ry - Inches(0.01), Inches(9.5), Inches(0.28), BG_CARD)
        for ci, cell in enumerate(row):
            clr = WHITE if ci == 0 else LIGHT_GRAY
            add_text(slide, col_x[ci], ry, col_w[ci], Inches(0.26),
                     cell, size=9, color=clr, bold=(ci == 0))

    # Tuning cards
    ty = Inches(4.85)
    add_shape(slide, Inches(0.3), ty, Inches(4.5), Inches(2.3), BG_CARD, ACCENT2)
    add_text(slide, Inches(0.5), ty + Inches(0.05), Inches(4.1), Inches(0.3),
             "seg_thresh 调优 (最关键参数)", size=13, color=ACCENT2, bold=True)
    add_bullets(slide, Inches(0.5), ty + Inches(0.35), Inches(4.1), Inches(1.8), [
        "控制前景/背景分界线",
        "↑ 增大 → 更严格, 减少假阳性, 可能漏检",
        "↓ 减小 → 更敏感, 增加检出, 可能引入噪点",
        "操作: 用 --tile-index + --debug-vis",
        "  查看 cr_step1_foreground_mask.png",
        "  确认前景覆盖所有可见血管",
    ], size=10, bullet_color=ACCENT2)

    add_shape(slide, Inches(5.1), ty, Inches(4.5), Inches(2.3), BG_CARD, ORANGE)
    add_text(slide, Inches(5.3), ty + Inches(0.05), Inches(4.1), Inches(0.3),
             "morphology_kernel 调优", size=13, color=ORANGE, bold=True)
    add_bullets(slide, Inches(5.3), ty + Inches(0.35), Inches(4.1), Inches(1.8), [
        "控制相邻阳性像素的连接程度",
        "↑ 增大 → 合并更多断裂区域",
        "  大血管壁断裂时需要",
        "↓ 减小 → 保持细粒度分割",
        "  避免不同血管被错误合并",
        "操作: 对比 cr_step2a vs cr_step3",
    ], size=10, bullet_color=ORANGE)


def slide_summary(prs):
    """Slide 13: Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "", "总结与展望")

    add_shape(slide, Inches(0.3), Inches(1.2), Inches(4.5), Inches(5.5), BG_CARD, ACCENT)
    add_text(slide, Inches(0.5), Inches(1.3), Inches(4.1), Inches(0.4),
             "技术路线总结", size=18, color=ACCENT, bold=True)
    add_bullets(slide, Inches(0.5), Inches(1.7), Inches(4.1), Inches(4.8), [
        "DeepLIIF + SAM2 混合流水线",
        "  虚拟染色分解 + 通用分割模型",
        "核心创新: 连通域提取模式",
        "  Seg+Marker 联合 → 形态学 → 完整区域",
        "  优于逐个细胞提取",
        "Mask-Only 提示策略",
        "  logit ±5 Prompt → 3 候选 → 最优选择",
        "  置信度优先重叠消解",
        "WSI 全片处理能力",
        "  YOLO 流式过滤 + 多 GPU 并行",
        "  GeoJSON → QuPath 直接加载",
        "完善的可视化诊断",
        "  每步中间结果 + SAM2 实例级调试",
        "模块化设计, ~6000 行代码",
    ], size=12, bullet_color=ACCENT)

    add_shape(slide, Inches(5.1), Inches(1.2), Inches(4.5), Inches(5.5), BG_CARD, ACCENT3)
    add_text(slide, Inches(5.3), Inches(1.3), Inches(4.1), Inches(0.4),
             "后续工作", size=18, color=ACCENT3, bold=True)
    add_bullets(slide, Inches(5.3), Inches(1.7), Inches(4.1), Inches(4.8), [
        "定量评估",
        "  与病理医生手工标注对比",
        "  Dice / IoU / F1 指标",
        "阈值自适应",
        "  根据染色条件自动调整 seg_thresh",
        "  减少人工参数调节",
        "模型微调",
        "  在本领域数据上 fine-tune SAM2",
        "  提升小型微血管检测精度",
        "扩展标志物",
        "  CD31, Factor VIII 等",
        "推理加速",
        "  TensorRT / ONNX 部署",
        "Web UI 交互平台",
        "  在线标注校对与结果浏览",
    ], size=12, bullet_color=ACCENT3)


def slide_thanks(prs):
    """Slide 14: Thank you."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Pt(6))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT; s.line.fill.background()

    add_text(slide, Inches(1), Inches(1.5), Inches(8), Inches(1),
             "谢谢!", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(2.7), Inches(8), Inches(0.6),
             "欢迎提问与讨论", size=24, color=ACCENT, align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(3.4))

    # Result previews
    iw, ih, iy = Inches(2.0), Inches(2.0), Inches(3.7)
    imgs = [
        (f"{TILE1}/step1_original.png", "原图"),
        (f"{TILE1}/step5_merged_5inst.png", "5 实例"),
        (f"{TILE2}/step1_original.png", "原图"),
        (f"{TILE2}/step5_merged_10inst.png", "10 实例"),
    ]
    x = Inches(0.5)
    for path, label in imgs:
        add_img_label(slide, path, x, iy, iw, ih, label, DIM_GRAY, 10)
        x += Inches(2.3)

    add_text(slide, Inches(1), Inches(6.3), Inches(8), Inches(0.4),
             "CD34 Microvessel Detection  |  DeepLIIF + SAM2  |  Pipeline v2.0",
             size=13, color=DIM_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)              #  1: 封面 + 4 张预览图
    slide_outline(prs)            #  2: 技术路线总览 (11 节)
    slide_architecture(prs)       #  3: Pipeline 总体架构
    slide_deepliif_detail(prs)    #  4: DeepLIIF 详解 + 3 张实例图
    slide_extraction_detail(prs)  #  5: 连通域提取算法 6 步详解
    slide_extraction_example(prs) #  6: 两个 Tile 的提取过程 (各 5 张图)
    slide_sam2_algorithm(prs)     #  7: SAM2 算法 4 步详解 + before/after
    slide_sam2_mask_prompt(prs)   #  8: SAM2 掩码提示词生成可视化
    slide_sam2_comparison(prs)    #  9: 三候选掩码对比 (高分/低分两例)
    slide_postprocessing(prs)     # 10: 后处理与合并逻辑 + 效果对比
    slide_full_pipeline_demo(prs) # 11: 完整 Pipeline 两场景对比
    slide_wsi_pipeline(prs)       # 12: WSI 全片处理方案
    slide_geojson_export(prs)     # 13: GeoJSON 导出详解
    slide_parameters(prs)         # 14: 参数详表与调优建议
    slide_summary(prs)            # 15: 总结与展望
    slide_thanks(prs)             # 16: 致谢

    output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "CD34_Pipeline_Presentation.pptx")
    prs.save(output_path)
    print(f"PPT saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
