#!/usr/bin/env python3
"""
生成 Tile 合并流程 PPT — 以示意图为主，尽量少纯文字
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 颜色 ──
BG       = RGBColor(0x1B, 0x1B, 0x2F)
CARD     = RGBColor(0x24, 0x24, 0x3E)
CARD2    = RGBColor(0x1E, 0x1E, 0x35)
BLUE     = RGBColor(0x4E, 0xC5, 0xF1)
ORANGE   = RGBColor(0xFF, 0xA5, 0x00)
GREEN    = RGBColor(0x2E, 0xCC, 0x71)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0xAA, 0xAA, 0xAA)
GRAY2    = RGBColor(0x77, 0x77, 0x77)
RED      = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW   = RGBColor(0xF1, 0xC4, 0x0F)
PURPLE   = RGBColor(0x9B, 0x59, 0xB6)
PINK     = RGBColor(0xE9, 0x1E, 0x63)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x50)
TEAL     = RGBColor(0x00, 0x96, 0x88)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── 基础绘图工具 ──

def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def rect(slide, l, t, w, h, fill, border=None, bw=Pt(1), radius=True):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s

def box(slide, l, t, w, h, txt, sz=18, clr=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = clr
    p.font.bold = bold; p.font.name = font; p.alignment = align
    return tb

def mbox(slide, l, t, w, h, lines, font="Microsoft YaHei"):
    """多行文本: lines = [(text, size, color, bold, align), ...]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, sz, clr, bld, al) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = clr
        p.font.bold = bld; p.font.name = font; p.alignment = al
        p.space_after = Pt(3)
    return tb

def arrow_r(slide, x, y, w=Inches(0.35), h=Inches(0.28), clr=BLUE):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y - h/2, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = clr; a.line.fill.background()

def arrow_d(slide, x, y, w=Inches(0.28), h=Inches(0.35), clr=BLUE):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x - w/2, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = clr; a.line.fill.background()

def circle(slide, x, y, d, fill, txt="", sz=20, txt_clr=BG):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    c.fill.solid(); c.fill.fore_color.rgb = fill; c.line.fill.background()
    if txt:
        box(slide, x, y + Pt(4), d, d - Pt(8), txt, sz, txt_clr, True, PP_ALIGN.CENTER)

def line_h(slide, x, y, w, clr=BLUE, thick=Pt(2)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thick, )
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()


# ── 画一个带标签的 tile 格子 ──

def draw_tile(slide, x, y, w, h, fill, label, label_clr=WHITE, border=None):
    rect(slide, x, y, w, h, fill, border or GRAY2, Pt(1))
    box(slide, x, y + h/2 - Inches(0.15), w, Inches(0.3),
        label, 11, label_clr, False, PP_ALIGN.CENTER)


# ================================================================
#  Slide 1: 封面
# ================================================================
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    box(s, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
        "Tile 合并算法", 48, BLUE, True, PP_ALIGN.CENTER)
    box(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
        "如何将数万个 tile 的分割结果拼成完整的 QuPath 标注",
        20, WHITE, False, PP_ALIGN.CENTER)
    line_h(s, Inches(4), Inches(4.3), Inches(5.3), BLUE, Pt(3))
    box(s, Inches(1), Inches(4.7), Inches(11), Inches(0.5),
        "export_geojson()  ·  tile_reconstruction.py", 15, GRAY, False, PP_ALIGN.CENTER)

    # 三个 pass 图标
    for i, (label, clr) in enumerate([
        ("Pass 1\n提取多边形", BLUE),
        ("Pass 2\nUnion-Find", ORANGE),
        ("Pass 3\n合并导出", GREEN),
    ]):
        cx = Inches(3.2 + i * 2.8)
        rect(s, cx, Inches(5.5), Inches(2.2), Inches(1.0), CARD, clr, Pt(2))
        circle(s, cx + Inches(0.8), Inches(5.2), Inches(0.5), clr, str(i+1), 18)
        box(s, cx, Inches(5.65), Inches(2.2), Inches(0.8),
            label, 13, clr, True, PP_ALIGN.CENTER)


# ================================================================
#  Slide 2: 问题 — 图解
# ================================================================
def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    box(s, Inches(0.8), Inches(0.3), Inches(6), Inches(0.7),
        "问题: 同一细胞被切到多个 Tile", 28, BLUE, True)

    # ── 左: 3x3 tile grid 示意 ──
    gx, gy = Inches(0.8), Inches(1.3)
    tw, th = Inches(1.55), Inches(1.55)
    overlap_px = Inches(0.30)  # 视觉 overlap 宽度
    stride_vis = tw - overlap_px  # 视觉 stride

    tile_colors = [
        [RGBColor(0x1A,0x3A,0x5C), RGBColor(0x1A,0x40,0x66), RGBColor(0x1A,0x46,0x70)],
        [RGBColor(0x20,0x44,0x64), RGBColor(0x20,0x4A,0x6E), RGBColor(0x20,0x50,0x78)],
        [RGBColor(0x26,0x4E,0x6C), RGBColor(0x26,0x54,0x76), RGBColor(0x26,0x5A,0x80)],
    ]

    for r in range(3):
        for c in range(3):
            tx = gx + c * stride_vis
            ty = gy + r * stride_vis
            draw_tile(s, tx, ty, tw, th, tile_colors[r][c],
                      f"Tile({r},{c})", WHITE, RGBColor(0x44,0x66,0x88))

    # 画 overlap 高亮条 (水平方向 — 在列间)
    for r in range(3):
        for c in range(2):
            ox = gx + (c+1) * stride_vis
            oy = gy + r * stride_vis
            rect(s, ox, oy, overlap_px, th, ORANGE, None, radius=False)

    # 画 overlap 高亮条 (垂直方向 — 在行间)
    for r in range(2):
        for c in range(3):
            ox = gx + c * stride_vis
            oy = gy + (r+1) * stride_vis
            rect(s, ox, oy, tw, overlap_px, YELLOW, None, radius=False)

    # 标注
    box(s, gx, gy + 3 * stride_vis + Inches(0.15), Inches(5), Inches(0.3),
        "橙/黄 = overlap 区域 (128px)", 12, ORANGE, True, PP_ALIGN.CENTER)

    # ── 右: zoom in 一个 overlap 区域 ──
    zx = Inches(6.5)

    box(s, zx, Inches(1.0), Inches(6.5), Inches(0.5),
        "Zoom In: 一个细胞跨越两个 Tile", 20, YELLOW, True)

    # Tile A
    ta_x, ta_y = zx + Inches(0.3), Inches(1.8)
    ta_w, ta_h = Inches(2.5), Inches(2.5)
    rect(s, ta_x, ta_y, ta_w, ta_h, RGBColor(0x1A,0x3A,0x5C), BLUE, Pt(2))
    box(s, ta_x, ta_y + Inches(0.05), ta_w, Inches(0.3),
        "Tile A", 14, BLUE, True, PP_ALIGN.CENTER)

    # Tile B
    tb_x = ta_x + ta_w - Inches(0.6)  # overlap
    rect(s, tb_x, ta_y, ta_w, ta_h, RGBColor(0x20,0x44,0x64), ORANGE, Pt(2))
    box(s, tb_x, ta_y + Inches(0.05), ta_w, Inches(0.3),
        "Tile B", 14, ORANGE, True, PP_ALIGN.CENTER)

    # overlap 区域
    ov_x = tb_x
    ov_w = ta_x + ta_w - tb_x
    rect(s, ov_x, ta_y, ov_w, ta_h, YELLOW, None, radius=False)
    box(s, ov_x - Inches(0.1), ta_y + ta_h + Inches(0.05), ov_w + Inches(0.2), Inches(0.25),
        "overlap", 11, YELLOW, True, PP_ALIGN.CENTER)

    # 画一个"细胞"跨越 overlap — 用椭圆
    cell_x = ov_x - Inches(0.4)
    cell_y = ta_y + Inches(1.0)
    cell_w = Inches(1.4)
    cell_h = Inches(0.9)
    cell = s.shapes.add_shape(MSO_SHAPE.OVAL, cell_x, cell_y, cell_w, cell_h)
    cell.fill.solid(); cell.fill.fore_color.rgb = RED
    cell.line.color.rgb = WHITE; cell.line.width = Pt(2)
    box(s, cell_x, cell_y + Inches(0.2), cell_w, Inches(0.5),
        "细胞", 14, WHITE, True, PP_ALIGN.CENTER)

    # 标注
    mbox(s, zx, Inches(4.6), Inches(6.5), Inches(2.5), [
        ("问题:", 18, RED, True, PP_ALIGN.LEFT),
        ("Tile A 中: 这个细胞的实例 ID = 3", 15, BLUE, False, PP_ALIGN.LEFT),
        ("Tile B 中: 同一个细胞的实例 ID = 1", 15, ORANGE, False, PP_ALIGN.LEFT),
        ("", 6, WHITE, False, PP_ALIGN.LEFT),
        ("两个 tile 各自独立推理，ID 不一致!", 15, YELLOW, True, PP_ALIGN.LEFT),
        ("→ 需要合并算法识别它们是同一个细胞", 15, GREEN, True, PP_ALIGN.LEFT),
    ])


# ================================================================
#  Slide 3: 三 Pass 流程总览 — 纯图
# ================================================================
def slide_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    box(s, Inches(0.8), Inches(0.2), Inches(11), Inches(0.6),
        "合并算法: 三 Pass 流程", 32, BLUE, True)

    # 大流程图: 输入 → Pass1 → Pass2 → Pass3 → 输出
    flow_y = Inches(1.2)
    nodes = [
        ("npy_masks/", GRAY,    "每个 tile 的\n实例 mask (.npy)"),
        ("Pass 1",     BLUE,    "逐 tile\n提取轮廓多边形"),
        ("Pass 2",     ORANGE,  "overlap 区域\n像素匹配\nUnion-Find 合并"),
        ("Pass 3",     GREEN,   "unary_union\n多边形融合"),
        (".geojson",   PURPLE,  "QuPath 标注\nCD34+ 检测"),
    ]

    bw = Inches(2.1)
    bh = Inches(1.5)
    gap = Inches(0.45)
    sx = Inches(0.4)

    for i, (title, clr, desc) in enumerate(nodes):
        nx = sx + i * (bw + gap)
        rect(s, nx, flow_y, bw, bh, CARD, clr, Pt(2))
        box(s, nx, flow_y + Inches(0.1), bw, Inches(0.4),
            title, 17, clr, True, PP_ALIGN.CENTER)
        box(s, nx, flow_y + Inches(0.55), bw, Inches(0.9),
            desc, 12, GRAY, False, PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            arrow_r(s, nx + bw + Inches(0.03), flow_y + bh/2)

    # ── 下面画三个 Pass 的关键示意小图 ──
    py = Inches(3.2)
    pw = Inches(3.8)
    ph = Inches(3.8)
    pgap = Inches(0.5)

    # -- Pass 1 示意: tile mask → polygon --
    p1x = Inches(0.4)
    rect(s, p1x, py, pw, ph, CARD, BLUE, Pt(1))
    box(s, p1x, py + Inches(0.05), pw, Inches(0.35),
        "Pass 1: 提取多边形", 15, BLUE, True, PP_ALIGN.CENTER)

    # 画一个 tile mask 示意 (方格 + 彩色实例)
    mx, my = p1x + Inches(0.3), py + Inches(0.6)
    mw, mh = Inches(1.2), Inches(1.2)
    rect(s, mx, my, mw, mh, DARK_BLUE, GRAY2, Pt(1), False)  # mask 背景
    # 两个"实例" 用彩色方块
    rect(s, mx + Inches(0.1), my + Inches(0.2), Inches(0.4), Inches(0.35),
         RED, None, radius=False)
    box(s, mx + Inches(0.1), my + Inches(0.22), Inches(0.4), Inches(0.3),
        "#1", 9, WHITE, True, PP_ALIGN.CENTER)
    rect(s, mx + Inches(0.6), my + Inches(0.6), Inches(0.5), Inches(0.4),
         TEAL, None, radius=False)
    box(s, mx + Inches(0.6), my + Inches(0.62), Inches(0.5), Inches(0.3),
        "#2", 9, WHITE, True, PP_ALIGN.CENTER)
    box(s, mx, my + mh + Inches(0.05), mw, Inches(0.2),
        "tile mask", 10, GRAY, False, PP_ALIGN.CENTER)

    arrow_r(s, mx + mw + Inches(0.1), my + mh/2, Inches(0.3), Inches(0.2))

    # 多边形示意
    px2 = mx + mw + Inches(0.5)
    # 用不规则四边形模拟 polygon
    poly1 = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                px2, my + Inches(0.15), Inches(0.5), Inches(0.4))
    poly1.fill.solid(); poly1.fill.fore_color.rgb = RED
    poly1.line.color.rgb = WHITE; poly1.line.width = Pt(1)

    poly2 = s.shapes.add_shape(MSO_SHAPE.HEXAGON,
                                px2 + Inches(0.1), my + Inches(0.65), Inches(0.55), Inches(0.45))
    poly2.fill.solid(); poly2.fill.fore_color.rgb = TEAL
    poly2.line.color.rgb = WHITE; poly2.line.width = Pt(1)
    box(s, px2 - Inches(0.1), my + mh + Inches(0.05), Inches(0.9), Inches(0.2),
        "Polygon", 10, GRAY, False, PP_ALIGN.CENTER)

    # 下方说明
    mbox(s, p1x + Inches(0.15), py + Inches(2.2), pw - Inches(0.3), Inches(1.5), [
        ("findContours → Polygon", 12, BLUE, True, PP_ALIGN.LEFT),
        ("坐标偏移到全局坐标系", 11, WHITE, False, PP_ALIGN.LEFT),
        ("poly_map[(r,c,id)] = [Poly]", 11, YELLOW, False, PP_ALIGN.LEFT),
    ])

    # -- Pass 2 示意: overlap matching --
    p2x = p1x + pw + pgap
    rect(s, p2x, py, pw, ph, CARD, ORANGE, Pt(1))
    box(s, p2x, py + Inches(0.05), pw, Inches(0.35),
        "Pass 2: Overlap 匹配", 15, ORANGE, True, PP_ALIGN.CENTER)

    # 两个 tile + overlap 示意
    t2a_x = p2x + Inches(0.2)
    t2a_y = py + Inches(0.6)
    t2w = Inches(1.3)
    t2h = Inches(1.0)
    rect(s, t2a_x, t2a_y, t2w, t2h, RGBColor(0x1A,0x3A,0x5C), BLUE, Pt(1), False)
    box(s, t2a_x, t2a_y + Inches(0.02), t2w, Inches(0.2), "Tile A", 10, BLUE, True, PP_ALIGN.CENTER)

    t2b_x = t2a_x + t2w - Inches(0.3)
    rect(s, t2b_x, t2a_y, t2w, t2h, RGBColor(0x20,0x44,0x64), ORANGE, Pt(1), False)
    box(s, t2b_x, t2a_y + Inches(0.02), t2w, Inches(0.2), "Tile B", 10, ORANGE, True, PP_ALIGN.CENTER)

    # overlap 区域高亮
    rect(s, t2b_x, t2a_y, Inches(0.3), t2h, YELLOW, None, radius=False)

    # 细胞跨越 overlap
    cell2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                t2b_x - Inches(0.15), t2a_y + Inches(0.35),
                                Inches(0.6), Inches(0.4))
    cell2.fill.solid(); cell2.fill.fore_color.rgb = RED
    cell2.line.color.rgb = WHITE; cell2.line.width = Pt(2)

    # ID 标注
    box(s, t2a_x + Inches(0.1), t2a_y + Inches(0.55), Inches(0.5), Inches(0.2),
        "ID=3", 10, BLUE, True)
    box(s, t2b_x + Inches(0.35), t2a_y + Inches(0.55), Inches(0.5), Inches(0.2),
        "ID=1", 10, ORANGE, True)

    # Union-Find 合并箭头
    arrow_d(s, t2a_x + t2w, t2a_y + t2h + Inches(0.1), Inches(0.2), Inches(0.3), YELLOW)

    # Union-Find 结果
    uf_y = t2a_y + t2h + Inches(0.5)
    rect(s, p2x + Inches(0.3), uf_y, Inches(3.2), Inches(0.6), CARD2, GREEN, Pt(1))
    box(s, p2x + Inches(0.3), uf_y + Inches(0.05), Inches(3.2), Inches(0.5),
        "Union-Find: (A,3) ∪ (B,1)\n→ 同一个细胞!", 12, GREEN, True, PP_ALIGN.CENTER)

    # 下方说明
    mbox(s, p2x + Inches(0.15), py + Inches(2.5), pw - Inches(0.3), Inches(1.3), [
        ("overlap 区域像素配对", 12, ORANGE, True, PP_ALIGN.LEFT),
        ("右邻 / 下邻 / 对角", 11, WHITE, False, PP_ALIGN.LEFT),
        ("Union-Find 路径压缩", 11, YELLOW, False, PP_ALIGN.LEFT),
    ])

    # -- Pass 3 示意: merge → GeoJSON --
    p3x = p2x + pw + pgap
    rect(s, p3x, py, pw, ph, CARD, GREEN, Pt(1))
    box(s, p3x, py + Inches(0.05), pw, Inches(0.35),
        "Pass 3: 合并导出", 15, GREEN, True, PP_ALIGN.CENTER)

    # "合并前" 两个 polygon
    bef_y = py + Inches(0.6)
    poly_a = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                 p3x + Inches(0.3), bef_y, Inches(0.8), Inches(0.6))
    poly_a.fill.solid(); poly_a.fill.fore_color.rgb = RED
    poly_a.line.color.rgb = WHITE; poly_a.line.width = Pt(1)

    poly_b = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                 p3x + Inches(0.8), bef_y + Inches(0.1), Inches(0.8), Inches(0.6))
    poly_b.fill.solid(); poly_b.fill.fore_color.rgb = RGBColor(0xC0, 0x39, 0x2B)
    poly_b.line.color.rgb = WHITE; poly_b.line.width = Pt(1)

    box(s, p3x + Inches(0.2), bef_y + Inches(0.7), Inches(1.6), Inches(0.2),
        "合并前 (2个片段)", 10, GRAY, False, PP_ALIGN.CENTER)

    # 箭头
    arrow_r(s, p3x + Inches(1.8), bef_y + Inches(0.35), Inches(0.3), Inches(0.2), GREEN)

    # "合并后" 一个完整 polygon
    poly_m = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 p3x + Inches(2.3), bef_y + Inches(0.05),
                                 Inches(1.1), Inches(0.7))
    poly_m.fill.solid(); poly_m.fill.fore_color.rgb = GREEN
    poly_m.line.color.rgb = WHITE; poly_m.line.width = Pt(2)
    box(s, p3x + Inches(2.3), bef_y + Inches(0.15), Inches(1.1), Inches(0.4),
        "完整", 12, WHITE, True, PP_ALIGN.CENTER)

    box(s, p3x + Inches(2.2), bef_y + Inches(0.8), Inches(1.3), Inches(0.2),
        "unary_union", 10, GREEN, True, PP_ALIGN.CENTER)

    # 下方: GeoJSON 输出
    arrow_d(s, p3x + pw/2, bef_y + Inches(1.1), Inches(0.2), Inches(0.25), PURPLE)

    gj_y = bef_y + Inches(1.5)
    rect(s, p3x + Inches(0.5), gj_y, Inches(2.8), Inches(0.8), CARD2, PURPLE, Pt(1))
    mbox(s, p3x + Inches(0.6), gj_y + Inches(0.05), Inches(2.6), Inches(0.7), [
        (".geojson 输出", 13, PURPLE, True, PP_ALIGN.CENTER),
        ("QuPath 直接导入", 11, GRAY, False, PP_ALIGN.CENTER),
        ("classification: CD34+", 10, GREEN, False, PP_ALIGN.CENTER),
    ])


# ================================================================
#  Slide 4: Pass 1 详解 — 图解
# ================================================================
def slide_pass1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    box(s, Inches(0.6), Inches(0.2), Inches(5), Inches(0.6),
        "Pass 1: 逐 Tile 提取多边形", 30, BLUE, True)

    # 流程图: npy → binary → contour → polygon → poly_map
    flow_items = [
        ("tile.npy\n实例 mask", GRAY,  Inches(0.3)),
        ("binary\nmask==id", BLUE,     Inches(2.8)),
        ("findContours", YELLOW,       Inches(5.3)),
        ("Shapely\nPolygon", GREEN,    Inches(7.8)),
        ("poly_map\n{(r,c,id):[P]}", ORANGE, Inches(10.3)),
    ]

    fy = Inches(1.2)
    fw, fh = Inches(2.0), Inches(1.2)

    for i, (label, clr, fx) in enumerate(flow_items):
        rect(s, fx, fy, fw, fh, CARD, clr, Pt(2))
        box(s, fx, fy + Inches(0.15), fw, Inches(0.9),
            label, 14, clr, True, PP_ALIGN.CENTER)
        if i < len(flow_items) - 1:
            arrow_r(s, fx + fw + Inches(0.05), fy + fh/2, Inches(0.25), Inches(0.2))

    # ── 下方: 大的视觉示意 ──

    # 1) 实例 mask 示意
    mx, my = Inches(0.5), Inches(2.8)
    mw, mh = Inches(2.8), Inches(2.8)
    rect(s, mx, my, mw, mh, DARK_BLUE, BLUE, Pt(1), False)
    box(s, mx, my + Inches(0.02), mw, Inches(0.3),
        "tile_mask.npy", 13, BLUE, True, PP_ALIGN.CENTER)
    # 三个彩色实例块
    inst_data = [
        (Inches(0.2), Inches(0.5), Inches(0.8), Inches(0.7), RED,  "#1"),
        (Inches(1.3), Inches(0.8), Inches(0.9), Inches(0.8), TEAL, "#2"),
        (Inches(0.5), Inches(1.8), Inches(1.0), Inches(0.6), PURPLE, "#3"),
    ]
    for dx, dy, dw, dh, clr, lbl in inst_data:
        r = s.shapes.add_shape(MSO_SHAPE.OVAL, mx+dx, my+dy, dw, dh)
        r.fill.solid(); r.fill.fore_color.rgb = clr
        r.line.color.rgb = WHITE; r.line.width = Pt(1)
        box(s, mx+dx, my+dy+dh/2-Inches(0.1), dw, Inches(0.25),
            lbl, 11, WHITE, True, PP_ALIGN.CENTER)

    # 大箭头
    arrow_r(s, mx + mw + Inches(0.2), my + mh/2, Inches(0.5), Inches(0.35))

    # 2) 二值化 + contour
    bx = mx + mw + Inches(0.9)
    rect(s, bx, my, mw, mh, DARK_BLUE, YELLOW, Pt(1), False)
    box(s, bx, my + Inches(0.02), mw, Inches(0.3),
        "findContours 逐实例", 13, YELLOW, True, PP_ALIGN.CENTER)
    # 三个轮廓（空心）
    for dx, dy, dw, dh, clr, lbl in inst_data:
        r = s.shapes.add_shape(MSO_SHAPE.OVAL, bx+dx, my+dy, dw, dh)
        r.fill.background()
        r.line.color.rgb = clr; r.line.width = Pt(3)

    # 大箭头
    arrow_r(s, bx + mw + Inches(0.2), my + mh/2, Inches(0.5), Inches(0.35))

    # 3) 全局坐标 polygon
    gx2 = bx + mw + Inches(0.9)
    rect(s, gx2, my, Inches(5.5), mh, CARD, GREEN, Pt(1))
    box(s, gx2, my + Inches(0.02), Inches(5.5), Inches(0.3),
        "全局坐标系 → poly_map", 13, GREEN, True, PP_ALIGN.CENTER)

    # 画一个大坐标系背景 + 多个 polygon
    coord_x, coord_y = gx2 + Inches(0.2), my + Inches(0.5)
    coord_w, coord_h = Inches(5.1), Inches(2.1)
    rect(s, coord_x, coord_y, coord_w, coord_h,
         RGBColor(0x15,0x15,0x28), GRAY2, Pt(1), False)

    # 模拟几个 tile 的 polygon 散布在全局坐标系中
    polys_vis = [
        (Inches(0.3), Inches(0.3), Inches(0.6), Inches(0.5), RED),
        (Inches(1.2), Inches(0.5), Inches(0.7), Inches(0.6), TEAL),
        (Inches(0.6), Inches(1.2), Inches(0.8), Inches(0.5), PURPLE),
        (Inches(2.5), Inches(0.2), Inches(0.5), Inches(0.4), BLUE),
        (Inches(3.5), Inches(0.8), Inches(0.6), Inches(0.5), ORANGE),
        (Inches(3.0), Inches(1.3), Inches(0.7), Inches(0.5), PINK),
    ]
    for dx, dy, dw, dh, clr in polys_vis:
        p = s.shapes.add_shape(MSO_SHAPE.OVAL, coord_x+dx, coord_y+dy, dw, dh)
        p.fill.background()
        p.line.color.rgb = clr; p.line.width = Pt(2)

    # 标签
    box(s, coord_x, coord_y + coord_h + Inches(0.05), coord_w, Inches(0.25),
        "每个轮廓已转换为全局坐标，存入 poly_map[(row, col, inst_id)]",
        11, GRAY, False, PP_ALIGN.CENTER)

    # 底部注释
    box(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.8),
        "关键: 坐标偏移 coords[i] += (x_offset, y_offset)  ·  "
        "CHAIN_APPROX_NONE 保留所有轮廓点  ·  "
        "make_valid() 修复无效几何",
        12, GRAY2)


# ================================================================
#  Slide 5: Pass 2 详解 — 图解
# ================================================================
def slide_pass2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    box(s, Inches(0.6), Inches(0.2), Inches(8), Inches(0.6),
        "Pass 2: Overlap 像素匹配 + Union-Find", 30, ORANGE, True)

    # ── 上半: 三个方向的邻居检查 ──
    box(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
        "对每个 tile，检查 3 个方向的邻居:", 16, WHITE, True)

    # 画三个 2-tile 示意
    dirs_data = [
        ("右邻 (row, col+1)", "mask_a[:, 384:]", "mask_b[:, :128]", Inches(0.3)),
        ("下邻 (row+1, col)", "mask_a[384:, :]", "mask_b[:128, :]", Inches(4.6)),
        ("右下 (row+1, col+1)", "mask_a[384:, 384:]", "mask_b[:128, :128]", Inches(8.9)),
    ]

    for dir_label, slice_a, slice_b, dx in dirs_data:
        dy = Inches(1.5)
        dw = Inches(3.8)
        dh = Inches(2.2)
        rect(s, dx, dy, dw, dh, CARD, GRAY2, Pt(1))
        box(s, dx, dy + Inches(0.05), dw, Inches(0.25),
            dir_label, 13, ORANGE, True, PP_ALIGN.CENTER)

        # 两个 tile 小方块
        t_w, t_h = Inches(1.1), Inches(1.0)
        is_vertical = "row+1" in dir_label and "col+1" not in dir_label
        is_diag = "row+1" in dir_label and "col+1" in dir_label

        if is_vertical:
            # 上下排列
            ax, ay = dx + Inches(0.5), dy + Inches(0.4)
            bx2, by2 = dx + Inches(0.5), dy + Inches(1.1)
            ov_x, ov_y = ax, by2
            ov_w, ov_h = t_w, Inches(0.3)
        elif is_diag:
            ax, ay = dx + Inches(0.3), dy + Inches(0.35)
            bx2, by2 = dx + Inches(0.9), dy + Inches(0.95)
            ov_x, ov_y = bx2, by2
            ov_w, ov_h = Inches(0.5), Inches(0.4)
        else:
            # 左右排列
            ax, ay = dx + Inches(0.3), dy + Inches(0.5)
            bx2, by2 = dx + Inches(1.1), dy + Inches(0.5)
            ov_x, ov_y = bx2, ay
            ov_w, ov_h = Inches(0.3), t_h

        rect(s, ax, ay, t_w, t_h, RGBColor(0x1A,0x3A,0x5C), BLUE, Pt(1), False)
        box(s, ax, ay + Inches(0.02), t_w, Inches(0.2), "A", 10, BLUE, True, PP_ALIGN.CENTER)

        rect(s, bx2, by2, t_w, t_h, RGBColor(0x20,0x44,0x64), ORANGE, Pt(1), False)
        box(s, bx2, by2 + Inches(0.02), t_w, Inches(0.2), "B", 10, ORANGE, True, PP_ALIGN.CENTER)

        # overlap
        rect(s, ov_x, ov_y, ov_w, ov_h, YELLOW, None, radius=False)

        # slice 标注
        tx = dx + Inches(1.9)
        mbox(s, tx, dy + Inches(0.45), Inches(1.8), Inches(1.5), [
            (slice_a, 9, BLUE, False, PP_ALIGN.LEFT),
            (slice_b, 9, ORANGE, False, PP_ALIGN.LEFT),
            ("", 4, WHITE, False, PP_ALIGN.LEFT),
            ("both_fg = A>0 & B>0", 9, YELLOW, True, PP_ALIGN.LEFT),
            ("→ union(gid_a, gid_b)", 9, GREEN, True, PP_ALIGN.LEFT),
        ], "Consolas")

    # ── 下半: Union-Find 示意 ──
    uf_y = Inches(3.9)
    rect(s, Inches(0.3), uf_y, Inches(12.7), Inches(3.2), CARD, ORANGE, Pt(1))
    box(s, Inches(0.3), uf_y + Inches(0.05), Inches(12.7), Inches(0.35),
        "Union-Find 并查集: 将跨 tile 的同一细胞关联起来", 17, ORANGE, True, PP_ALIGN.CENTER)

    # 画 Union-Find 树状结构示意
    # 合并前: 4 个独立节点
    bf_x = Inches(0.8)
    bf_y = uf_y + Inches(0.7)
    box(s, bf_x, bf_y, Inches(2), Inches(0.3), "合并前:", 14, WHITE, True)

    nodes_before = [
        ("(0,0,3)", BLUE,   Inches(0.6)),
        ("(0,1,1)", ORANGE, Inches(1.8)),
        ("(1,0,2)", TEAL,   Inches(3.0)),
        ("(1,1,1)", PURPLE, Inches(4.2)),
    ]
    nb_y = bf_y + Inches(0.5)
    for label, clr, ndx in nodes_before:
        circle(s, bf_x + ndx, nb_y, Inches(0.45), clr, "", 10)
        box(s, bf_x + ndx - Inches(0.15), nb_y + Inches(0.5), Inches(0.75), Inches(0.3),
            label, 9, clr, False, PP_ALIGN.CENTER, "Consolas")

    # 大箭头
    arrow_r(s, bf_x + Inches(5.2), nb_y + Inches(0.2), Inches(0.5), Inches(0.3), YELLOW)

    # 合并后: 分成两组
    af_x = bf_x + Inches(5.9)
    box(s, af_x, bf_y, Inches(2), Inches(0.3), "合并后:", 14, WHITE, True)

    # Group 1: (0,0,3) 和 (0,1,1) 连接
    g1_root_x = af_x + Inches(0.8)
    g1_root_y = nb_y
    circle(s, g1_root_x, g1_root_y, Inches(0.45), BLUE, "", 10)
    box(s, g1_root_x - Inches(0.2), g1_root_y + Inches(0.5), Inches(0.85), Inches(0.25),
        "(0,0,3)", 9, BLUE, False, PP_ALIGN.CENTER, "Consolas")

    g1_child_x = af_x + Inches(2.0)
    circle(s, g1_child_x, g1_root_y, Inches(0.45), ORANGE, "", 10)
    box(s, g1_child_x - Inches(0.2), g1_root_y + Inches(0.5), Inches(0.85), Inches(0.25),
        "(0,1,1)", 9, ORANGE, False, PP_ALIGN.CENTER, "Consolas")

    # 连线 (用细长矩形模拟)
    line_h(s, g1_root_x + Inches(0.45), g1_root_y + Inches(0.2),
           g1_child_x - g1_root_x - Inches(0.45), GREEN, Pt(3))

    # Group label
    rect(s, af_x + Inches(0.5), g1_root_y + Inches(0.8), Inches(2.2), Inches(0.35),
         CARD2, GREEN, Pt(1))
    box(s, af_x + Inches(0.5), g1_root_y + Inches(0.82), Inches(2.2), Inches(0.3),
        "同一个细胞!", 11, GREEN, True, PP_ALIGN.CENTER)

    # Group 2
    g2_root_x = af_x + Inches(3.5)
    circle(s, g2_root_x, g1_root_y, Inches(0.45), TEAL, "", 10)
    box(s, g2_root_x - Inches(0.2), g1_root_y + Inches(0.5), Inches(0.85), Inches(0.25),
        "(1,0,2)", 9, TEAL, False, PP_ALIGN.CENTER, "Consolas")

    g2_child_x = af_x + Inches(4.7)
    circle(s, g2_child_x, g1_root_y, Inches(0.45), PURPLE, "", 10)
    box(s, g2_child_x - Inches(0.2), g1_root_y + Inches(0.5), Inches(0.85), Inches(0.25),
        "(1,1,1)", 9, PURPLE, False, PP_ALIGN.CENTER, "Consolas")

    line_h(s, g2_root_x + Inches(0.45), g1_root_y + Inches(0.2),
           g2_child_x - g2_root_x - Inches(0.45), GREEN, Pt(3))

    rect(s, af_x + Inches(3.2), g1_root_y + Inches(0.8), Inches(2.2), Inches(0.35),
         CARD2, GREEN, Pt(1))
    box(s, af_x + Inches(3.2), g1_root_y + Inches(0.82), Inches(2.2), Inches(0.3),
        "同一个细胞!", 11, GREEN, True, PP_ALIGN.CENTER)

    # 底部注释
    box(s, Inches(0.5), uf_y + Inches(2.8), Inches(12), Inches(0.3),
        "路径压缩: O(alpha(n)) 近常数时间  ·  小编号做 root 保持稳定",
        12, GRAY2)


# ================================================================
#  Slide 6: Pass 3 详解 — 图解
# ================================================================
def slide_pass3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    box(s, Inches(0.6), Inches(0.2), Inches(8), Inches(0.6),
        "Pass 3: 多边形合并 + GeoJSON 导出", 30, GREEN, True)

    # ── 上: 流程图 ──
    flow_y = Inches(1.1)
    flow_items = [
        ("按 root 分组", GREEN,   "groups[root]\n.extend(polys)",   Inches(0.3)),
        ("unary_union",  YELLOW,  "合并同组内\n所有多边形",         Inches(3.6)),
        ("过滤 + 简化",  ORANGE,  "面积 < min_area\n→ 丢弃",       Inches(6.9)),
        ("GeoJSON",      PURPLE,  "QuPath Feature\nCD34+ 检测",     Inches(10.2)),
    ]
    fw2, fh2 = Inches(2.8), Inches(1.1)
    for i, (label, clr, desc, fx) in enumerate(flow_items):
        rect(s, fx, flow_y, fw2, fh2, CARD, clr, Pt(2))
        box(s, fx, flow_y + Inches(0.05), fw2, Inches(0.35),
            label, 16, clr, True, PP_ALIGN.CENTER)
        box(s, fx, flow_y + Inches(0.4), fw2, Inches(0.6),
            desc, 12, GRAY, False, PP_ALIGN.CENTER)
        if i < len(flow_items) - 1:
            arrow_r(s, fx + fw2 + Inches(0.05), flow_y + fh2/2)

    # ── 中: unary_union 大图 ──
    uy = Inches(2.6)
    rect(s, Inches(0.3), uy, Inches(12.7), Inches(4.5), CARD, GREEN, Pt(1))
    box(s, Inches(0.3), uy + Inches(0.05), Inches(12.7), Inches(0.35),
        "unary_union 效果示意", 17, GREEN, True, PP_ALIGN.CENTER)

    # 合并前: 2 个多边形有重叠
    bef_label_y = uy + Inches(0.5)
    box(s, Inches(0.6), bef_label_y, Inches(4), Inches(0.3),
        "合并前 (来自 Tile A 和 Tile B):", 14, WHITE, True)

    # Polygon A
    pa_x, pa_y = Inches(0.8), uy + Inches(1.0)
    pa = s.shapes.add_shape(MSO_SHAPE.OVAL, pa_x, pa_y, Inches(2.5), Inches(1.8))
    pa.fill.solid(); pa.fill.fore_color.rgb = RGBColor(0x2C, 0x6F, 0xA0)
    pa.line.color.rgb = BLUE; pa.line.width = Pt(3)
    box(s, pa_x + Inches(0.2), pa_y + Inches(0.6), Inches(1.2), Inches(0.5),
        "Poly A\n(Tile A)", 12, WHITE, True, PP_ALIGN.CENTER)

    # Polygon B (重叠)
    pb_x = pa_x + Inches(1.5)
    pb = s.shapes.add_shape(MSO_SHAPE.OVAL, pb_x, pa_y + Inches(0.2),
                             Inches(2.5), Inches(1.8))
    pb.fill.solid(); pb.fill.fore_color.rgb = RGBColor(0x8B, 0x44, 0x13)
    pb.line.color.rgb = ORANGE; pb.line.width = Pt(3)
    box(s, pb_x + Inches(1.0), pa_y + Inches(0.8), Inches(1.2), Inches(0.5),
        "Poly B\n(Tile B)", 12, WHITE, True, PP_ALIGN.CENTER)

    # 重叠区域标注
    box(s, pa_x + Inches(1.5), pa_y + Inches(1.9), Inches(1.5), Inches(0.3),
        "← 重叠区域 →", 11, YELLOW, True, PP_ALIGN.CENTER)

    # 大箭头
    arrow_r(s, Inches(5.0), pa_y + Inches(0.9), Inches(0.8), Inches(0.5), GREEN)

    box(s, Inches(5.0), pa_y + Inches(1.5), Inches(0.8), Inches(0.3),
        "unary_\nunion", 11, GREEN, True, PP_ALIGN.CENTER)

    # 合并后: 一个完整的 polygon
    box(s, Inches(6.2), bef_label_y, Inches(4), Inches(0.3),
        "合并后:", 14, WHITE, True)

    pm = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), pa_y,
                             Inches(3.5), Inches(2.0))
    pm.fill.solid(); pm.fill.fore_color.rgb = RGBColor(0x1B, 0x7E, 0x4E)
    pm.line.color.rgb = GREEN; pm.line.width = Pt(3)
    box(s, Inches(6.8), pa_y + Inches(0.5), Inches(2.8), Inches(0.9),
        "完整的 Polygon\n无缝融合\n保持轮廓精度", 14, WHITE, True, PP_ALIGN.CENTER)

    # GeoJSON 输出示意
    gj_x = Inches(10.5)
    arrow_r(s, Inches(10.2), pa_y + Inches(0.9), Inches(0.3), Inches(0.3), PURPLE)

    rect(s, gj_x, pa_y + Inches(0.2), Inches(2.3), Inches(1.8), CARD2, PURPLE, Pt(2))
    mbox(s, gj_x + Inches(0.1), pa_y + Inches(0.3), Inches(2.1), Inches(1.6), [
        (".geojson", 15, PURPLE, True, PP_ALIGN.CENTER),
        ("", 4, WHITE, False, PP_ALIGN.LEFT),
        ("type: Feature", 10, GRAY, False, PP_ALIGN.LEFT),
        ("geometry: Polygon", 10, GRAY, False, PP_ALIGN.LEFT),
        ("classification:", 10, GRAY, False, PP_ALIGN.LEFT),
        ("  name: CD34+", 10, GREEN, True, PP_ALIGN.LEFT),
        ("  color: [200,50,50]", 10, RED, False, PP_ALIGN.LEFT),
    ], "Consolas")


# ================================================================
#  Slide 7: 整体流程图
# ================================================================
def slide_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
    box(s, Inches(0.6), Inches(0.2), Inches(11), Inches(0.6),
        "WSI Pipeline 完整流程", 30, BLUE, True)

    # ── 上部: 完整 pipeline 横向 ──
    steps = [
        ("WSI",        GRAY,    ".ndpi/.svs"),
        ("切 Tile",    BLUE,    "512x512\noverlap=128"),
        ("YOLO",       YELLOW,  "过滤背景"),
        ("DeepLIIF",   ORANGE,  "细胞分割"),
        ("SAM2",       GREEN,   "实例分割"),
        ("合并",       RED,     "3-Pass"),
        ("GeoJSON",    PURPLE,  "QuPath"),
    ]

    bw3 = Inches(1.5)
    bh3 = Inches(1.3)
    sx3 = Inches(0.2)
    sy3 = Inches(1.0)
    gap3 = Inches(0.15)

    for i, (title, clr, desc) in enumerate(steps):
        nx = sx3 + i * (bw3 + gap3)
        rect(s, nx, sy3, bw3, bh3, CARD, clr, Pt(2))
        box(s, nx, sy3 + Inches(0.1), bw3, Inches(0.4),
            title, 15, clr, True, PP_ALIGN.CENTER)
        box(s, nx, sy3 + Inches(0.55), bw3, Inches(0.6),
            desc, 11, GRAY, False, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow_r(s, nx + bw3 + Inches(0.01), sy3 + bh3/2,
                    Inches(0.13), Inches(0.16))

    # 高亮"合并"步骤
    # 用大括号连到下面
    merge_x = sx3 + 5 * (bw3 + gap3)
    arrow_d(s, merge_x + bw3/2, sy3 + bh3 + Inches(0.1),
            Inches(0.2), Inches(0.3), RED)

    # ── 下部: 合并阶段 zoom in ──
    zy = Inches(3.0)
    rect(s, Inches(0.3), zy, Inches(12.7), Inches(4.2), CARD, RED, Pt(2))
    box(s, Inches(0.3), zy + Inches(0.05), Inches(12.7), Inches(0.35),
        "Tile 合并阶段 (export_geojson) — Zoom In", 17, RED, True, PP_ALIGN.CENTER)

    # 输入
    inp_x, inp_y = Inches(0.6), zy + Inches(0.7)
    inp_w, inp_h = Inches(1.8), Inches(1.3)
    rect(s, inp_x, inp_y, inp_w, inp_h, CARD2, GRAY, Pt(1))
    # 画几个小 npy 文件
    for j in range(4):
        fy2 = inp_y + Inches(0.1 + j * 0.3)
        rect(s, inp_x + Inches(0.15), fy2, Inches(1.5), Inches(0.22),
             DARK_BLUE, GRAY2, Pt(1), False)
        box(s, inp_x + Inches(0.15), fy2, Inches(1.5), Inches(0.22),
            f"tile_{j}_0.npy", 8, GRAY, False, PP_ALIGN.CENTER, "Consolas")
    box(s, inp_x, inp_y + inp_h + Inches(0.02), inp_w, Inches(0.2),
        "npy_masks/", 11, GRAY, True, PP_ALIGN.CENTER)

    arrow_r(s, inp_x + inp_w + Inches(0.05), inp_y + inp_h/2)

    # Pass 1
    p1x = inp_x + inp_w + Inches(0.5)
    p1w = Inches(2.5)
    rect(s, p1x, inp_y, p1w, inp_h, CARD2, BLUE, Pt(2))
    circle(s, p1x + Inches(0.9), inp_y + Inches(0.1), Inches(0.5), BLUE, "1", 16)
    box(s, p1x, inp_y + Inches(0.7), p1w, Inches(0.3),
        "提取多边形", 14, BLUE, True, PP_ALIGN.CENTER)
    box(s, p1x, inp_y + Inches(1.0), p1w, Inches(0.25),
        "findContours → Polygon", 10, GRAY, False, PP_ALIGN.CENTER)

    arrow_r(s, p1x + p1w + Inches(0.05), inp_y + inp_h/2)

    # Pass 2
    p2x = p1x + p1w + Inches(0.5)
    rect(s, p2x, inp_y, p1w, inp_h, CARD2, ORANGE, Pt(2))
    circle(s, p2x + Inches(0.9), inp_y + Inches(0.1), Inches(0.5), ORANGE, "2", 16)
    box(s, p2x, inp_y + Inches(0.7), p1w, Inches(0.3),
        "Union-Find", 14, ORANGE, True, PP_ALIGN.CENTER)
    box(s, p2x, inp_y + Inches(1.0), p1w, Inches(0.25),
        "overlap 像素匹配", 10, GRAY, False, PP_ALIGN.CENTER)

    arrow_r(s, p2x + p1w + Inches(0.05), inp_y + inp_h/2)

    # Pass 3
    p3x = p2x + p1w + Inches(0.5)
    rect(s, p3x, inp_y, p1w, inp_h, CARD2, GREEN, Pt(2))
    circle(s, p3x + Inches(0.9), inp_y + Inches(0.1), Inches(0.5), GREEN, "3", 16)
    box(s, p3x, inp_y + Inches(0.7), p1w, Inches(0.3),
        "合并导出", 14, GREEN, True, PP_ALIGN.CENTER)
    box(s, p3x, inp_y + Inches(1.0), p1w, Inches(0.25),
        "unary_union → GeoJSON", 10, GRAY, False, PP_ALIGN.CENTER)

    # 输出
    arrow_r(s, p3x + p1w + Inches(0.05), inp_y + inp_h/2, Inches(0.2), Inches(0.2), PURPLE)
    out_x = p3x + p1w + Inches(0.35)
    rect(s, out_x, inp_y + Inches(0.2), Inches(1.3), Inches(0.9), CARD2, PURPLE, Pt(2))
    mbox(s, out_x + Inches(0.05), inp_y + Inches(0.25), Inches(1.2), Inches(0.8), [
        (".geojson", 13, PURPLE, True, PP_ALIGN.CENTER),
        ("", 3, WHITE, False, PP_ALIGN.LEFT),
        ("QuPath", 11, GRAY, False, PP_ALIGN.CENTER),
    ])

    # 底部: 关键设计
    ky = zy + Inches(2.4)
    items = [
        ("128px overlap", BLUE, "边界细胞不丢失"),
        ("Union-Find", ORANGE, "O(alpha(n)) 合并"),
        ("不构建全局数组", GREEN, "内存友好"),
        ("直接 QuPath 导入", PURPLE, "CD34+ 检测标注"),
    ]
    for i, (title, clr, desc) in enumerate(items):
        ix = Inches(0.6 + i * 3.15)
        iw = Inches(2.8)
        rect(s, ix, ky, iw, Inches(1.3), CARD2, clr, Pt(1))
        circle(s, ix + Inches(1.05), ky + Inches(0.1), Inches(0.45), clr, "", 10)
        box(s, ix, ky + Inches(0.6), iw, Inches(0.3),
            title, 13, clr, True, PP_ALIGN.CENTER)
        box(s, ix, ky + Inches(0.9), iw, Inches(0.3),
            desc, 11, GRAY, False, PP_ALIGN.CENTER)


# ================================================================
#  Slide 8: Thank You
# ================================================================
def slide_end(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)

    box(s, Inches(1), Inches(2.5), Inches(11), Inches(1),
        "Thank You", 48, BLUE, True, PP_ALIGN.CENTER)

    line_h(s, Inches(4.5), Inches(3.8), Inches(4.3), BLUE, Pt(3))

    box(s, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
        "export_geojson()  ·  cd34_pipeline/io/tile_reconstruction.py",
        16, GRAY, False, PP_ALIGN.CENTER)


# ================================================================
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)     # 1
    slide_problem(prs)   # 2
    slide_overview(prs)  # 3
    slide_pass1(prs)     # 4
    slide_pass2(prs)     # 5
    slide_pass3(prs)     # 6
    slide_pipeline(prs)  # 7
    slide_end(prs)       # 8

    out = "/local1/yangwenjie/sam2/docs/tile_merge_algorithm.pptx"
    prs.save(out)
    print(f"PPT saved: {out}")


if __name__ == "__main__":
    main()
