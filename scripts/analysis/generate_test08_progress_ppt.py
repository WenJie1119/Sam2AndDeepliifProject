"""Generate a progress-report PPT for the CD34 test08 debug output."""

from __future__ import annotations

import csv
import json
import statistics
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
TEST_DIR = ROOT / "debug_output" / "test08"
REGION_DIR = TEST_DIR / "debug_region"
TILE_DIR = TEST_DIR / "debug_vis" / "tile_37_13_4992_14208"
OUT_PPT = ROOT / "docs" / "CD34_test08_project_progress.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(248, 249, 250)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(86, 96, 111)
ACCENT = RGBColor(160, 55, 45)
ACCENT_DARK = RGBColor(112, 35, 30)
BLUE = RGBColor(46, 103, 178)
GREEN = RGBColor(45, 140, 96)
GOLD = RGBColor(183, 122, 36)
LIGHT_RED = RGBColor(252, 236, 232)
LIGHT_BLUE = RGBColor(234, 242, 255)
LIGHT_GREEN = RGBColor(232, 247, 239)
LINE = RGBColor(224, 229, 235)
WHITE = RGBColor(255, 255, 255)

FONT = "Microsoft YaHei"
FONT_EN = "Aptos"


def load_json(path: Path) -> dict:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def load_area_stats() -> dict:
    summary = {}
    with (TEST_DIR / "DC2200155 A3 CD34_summary.csv").open("r", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            summary[row["metric"]] = float(row["value"])

    areas = []
    with (TEST_DIR / "DC2200155 A3 CD34_statistics.csv").open("r", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            areas.append(float(row["area_px2"]))
    areas_sorted = sorted(areas)

    def pct(p: float) -> float:
        if not areas_sorted:
            return 0.0
        k = (len(areas_sorted) - 1) * p / 100
        lo = int(k)
        hi = min(lo + 1, len(areas_sorted) - 1)
        frac = k - lo
        return areas_sorted[lo] * (1 - frac) + areas_sorted[hi] * frac

    return {
        "count": int(summary["count"]),
        "mean": summary["area_mean"],
        "std": summary["area_std"],
        "min": summary["area_min"],
        "max": summary["area_max"],
        "sum": sum(areas),
        "median": statistics.median(areas) if areas else 0,
        "p75": pct(75),
        "p90": pct(90),
        "p95": pct(95),
        "lt5": sum(a < 5 for a in areas),
        "lt20": sum(a < 20 for a in areas),
        "lt100": sum(a < 100 for a in areas),
        "lt500": sum(a < 500 for a in areas),
    }


def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()


def set_text_style(paragraph, size=18, bold=False, color=TEXT, font=FONT):
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def text_box(slide, x, y, w, h, text="", size=18, bold=False, color=TEXT, align=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    p = frame.paragraphs[0]
    p.text = text
    set_text_style(p, size=size, bold=bold, color=color)
    if align:
        p.alignment = align
    return box


def add_title(slide, title, subtitle=None):
    text_box(slide, Inches(0.55), Inches(0.28), Inches(8.8), Inches(0.46), title, 25, True, TEXT)
    if subtitle:
        text_box(slide, Inches(0.58), Inches(0.78), Inches(10.2), Inches(0.28), subtitle, 10.5, False, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.13), Inches(12.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def rounded_box(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def metric_card(slide, x, y, w, h, value, label, fill=WHITE, value_color=ACCENT):
    rounded_box(slide, x, y, w, h, fill=fill, line=LINE)
    text_box(slide, x + Inches(0.15), y + Inches(0.18), w - Inches(0.3), Inches(0.42), value, 22, True, value_color)
    text_box(slide, x + Inches(0.15), y + Inches(0.68), w - Inches(0.3), Inches(0.42), label, 10.5, False, MUTED)


def bullet_list(slide, x, y, w, h, items, size=14, color=TEXT, leading=1.12):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.space_after = Pt(5)
        p.line_spacing = leading
        set_text_style(p, size=size, color=color)
    return box


def labeled_image(slide, img_path, x, y, w, h, label, tempdir, label_color=MUTED):
    prepared = prepare_image(img_path, tempdir)
    add_image_fit(slide, prepared, x, y, w, h)
    text_box(slide, x, y + h + Inches(0.05), w, Inches(0.22), label, 9.5, False, label_color, PP_ALIGN.CENTER)


def prepare_image(path: Path, tempdir: Path, max_px=1900) -> Path:
    img = Image.open(path).convert("RGB")
    img.thumbnail((max_px, max_px), Image.Resampling.LANCZOS)
    out = tempdir / f"{path.stem}.jpg"
    img.save(out, "JPEG", quality=88, optimize=True)
    return out


def add_image_fit(slide, img_path: Path, x, y, w, h):
    img = Image.open(img_path)
    iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        pic_w = w
        pic_h = w / img_ratio
        pic_x = x
        pic_y = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * img_ratio
        pic_x = x + (w - pic_w) / 2
        pic_y = y
    slide.shapes.add_picture(str(img_path), pic_x, pic_y, width=pic_w, height=pic_h)


def pipeline_step(slide, x, y, w, h, title, detail, color):
    rounded_box(slide, x, y, w, h, fill=WHITE, line=color)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    text_box(slide, x + Inches(0.18), y + Inches(0.14), w - Inches(0.28), Inches(0.25), title, 13, True, TEXT)
    text_box(slide, x + Inches(0.18), y + Inches(0.48), w - Inches(0.28), Inches(0.55), detail, 9.5, False, MUTED)


def connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = RGBColor(170, 178, 188)
    line.line.width = Pt(1.2)


def progress_pill(slide, x, y, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.12), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = text
    set_text_style(p, 10.5, True, WHITE)
    p.alignment = PP_ALIGN.CENTER


def add_small_caption(slide, text):
    text_box(slide, Inches(0.62), Inches(7.1), Inches(11.9), Inches(0.22), text, 8.5, False, MUTED, PP_ALIGN.RIGHT)


def build_ppt():
    metadata = load_json(REGION_DIR / "metadata.json")
    region_summary = load_json(REGION_DIR / "region_summary.json")
    stitched_meta = load_json(REGION_DIR / "stitched_deepliif_metadata.json")
    prompt_summary = load_json(TILE_DIR / "step3_34_weighted_prompt_summary.json")
    sam_summary = load_json(TILE_DIR / "step4_weighted_summary.json")
    merge_summary = load_json(TILE_DIR / "step5_06_merge_filter_summary.json")
    area = load_area_stats()
    mpp = metadata["mpp"]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as td:
        tempdir = Path(td)

        # 1. Title
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        text_box(slide, Inches(0.7), Inches(0.78), Inches(10.6), Inches(0.75), "CD34 微血管识别项目进度汇报", 31, True, TEXT)
        text_box(slide, Inches(0.72), Inches(1.6), Inches(10.2), Inches(0.42), "基于 debug_output/test08 的区域验证结果", 17, False, MUTED)
        progress_pill(slide, Inches(0.72), Inches(2.25), "DeepLIIF + SAM2", ACCENT)
        progress_pill(slide, Inches(3.05), Inches(2.25), "weighted-points", BLUE)
        progress_pill(slide, Inches(5.38), Inches(2.25), "GeoJSON 导出", GREEN)
        labeled_image(
            slide,
            REGION_DIR / "05_region_geojson_overlay.png",
            Inches(7.35),
            Inches(0.65),
            Inches(5.15),
            Inches(5.92),
            "test08 最终 GeoJSON 覆盖结果",
            tempdir,
        )
        metric_card(slide, Inches(0.72), Inches(3.05), Inches(1.75), Inches(1.0), str(region_summary["tile_count"]), "验证 tiles", LIGHT_BLUE, BLUE)
        metric_card(slide, Inches(2.65), Inches(3.05), Inches(1.75), Inches(1.0), str(region_summary["final_geojson_features"]), "最终区域", LIGHT_RED, ACCENT)
        metric_card(slide, Inches(4.58), Inches(3.05), Inches(1.75), Inches(1.0), str(region_summary["accepted_stitch_matches_drawn"]), "拼接匹配", LIGHT_GREEN, GREEN)
        bullet_list(
            slide,
            Inches(0.78),
            Inches(4.45),
            Inches(5.7),
            Inches(1.5),
            [
                "已跑通从 WSI tile 到微血管实例 GeoJSON 的闭环。",
                "当前重点从“能生成结果”推进到“边界质量、碎片过滤、跨 tile 合并稳定性”的验证阶段。",
                "test08 已形成可复查的中间图、统计表和性能曲线。",
            ],
            14,
        )
        add_small_caption(slide, "样本：DC2200155 A3 CD34 | 生成日期：2026-07-02")

        # 2. Progress overview
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "当前完成度概览", "以 test08 结果为依据，项目已具备端到端处理与区域级调试能力")
        metric_card(slide, Inches(0.7), Inches(1.55), Inches(2.55), Inches(1.08), "30 / 30", "选中 tile 全部完成 SAM2", LIGHT_BLUE, BLUE)
        metric_card(slide, Inches(3.55), Inches(1.55), Inches(2.55), Inches(1.08), "569 -> 404", "tile 内实例合并到最终区域", LIGHT_RED, ACCENT)
        metric_card(slide, Inches(6.4), Inches(1.55), Inches(2.55), Inches(1.08), "0", "DeepLIIF 拼接空洞像素", LIGHT_GREEN, GREEN)
        metric_card(slide, Inches(9.25), Inches(1.55), Inches(2.55), Inches(1.08), "0.226 um/px", "当前 WSI 像素分辨率", WHITE, GOLD)
        pipeline_step(slide, Inches(0.8), Inches(3.05), Inches(2.15), Inches(1.05), "输入与 ROI", "支持 WSI tile 枚举、ROI 框选、邻域 tile 自动扩展。", BLUE)
        pipeline_step(slide, Inches(3.2), Inches(3.05), Inches(2.15), Inches(1.05), "DeepLIIF", "生成 Seg / Marker / DAPI 等中间通道，用于阳性提示构建。", GREEN)
        pipeline_step(slide, Inches(5.6), Inches(3.05), Inches(2.15), Inches(1.05), "Prompt 构建", "Seg/Marker 加权、DAB 强度过滤、腔隙补点、伪影过滤。", GOLD)
        pipeline_step(slide, Inches(8.0), Inches(3.05), Inches(2.15), Inches(1.05), "SAM2 分割", "weighted-points 模式已作为主路径，输出候选 mask 和评分。", ACCENT)
        pipeline_step(slide, Inches(10.4), Inches(3.05), Inches(2.15), Inches(1.05), "拼接导出", "center-valid 区域拼接、跨 tile 匹配、GeoJSON 导出。", ACCENT_DARK)
        for i in range(4):
            connector(slide, Inches(2.95 + 2.4 * i), Inches(3.58), Inches(3.2 + 2.4 * i), Inches(3.58))
        bullet_list(
            slide,
            Inches(0.9),
            Inches(4.75),
            Inches(11.6),
            Inches(1.35),
            [
                "工程侧：调试输出已覆盖原图、DeepLIIF 中间通道、prompt 热图、SAM2 候选、合并结果和全区域拼接图。",
                "算法侧：当前主线从“逐连通域二值 mask”收敛到 weighted-points，提高了对弱阳性和不连续血管壁的覆盖能力。",
                "验证侧：test08 的结果已经能支持人工逐层复查，为后续参数迭代提供依据。",
            ],
            13.5,
        )

        # 3. Technical route
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "技术路线与数据流", "当前主流程：DeepLIIF 提供免疫组化线索，SAM2 负责边界精细化，center-valid 负责跨 tile 稳定拼接")
        steps = [
            ("WSI / ROI", "读取区域内 tile，并加入一圈邻域 tile 保证边界上下文"),
            ("DeepLIIF", "生成 Seg / Marker / DAPI；提取 CD34 阳性、细胞核和组织结构线索"),
            ("weighted-points", "Seg/Marker 加权 mask + DAB 强度约束 + lumen 点 + artifact 过滤"),
            ("SAM2", "使用 mask input 和阳性点提示，选择最佳候选 mask"),
            ("后处理", "连通域拆分、碎片合并/过滤、tile 内实例生成"),
            ("拼接导出", "center-valid 裁剪、跨 tile overlap matching、GeoJSON 输出"),
        ]
        x0 = Inches(0.7)
        y0 = Inches(1.55)
        box_w = Inches(3.75)
        box_h = Inches(0.82)
        for idx, (name, detail) in enumerate(steps):
            col = idx % 2
            row = idx // 2
            x = x0 + Inches(4.4) * col
            y = y0 + Inches(1.45) * row
            pipeline_step(slide, x, y, box_w, box_h, name, detail, [BLUE, GREEN, GOLD, ACCENT, ACCENT_DARK, RGBColor(96, 80, 160)][idx])
        connector(slide, Inches(4.45), Inches(1.96), Inches(5.1), Inches(1.96))
        connector(slide, Inches(2.58), Inches(2.37), Inches(2.58), Inches(3.0))
        connector(slide, Inches(4.45), Inches(3.41), Inches(5.1), Inches(3.41))
        connector(slide, Inches(2.58), Inches(3.82), Inches(2.58), Inches(4.45))
        connector(slide, Inches(4.45), Inches(4.86), Inches(5.1), Inches(4.86))
        rounded_box(slide, Inches(9.9), Inches(1.58), Inches(2.55), Inches(3.92), LIGHT_RED, LINE)
        text_box(slide, Inches(10.12), Inches(1.82), Inches(2.1), Inches(0.32), "当前关键变化", 14, True, ACCENT)
        bullet_list(
            slide,
            Inches(10.1),
            Inches(2.33),
            Inches(2.08),
            Inches(2.55),
            [
                "主流程统一为 weighted-points。",
                "加入 DAB 强度门控，减少低强度噪声提示。",
                "增加 lumen 相关补点，改善管腔结构召回。",
                "输出完整 debug artifacts，便于定位误差来源。",
            ],
            11.2,
        )

        # 4. Experiment setup
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "test08 验证设置", "本次结果来自一个局部 ROI 区域，重点检查拼接、提示构建和最终 GeoJSON 质量")
        labeled_image(
            slide,
            REGION_DIR / "01_region_original_mosaic.png",
            Inches(0.7),
            Inches(1.42),
            Inches(5.35),
            Inches(5.35),
            "ROI 原图拼接与 tile 网格",
            tempdir,
        )
        x = Inches(6.45)
        y = Inches(1.5)
        rows = [
            ("样本", "DC2200155 A3 CD34"),
            ("ROI bbox(level0)", f"{metadata['debug_bbox_level0'][0]}, {metadata['debug_bbox_level0'][1]} -> {metadata['debug_bbox_level0'][2]}, {metadata['debug_bbox_level0'][3]}"),
            ("像素分辨率", f"{mpp:.4f} um/px"),
            ("tile / overlap / stride", f"{metadata['tile_size']} / {metadata['overlap']} / {metadata['stride']} px"),
            ("选择 tile", f"{metadata['selected_tile_count']} 个：core {metadata['core_tile_count']} + neighbor {metadata['neighbor_tile_count']}"),
            ("prompt mode", metadata["sam_prompt_mode"]),
            ("DeepLIIF 拼接画布", f"{stitched_meta['canvas_size_wh'][0]} x {stitched_meta['canvas_size_wh'][1]} px"),
            ("结果裁剪", metadata["result_clipping"]),
        ]
        rounded_box(slide, x, y, Inches(5.85), Inches(4.85), WHITE, LINE)
        for idx, (k, v) in enumerate(rows):
            yy = y + Inches(0.25 + idx * 0.54)
            text_box(slide, x + Inches(0.25), yy, Inches(1.85), Inches(0.25), k, 11.5, True, MUTED)
            text_box(slide, x + Inches(2.15), yy, Inches(3.45), Inches(0.25), v, 11.5, False, TEXT)
        text_box(
            slide,
            x,
            Inches(6.55),
            Inches(5.85),
            Inches(0.35),
            "说明：neighbor tile 参与上下文和边界拼接，但最终汇总以 ROI / valid 区域为准。",
            10.5,
            False,
            MUTED,
            PP_ALIGN.CENTER,
        )

        # 5. Final visual result
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "区域级结果：已能生成可复查的微血管实例", "红色覆盖为最终 GeoJSON 实例，网格用于检查 tile 边界和跨 tile 一致性")
        labeled_image(
            slide,
            REGION_DIR / "01_region_original_mosaic.png",
            Inches(0.65),
            Inches(1.42),
            Inches(5.7),
            Inches(5.35),
            "原始 ROI mosaic",
            tempdir,
        )
        labeled_image(
            slide,
            REGION_DIR / "05_region_geojson_overlay.png",
            Inches(6.92),
            Inches(1.42),
            Inches(5.7),
            Inches(5.35),
            "最终 GeoJSON overlay",
            tempdir,
        )
        metric_card(slide, Inches(5.66), Inches(2.02), Inches(1.95), Inches(0.86), str(area["count"]), "final features", LIGHT_RED, ACCENT)
        metric_card(slide, Inches(5.66), Inches(3.05), Inches(1.95), Inches(0.86), f"{area['sum'] * mpp * mpp / 1000:.1f}k", "总面积 um2", LIGHT_GREEN, GREEN)
        metric_card(slide, Inches(5.66), Inches(4.08), Inches(1.95), Inches(0.86), str(region_summary["accepted_stitch_matches_drawn"]), "accepted matches", LIGHT_BLUE, BLUE)

        # 6. DeepLIIF intermediate
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "中间结果：DeepLIIF 通道已接入主流程", "DeepLIIF 输出为后续加权 prompt 提供结构、阳性染色和细胞核信息")
        labeled_image(slide, REGION_DIR / "08_stitched_deepliif_seg.png", Inches(0.65), Inches(1.42), Inches(3.65), Inches(4.3), "Seg：组织结构/细胞边界", tempdir)
        labeled_image(slide, REGION_DIR / "09_stitched_deepliif_marker.png", Inches(4.85), Inches(1.42), Inches(3.65), Inches(4.3), "Marker：CD34 阳性线索", tempdir)
        labeled_image(slide, REGION_DIR / "12_stitched_positive_regions.png", Inches(9.05), Inches(1.42), Inches(3.65), Inches(4.3), "Positive regions：阳性区域编号", tempdir)
        bullet_list(
            slide,
            Inches(0.82),
            Inches(6.15),
            Inches(11.6),
            Inches(0.68),
            [
                "当前能够将 DeepLIIF 的区域输出拼成全 ROI 画布，并保持 hole_pixels=0，说明 center-valid 拼接覆盖完整。",
                "后续误差分析可以直接在 Seg / Marker / Positive regions 三层之间定位：是染色线索不足、prompt 构建偏差，还是 SAM2 边界选择问题。",
            ],
            11.2,
        )

        # 7. Prompt and SAM2 case
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "单 tile 调试链路：prompt 到 SAM2 的闭环已打通", "示例 tile_37_13_4992_14208：从原图、加权提示到合并后实例均有可视化输出")
        labeled_image(slide, TILE_DIR / "step1_original.png", Inches(0.65), Inches(1.45), Inches(2.95), Inches(2.95), "原图", tempdir)
        labeled_image(slide, TILE_DIR / "step3_21_weighted_final_overlay.png", Inches(3.85), Inches(1.45), Inches(2.95), Inches(2.95), "weighted prompt overlay", tempdir)
        labeled_image(slide, TILE_DIR / "step5_05_merge_filter_merged_overlay.png", Inches(7.05), Inches(1.45), Inches(2.95), Inches(2.95), "SAM2 merge overlay", tempdir)
        labeled_image(slide, TILE_DIR / "step7_sam2_merge_diff.png", Inches(10.25), Inches(1.45), Inches(2.25), Inches(2.95), "merge diff", tempdir)
        metric_card(slide, Inches(0.8), Inches(5.05), Inches(2.15), Inches(0.86), str(sam_summary["point_count"]), "positive / lumen points", LIGHT_BLUE, BLUE)
        metric_card(slide, Inches(3.25), Inches(5.05), Inches(2.15), Inches(0.86), f"{sam_summary['best_score']:.3f}", "best SAM2 score", LIGHT_GREEN, GREEN)
        metric_card(slide, Inches(5.7), Inches(5.05), Inches(2.15), Inches(0.86), str(merge_summary["final_regions"]), "tile final regions", LIGHT_RED, ACCENT)
        metric_card(slide, Inches(8.15), Inches(5.05), Inches(2.15), Inches(0.86), str(prompt_summary["dab_lumen_accepted_count"]), "accepted lumen points", WHITE, GOLD)
        bullet_list(
            slide,
            Inches(0.9),
            Inches(6.2),
            Inches(11.1),
            Inches(0.52),
            [
                f"DAB 强度过滤启用：marker 阈值 {prompt_summary['marker_thresh']}，DAB min intensity {prompt_summary['dab_min_intensity']}。",
                f"SAM2 产出 3 个候选，当前选择第 {sam_summary['best_idx']} 个候选，面积 {sam_summary['best_area']:,} px2。",
            ],
            10.5,
        )

        # 8. Stitching
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "跨 tile 拼接验证", "center-valid 拼接减少 tile 边缘重复；overlap matching 用于把跨边界实例合并")
        labeled_image(slide, REGION_DIR / "03_region_tile_merged_mosaic.png", Inches(0.65), Inches(1.42), Inches(3.55), Inches(4.42), "tile 内合并结果", tempdir)
        labeled_image(slide, REGION_DIR / "07_region_overlap_matches.png", Inches(4.9), Inches(1.42), Inches(3.55), Inches(4.42), "overlap matches", tempdir)
        labeled_image(slide, REGION_DIR / "06_region_tile_vs_geojson_diff.png", Inches(9.15), Inches(1.42), Inches(3.55), Inches(4.42), "tile vs GeoJSON diff", tempdir)
        metric_card(slide, Inches(1.0), Inches(6.15), Inches(2.1), Inches(0.72), str(region_summary["tile_merged_instances"]), "tile_merged instances", LIGHT_BLUE, BLUE)
        metric_card(slide, Inches(4.25), Inches(6.15), Inches(2.1), Inches(0.72), str(region_summary["final_geojson_features"]), "final GeoJSON features", LIGHT_RED, ACCENT)
        metric_card(slide, Inches(7.5), Inches(6.15), Inches(2.1), Inches(0.72), str(region_summary["accepted_stitch_matches_drawn"]), "accepted matches", LIGHT_GREEN, GREEN)
        metric_card(slide, Inches(10.75), Inches(6.15), Inches(1.3), Inches(0.72), "0", "holes", WHITE, GOLD)

        # 9. Quantitative statistics
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "量化统计：结果数量已稳定输出，但小碎片仍需优化", "面积分布显示当前存在较多极小区域，后处理阈值和伪影过滤是下一阶段重点")
        labeled_image(slide, TEST_DIR / "DC2200155 A3 CD34_area_histogram.png", Inches(0.65), Inches(1.45), Inches(7.5), Inches(2.72), "区域面积分布", tempdir)
        metric_card(slide, Inches(8.55), Inches(1.55), Inches(1.8), Inches(0.78), str(area["count"]), "count", LIGHT_RED, ACCENT)
        metric_card(slide, Inches(10.55), Inches(1.55), Inches(1.8), Inches(0.78), f"{area['mean']:.0f}", "mean px2", WHITE, GOLD)
        metric_card(slide, Inches(8.55), Inches(2.55), Inches(1.8), Inches(0.78), f"{area['median']:.1f}", "median px2", LIGHT_BLUE, BLUE)
        metric_card(slide, Inches(10.55), Inches(2.55), Inches(1.8), Inches(0.78), f"{area['max']:.0f}", "max px2", LIGHT_GREEN, GREEN)
        rounded_box(slide, Inches(0.75), Inches(4.75), Inches(11.75), Inches(1.48), WHITE, LINE)
        text_box(slide, Inches(1.0), Inches(5.0), Inches(2.55), Inches(0.3), f"<5 px2：{area['lt5']} 个", 13, True, ACCENT)
        text_box(slide, Inches(3.7), Inches(5.0), Inches(2.55), Inches(0.3), f"<20 px2：{area['lt20']} 个", 13, True, ACCENT)
        text_box(slide, Inches(6.4), Inches(5.0), Inches(2.55), Inches(0.3), f"<100 px2：{area['lt100']} 个", 13, True, ACCENT)
        text_box(slide, Inches(9.1), Inches(5.0), Inches(2.55), Inches(0.3), f"<500 px2：{area['lt500']} 个", 13, True, ACCENT)
        bullet_list(
            slide,
            Inches(1.0),
            Inches(5.55),
            Inches(10.7),
            Inches(0.48),
            [
                "结论：大区域可以被连续追踪，但大量小面积实例会拉低统计可信度；需要加入面积、形态和边界置信度联合过滤。",
            ],
            10.8,
        )

        # 10. Runtime and outputs
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "工程输出与性能监控", "test08 已输出结果文件、调试图和运行曲线，便于复现实验和定位瓶颈")
        labeled_image(slide, TEST_DIR / "metrics_progress.png", Inches(0.7), Inches(1.45), Inches(3.7), Inches(2.05), "progress", tempdir)
        labeled_image(slide, TEST_DIR / "metrics_speed.png", Inches(4.75), Inches(1.45), Inches(3.7), Inches(2.05), "speed", tempdir)
        labeled_image(slide, TEST_DIR / "metrics_throughput.png", Inches(8.8), Inches(1.45), Inches(3.7), Inches(2.05), "throughput", tempdir)
        rounded_box(slide, Inches(0.85), Inches(4.35), Inches(5.3), Inches(1.75), WHITE, LINE)
        text_box(slide, Inches(1.1), Inches(4.58), Inches(4.8), Inches(0.28), "主要输出", 14, True, TEXT)
        bullet_list(
            slide,
            Inches(1.1),
            Inches(4.98),
            Inches(4.7),
            Inches(0.78),
            [
                "GeoJSON：DC2200155 A3 CD34.geojson",
                "统计表：summary.csv / statistics.csv",
                "区域调试图：01-13 系列 PNG",
                "单 tile 调试链：step1-step7 系列 PNG/JSON",
            ],
            10.8,
        )
        rounded_box(slide, Inches(6.75), Inches(4.35), Inches(5.35), Inches(1.75), WHITE, LINE)
        text_box(slide, Inches(7.0), Inches(4.58), Inches(4.8), Inches(0.28), "当前工程状态", 14, True, TEXT)
        bullet_list(
            slide,
            Inches(7.0),
            Inches(4.98),
            Inches(4.7),
            Inches(0.78),
            [
                "区域级 debug artifacts 已比较完整。",
                "性能曲线已接入，但 test08 规模较小，尚不足以代表全 WSI 吞吐。",
                "下一步需要在多 ROI / 多 WSI 上做批量稳定性评估。",
            ],
            10.8,
        )

        # 11. Issues and next steps
        slide = prs.slides.add_slide(blank)
        add_bg(slide)
        add_title(slide, "当前判断与下一步计划", "项目已进入结果质量迭代阶段，重点从 pipeline 可用性转向统计可信度和泛化验证")
        rounded_box(slide, Inches(0.8), Inches(1.55), Inches(3.55), Inches(4.85), LIGHT_GREEN, LINE)
        text_box(slide, Inches(1.05), Inches(1.85), Inches(3.05), Inches(0.35), "已完成", 17, True, GREEN)
        bullet_list(
            slide,
            Inches(1.05),
            Inches(2.35),
            Inches(2.9),
            Inches(3.35),
            [
                "DeepLIIF + SAM2 端到端流程跑通。",
                "weighted-points 成为主 prompt 路径。",
                "ROI 级 center-valid 拼接和 GeoJSON 导出完成。",
                "单 tile 和区域级 debug 可视化已形成闭环。",
            ],
            12.5,
        )
        rounded_box(slide, Inches(4.9), Inches(1.55), Inches(3.55), Inches(4.85), LIGHT_RED, LINE)
        text_box(slide, Inches(5.15), Inches(1.85), Inches(3.05), Inches(0.35), "待解决问题", 17, True, ACCENT)
        bullet_list(
            slide,
            Inches(5.15),
            Inches(2.35),
            Inches(2.9),
            Inches(3.35),
            [
                "极小碎片数量偏多，需要更严格的后处理规则。",
                "部分血管边界过宽或粘连，需要结合形态学和强度约束细化。",
                "当前只有区域验证，尚缺少多样本统计和人工标注对照。",
                "性能曲线需要在全 WSI 或更大 ROI 上复测。",
            ],
            12.5,
        )
        rounded_box(slide, Inches(9.0), Inches(1.55), Inches(3.55), Inches(4.85), LIGHT_BLUE, LINE)
        text_box(slide, Inches(9.25), Inches(1.85), Inches(3.05), Inches(0.35), "下一步", 17, True, BLUE)
        bullet_list(
            slide,
            Inches(9.25),
            Inches(2.35),
            Inches(2.9),
            Inches(3.35),
            [
                "加入面积、长宽比、边界支持度联合过滤。",
                "建立人工复核样本集，计算召回、误检和边界质量。",
                "扩展到多 ROI / 多 WSI 批量验证。",
                "整理参数配置和自动汇报脚本，固定可复现实验流程。",
            ],
            12.5,
        )
        text_box(slide, Inches(0.9), Inches(6.85), Inches(11.8), Inches(0.28), "阶段结论：当前 pipeline 已具备项目演示和区域级调试条件，但还不应作为最终定量结果直接交付。", 12.5, True, ACCENT_DARK, PP_ALIGN.CENTER)

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPT)
    return OUT_PPT


if __name__ == "__main__":
    out = build_ppt()
    print(out)
