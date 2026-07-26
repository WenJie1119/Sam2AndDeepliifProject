"""Generate a result-first, step-by-step pipeline explanation PPT for test08."""

from __future__ import annotations

import csv
import json
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
TEST_DIR = ROOT / "debug_output" / "test08"
REGION_DIR = TEST_DIR / "debug_region"
TILE_DIR = TEST_DIR / "debug_vis" / "tile_37_13_4992_14208"
OUT_PPT = ROOT / "docs" / "reports" / "CD34_test08_pipeline_logic_explanation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(248, 249, 250)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(92, 102, 117)
LINE = RGBColor(223, 228, 235)
RED = RGBColor(158, 53, 45)
RED_LIGHT = RGBColor(252, 235, 232)
BLUE = RGBColor(43, 101, 176)
BLUE_LIGHT = RGBColor(232, 241, 255)
GREEN = RGBColor(42, 135, 91)
GREEN_LIGHT = RGBColor(229, 246, 237)
GOLD = RGBColor(182, 118, 34)
GOLD_LIGHT = RGBColor(255, 244, 226)
DARK = RGBColor(17, 24, 39)

FONT = "Microsoft YaHei"


def load_json(path: Path) -> dict:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def load_summary() -> dict:
    summary = {}
    with (TEST_DIR / "DC2200155 A3 CD34_summary.csv").open("r", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            summary[row["metric"]] = float(row["value"])
    areas = []
    with (TEST_DIR / "DC2200155 A3 CD34_statistics.csv").open("r", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            areas.append(float(row["area_px2"]))
    areas.sort()
    median = areas[len(areas) // 2] if areas else 0
    return {
        "count": int(summary["count"]),
        "mean": summary["area_mean"],
        "max": summary["area_max"],
        "median": median,
        "lt20": sum(a < 20 for a in areas),
        "lt100": sum(a < 100 for a in areas),
    }


def prepare_image(path: Path, tempdir: Path, max_px: int = 1700) -> Path:
    img = Image.open(path).convert("RGB")
    img.thumbnail((max_px, max_px), Image.Resampling.LANCZOS)
    out = tempdir / f"{path.stem}.jpg"
    img.save(out, "JPEG", quality=88, optimize=True)
    return out


def add_image_fit(slide, path: Path, x, y, w, h, tempdir: Path):
    img_path = prepare_image(path, tempdir)
    img = Image.open(img_path)
    iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        pic_w = w
        pic_h = w / img_ratio
        pic_x = x
        pic_y = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * img_ratio
        pic_x = x + (w - pic_w) / 2
        pic_y = y
    slide.shapes.add_picture(str(img_path), pic_x, pic_y, width=pic_w, height=pic_h)


def set_p_style(p, size=14, bold=False, color=TEXT):
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def bg(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid()
    s.fill.fore_color.rgb = BG
    s.line.fill.background()


def textbox(slide, x, y, w, h, text="", size=14, bold=False, color=TEXT, align=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    p = frame.paragraphs[0]
    p.text = text
    set_p_style(p, size, bold, color)
    if align is not None:
        p.alignment = align
    return box


def box(slide, x, y, w, h, fill=WHITE, line=LINE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(0.8)
    return s


def title(slide, main, sub=""):
    textbox(slide, Inches(0.55), Inches(0.25), Inches(11.9), Inches(0.43), main, 24, True, DARK)
    if sub:
        textbox(slide, Inches(0.58), Inches(0.72), Inches(11.8), Inches(0.28), sub, 10.5, False, MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.08), Inches(12.25), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def bullets(slide, x, y, w, h, items, size=12.2, color=TEXT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(4)
        p.line_spacing = 1.08
        set_p_style(p, size, False, color)
    return tb


def card(slide, x, y, w, h, header, items, fill=WHITE, accent=BLUE):
    box(slide, x, y, w, h, fill, LINE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.07), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    textbox(slide, x + Inches(0.2), y + Inches(0.16), w - Inches(0.35), Inches(0.28), header, 13, True, accent)
    bullets(slide, x + Inches(0.2), y + Inches(0.55), w - Inches(0.35), h - Inches(0.65), items, 10.9)


def metric(slide, x, y, w, h, value, label, color=BLUE, fill=WHITE):
    box(slide, x, y, w, h, fill, LINE)
    textbox(slide, x + Inches(0.12), y + Inches(0.13), w - Inches(0.24), Inches(0.34), value, 18, True, color, PP_ALIGN.CENTER)
    textbox(slide, x + Inches(0.12), y + Inches(0.51), w - Inches(0.24), Inches(0.25), label, 9.5, False, MUTED, PP_ALIGN.CENTER)


def image_panel(slide, path, x, y, w, h, label, tempdir):
    box(slide, x, y, w, h + Inches(0.34), WHITE, LINE)
    add_image_fit(slide, path, x + Inches(0.08), y + Inches(0.08), w - Inches(0.16), h - Inches(0.06), tempdir)
    textbox(slide, x + Inches(0.08), y + h + Inches(0.09), w - Inches(0.16), Inches(0.22), label, 9.3, False, MUTED, PP_ALIGN.CENTER)


def io_bar(slide, y, input_text, output_text):
    box(slide, Inches(0.75), y, Inches(5.45), Inches(0.72), BLUE_LIGHT, LINE)
    textbox(slide, Inches(0.95), y + Inches(0.14), Inches(1.0), Inches(0.22), "输入", 11.5, True, BLUE)
    textbox(slide, Inches(1.75), y + Inches(0.14), Inches(4.15), Inches(0.38), input_text, 10.3, False, TEXT)
    box(slide, Inches(7.1), y, Inches(5.45), Inches(0.72), GREEN_LIGHT, LINE)
    textbox(slide, Inches(7.3), y + Inches(0.14), Inches(1.0), Inches(0.22), "输出", 11.5, True, GREEN)
    textbox(slide, Inches(8.1), y + Inches(0.14), Inches(4.15), Inches(0.38), output_text, 10.3, False, TEXT)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.38), y + Inches(0.2), Inches(0.48), Inches(0.32))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MUTED
    arrow.line.fill.background()


def code_box(slide, x, y, w, h, lines, title_text="代码逻辑"):
    box(slide, x, y, w, h, RGBColor(253, 253, 253), LINE)
    textbox(slide, x + Inches(0.18), y + Inches(0.14), w - Inches(0.35), Inches(0.25), title_text, 12.5, True, DARK)
    content = "\n".join(lines)
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.46), w - Inches(0.35), h - Inches(0.58))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.name = "Consolas"
    p.font.size = Pt(9.2)
    p.font.color.rgb = RGBColor(45, 55, 72)
    return tb


def footer(slide, text="debug_output/test08"):
    textbox(slide, Inches(0.65), Inches(7.12), Inches(12.0), Inches(0.2), text, 8.5, False, MUTED, PP_ALIGN.RIGHT)


def build():
    metadata = load_json(REGION_DIR / "metadata.json")
    region_summary = load_json(REGION_DIR / "region_summary.json")
    stitched = load_json(REGION_DIR / "stitched_deepliif_metadata.json")
    prompt = load_json(TILE_DIR / "step3_34_weighted_prompt_summary.json")
    sam = load_json(TILE_DIR / "step4_weighted_summary.json")
    merge = load_json(TILE_DIR / "step5_06_merge_filter_summary.json")
    summary = load_summary()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)

        # 1. Result first
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "当前结果先看结论", "test08 已经能从 WSI 局部 ROI 自动生成 CD34 微血管实例 GeoJSON；下面再拆解每一步的输入、输出和代码逻辑")
        image_panel(s, REGION_DIR / "05_region_geojson_overlay.png", Inches(0.65), Inches(1.34), Inches(7.0), Inches(5.35), "最终结果：GeoJSON 微血管实例覆盖到原图", tmp)
        metric(s, Inches(8.05), Inches(1.52), Inches(1.85), Inches(0.82), str(region_summary["tile_count"]), "处理 tile", BLUE, BLUE_LIGHT)
        metric(s, Inches(10.2), Inches(1.52), Inches(1.85), Inches(0.82), str(summary["count"]), "最终实例", RED, RED_LIGHT)
        metric(s, Inches(8.05), Inches(2.58), Inches(1.85), Inches(0.82), str(region_summary["tile_merged_instances"]), "tile 内实例", GOLD, GOLD_LIGHT)
        metric(s, Inches(10.2), Inches(2.58), Inches(1.85), Inches(0.82), str(region_summary["accepted_stitch_matches_drawn"]), "跨 tile 合并", GREEN, GREEN_LIGHT)
        card(s, Inches(8.0), Inches(3.85), Inches(4.2), Inches(1.45), "当前结果说明", [
            "红色区域是最终导出的微血管实例轮廓。",
            "结果来自 30 个 tile：12 个 core + 18 个 neighbor。",
            "最终输出包含 GeoJSON、统计 CSV、区域级和单 tile 调试图。",
        ], WHITE, RED)
        card(s, Inches(8.0), Inches(5.55), Inches(4.2), Inches(1.15), "现阶段重点", [
            "讲清楚每一步如何把图像变成 prompt、mask 和 GeoJSON。",
            "同时指出小碎片较多，后处理仍需继续优化。",
        ], WHITE, BLUE)
        footer(s)

        # 2. Overall dataflow
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "整体流程：从 ROI tile 到 GeoJSON", "主线代码在 cell/main.py：Producer 负责 DeepLIIF + prompt，Consumer 负责 SAM2 + 后处理")
        y = Inches(1.55)
        steps = [
            ("1 选 tile", "ROI/debug region\n选中原图 tile"),
            ("2 DeepLIIF", "原图 tile -> Seg / Marker / DAPI"),
            ("3 weighted prompt", "Seg/Marker/DAB/DAPI\n-> mask_input + points"),
            ("4 SAM2", "mask_input + points\n-> raw mask"),
            ("5 merge/postprocess", "raw mask\n-> tile 实例 mask"),
            ("6 stitching/export", "tile 实例\n-> GeoJSON + stats"),
        ]
        colors = [BLUE, GREEN, GOLD, RED, RGBColor(107, 79, 169), DARK]
        for i, (name, desc) in enumerate(steps):
            x = Inches(0.68 + i * 2.07)
            box(s, x, y, Inches(1.72), Inches(1.02), WHITE, colors[i])
            textbox(s, x + Inches(0.1), y + Inches(0.14), Inches(1.5), Inches(0.24), name, 12, True, colors[i], PP_ALIGN.CENTER)
            textbox(s, x + Inches(0.1), y + Inches(0.44), Inches(1.5), Inches(0.42), desc, 8.8, False, TEXT, PP_ALIGN.CENTER)
            if i < len(steps) - 1:
                arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.76), y + Inches(0.38), Inches(0.26), Inches(0.22))
                arr.fill.solid()
                arr.fill.fore_color.rgb = MUTED
                arr.line.fill.background()
        image_panel(s, REGION_DIR / "01_region_original_mosaic.png", Inches(0.75), Inches(3.25), Inches(3.0), Inches(2.7), "输入 ROI 原图", tmp)
        image_panel(s, REGION_DIR / "08_stitched_deepliif_seg.png", Inches(3.98), Inches(3.25), Inches(3.0), Inches(2.7), "DeepLIIF Seg", tmp)
        image_panel(s, TILE_DIR / "step3_22_weighted_mask_input_256.png", Inches(7.2), Inches(3.25), Inches(2.4), Inches(2.7), "SAM2 mask_input", tmp)
        image_panel(s, REGION_DIR / "05_region_geojson_overlay.png", Inches(9.85), Inches(3.25), Inches(2.8), Inches(2.7), "最终 GeoJSON", tmp)
        footer(s, "核心编排：cell/main.py Producer._extract_one(), Consumer._run_impl(), PostProcessor.export_geojson()")

        # 3. Step 1 tile selection
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "步骤 1：ROI / debug region 选 tile", "功能：只处理目标区域附近的 tile，并加入邻域 tile 解决边界上下文不足")
        io_bar(s, Inches(1.32), "WSI 路径 + ROI/debug-region 坐标 + tile_size/overlap", "selected_tiles.csv + metadata.json + 后续待处理 tile 列表")
        image_panel(s, REGION_DIR / "01_region_original_mosaic.png", Inches(0.75), Inches(2.35), Inches(5.15), Inches(3.9), "区域原图 mosaic：可看到 tile 网格和 ROI 框", tmp)
        card(s, Inches(6.3), Inches(2.35), Inches(2.85), Inches(1.65), "核心功能", [
            "按 ROI bbox 找到相交 tile。",
            "加入 neighbor_radius=1 的外圈 tile。",
            "记录 row/col/x/y/level0 坐标，后面用于拼接和导出全局坐标。",
        ], WHITE, BLUE)
        code_box(s, Inches(9.45), Inches(2.35), Inches(2.95), Inches(1.65), [
            "enumerate_debug_region_tiles(...)",
            "write selected_tiles.csv",
            "metadata: tile_size, overlap,",
            "          stride, crop_origin"
        ], "相关代码")
        card(s, Inches(6.3), Inches(4.3), Inches(6.1), Inches(1.5), "test08 参数", [
            f"tile_size={metadata['tile_size']}，overlap={metadata['overlap']}，stride={metadata['stride']}。",
            f"选中 {metadata['selected_tile_count']} 个 tile，其中 core={metadata['core_tile_count']}，neighbor={metadata['neighbor_tile_count']}。",
            f"mpp={metadata['mpp']:.4f} um/px，后续 GeoJSON 坐标会回到 level-0 坐标系。",
        ], WHITE, GREEN)
        footer(s, "输入输出文件：debug_region/selected_tiles.csv, debug_region/metadata.json")

        # 4. Step 2 DeepLIIF
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "步骤 2：DeepLIIF 生成结构和染色线索", "功能：把原始 CD34 IHC tile 转成模型后续可用的 Seg、Marker、DAPI 通道")
        io_bar(s, Inches(1.25), "PIL tile 图像 batch", "每个 tile 的 Seg / Marker / DAPI 图像")
        image_panel(s, TILE_DIR / "step1_original.png", Inches(0.75), Inches(2.25), Inches(2.55), Inches(2.7), "输入：原图 tile", tmp)
        image_panel(s, TILE_DIR / "step2_01_deepliif_Seg.png", Inches(3.55), Inches(2.25), Inches(2.55), Inches(2.7), "输出：Seg", tmp)
        image_panel(s, TILE_DIR / "step2_02_deepliif_Marker.png", Inches(6.35), Inches(2.25), Inches(2.55), Inches(2.7), "输出：Marker", tmp)
        image_panel(s, TILE_DIR / "step2_03_deepliif_DAPI.png", Inches(9.15), Inches(2.25), Inches(2.55), Inches(2.7), "输出：DAPI", tmp)
        card(s, Inches(0.9), Inches(5.55), Inches(4.2), Inches(1.0), "每个输出的作用", [
            "Seg：提供组织/细胞边界和结构阳性线索。",
            "Marker：提供 CD34 阳性染色强度线索。",
            "DAPI：辅助识别无核管腔/腔隙区域。",
        ], WHITE, GREEN)
        code_box(s, Inches(5.45), Inches(5.55), Inches(3.2), Inches(1.0), [
            "DeepLIIFProcessor.process_batch(",
            "    tile_pils, batch_size, resolution)",
            "return [{'Seg','Marker','DAPI'}]"
        ], "cell/deepliif.py")
        card(s, Inches(8.95), Inches(5.55), Inches(3.15), Inches(1.0), "关键点", [
            "DeepLIIF 不直接给最终血管实例，只给后续 prompt 构建所需的中间证据。",
            "这些中间图会保存为 debug_vis step2。",
        ], WHITE, GOLD)
        footer(s, "主调用：cell/main.py Producer._extract_one() -> DeepLIIFProcessor.process_batch()")

        # 5. Step 3 positive masks
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "步骤 3：从 Seg / Marker 提取阳性候选", "功能：把 DeepLIIF 的彩色通道变成可计算的阳性区域，为 weighted prompt 打底")
        io_bar(s, Inches(1.2), "Seg RGB + Marker 灰度/颜色强度", "seg_logits + marker_logits + 初始阳性支持区域")
        image_panel(s, TILE_DIR / "step2_03_seg_positive_r_intensity_curve.png", Inches(0.75), Inches(2.18), Inches(3.35), Inches(2.55), "Seg 红色强度分布", tmp)
        image_panel(s, TILE_DIR / "step2_04_marker_nonzero_intensity_curve.png", Inches(4.35), Inches(2.18), Inches(3.35), Inches(2.55), "Marker 非零强度分布", tmp)
        image_panel(s, TILE_DIR / "step3_01_weighted_raw_heatmap.png", Inches(7.95), Inches(2.18), Inches(3.35), Inches(2.55), "Seg/Marker 融合初始热图", tmp)
        card(s, Inches(0.9), Inches(5.25), Inches(3.65), Inches(1.2), "Seg 逻辑", [
            "foreground = (R+B > seg_thresh) 且 G 不太高。",
            "R >= B 认为是阳性方向，B > R 认为偏阴性。",
            "按红色强度映射到 -5..5 logits。",
        ], WHITE, BLUE)
        card(s, Inches(4.85), Inches(5.25), Inches(3.65), Inches(1.2), "Marker 逻辑", [
            "marker_thresh 不固定时，用 two-stage Multi-Otsu 自动估计。",
            "Marker > threshold 的像素按强度映射到 1..5 logits。",
            f"示例 tile 阈值={prompt['marker_thresh']}，来源={prompt['marker_threshold_source']}。",
        ], WHITE, GREEN)
        code_box(s, Inches(8.8), Inches(5.25), Inches(3.55), Inches(1.2), [
            "_seg_logits(seg, config)",
            "_marker_logits(marker, config)",
            "pre_dab_raw = maximum(",
            "    seg_logits, marker_logits)"
        ], "weighted_prompt.py")
        footer(s)

        # 6. Step 4 DAB + lumen + artifact
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "步骤 4：构建 weighted prompt 的核心逻辑", "功能：用原图 DAB 强度约束假阳性，同时补回强阳性和管腔相关提示")
        io_bar(s, Inches(1.14), "初始 logits + 原图 RGB + DAPI", "清洗后的 logits、256x256 mask_input、positive point prompts")
        image_panel(s, TILE_DIR / "step3_02_weighted_dab_intensity.png", Inches(0.65), Inches(2.05), Inches(2.35), Inches(2.25), "DAB 强度", tmp)
        image_panel(s, TILE_DIR / "step3_06_weighted_dab_filter_overlay.png", Inches(3.25), Inches(2.05), Inches(2.35), Inches(2.25), "DAB 过滤 overlay", tmp)
        image_panel(s, TILE_DIR / "step3_16_weighted_artifact_decisions.png", Inches(5.85), Inches(2.05), Inches(2.35), Inches(2.25), "伪影判断", tmp)
        image_panel(s, TILE_DIR / "step3_21_weighted_final_overlay.png", Inches(8.45), Inches(2.05), Inches(2.35), Inches(2.25), "最终 prompt overlay", tmp)
        image_panel(s, TILE_DIR / "step3_22_weighted_mask_input_256.png", Inches(11.05), Inches(2.05), Inches(1.55), Inches(2.25), "mask_input", tmp)
        card(s, Inches(0.75), Inches(4.92), Inches(3.0), Inches(1.48), "DAB 过滤", [
            f"把原图 RGB 转 HED-DAB，并按 {prompt['dab_normalization_percentile']} 分位归一化。",
            f"DAB < {prompt['dab_min_intensity']} 的 prompt 像素被压掉。",
            f"本 tile 移除 {prompt['dab_prompt_removed_px']} px，保留 {prompt['dab_prompt_kept_px']} px。",
        ], WHITE, RED)
        card(s, Inches(4.0), Inches(4.92), Inches(3.0), Inches(1.48), "补强逻辑", [
            "强 DAB 像素可以新增或升级 prompt logit。",
            f"本 tile 新增 {prompt['dab_prompt_added_px']} px，升级 {prompt['dab_prompt_upgraded_px']} px。",
            "DAPI/DAB lumen 检测用于给管腔内部加点或弱填充。",
        ], WHITE, GREEN)
        code_box(s, Inches(7.25), Inches(4.92), Inches(5.05), Inches(1.48), [
            "filtered = _suppress_weak_dab_prompt(...)",
            "raw = _add_strong_dab_prompt_support(...)",
            "artifact_filtered = _suppress_artifacts(...)",
            "cleaned = _suppress_small_fragments(...)",
            "lowres_logits = _max_pool_lowres(final, 256)",
            "points = _strong_positive_points(final, ...)"
        ], "build_weighted_prompt()")
        footer(s, "重点：SAM2 输入不是二值图，而是 -5..5 的加权 logits + 少量强阳性点")

        # 7. Step 5 SAM2
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "步骤 5：SAM2 根据 weighted prompt 精细化边界", "功能：用 mask_input 给 SAM2 一个形状先验，再用 positive points 固定强阳性位置")
        io_bar(s, Inches(1.18), "tile RGB + mask_input(1x256x256 logits) + point_coords/labels", "SAM2 candidates + best mask + scores")
        image_panel(s, TILE_DIR / "step3_22_weighted_mask_input_256.png", Inches(0.75), Inches(2.15), Inches(2.1), Inches(2.2), "输入 mask_input", tmp)
        image_panel(s, TILE_DIR / "step3_33_weighted_positive_points.png", Inches(3.15), Inches(2.15), Inches(2.45), Inches(2.2), "输入 points", tmp)
        image_panel(s, TILE_DIR / "step4_weighted_candidate_0_score_0.9157.png", Inches(5.9), Inches(2.15), Inches(2.2), Inches(2.2), "candidate 0", tmp)
        image_panel(s, TILE_DIR / "step4_weighted_candidate_1_score_0.6561.png", Inches(8.35), Inches(2.15), Inches(2.2), Inches(2.2), "candidate 1", tmp)
        image_panel(s, TILE_DIR / "step4_weighted_candidate_2_score_0.9497.png", Inches(10.8), Inches(2.15), Inches(2.2), Inches(2.2), "candidate 2 / best", tmp)
        card(s, Inches(0.85), Inches(4.95), Inches(3.25), Inches(1.25), "SAM2 调用逻辑", [
            "predictor.set_image(tile_np)。",
            "predict(mask_input=..., multimask_output=True)。",
            "如果有 point_coords，则同时传 point_coords 和 point_labels。",
        ], WHITE, BLUE)
        card(s, Inches(4.35), Inches(4.95), Inches(3.25), Inches(1.25), "候选选择", [
            "SAM2 返回多个候选 masks 和 scores。",
            "当前取 score 最大的候选作为 best_mask。",
            f"示例 best_idx={sam['best_idx']}，best_score={sam['best_score']:.3f}，best_area={sam['best_area']:,} px。",
        ], WHITE, GREEN)
        code_box(s, Inches(7.85), Inches(4.95), Inches(4.35), Inches(1.25), [
            "masks, scores, low_res = predictor.predict(...)",
            "best_idx = argmax(scores)",
            "best_mask = masks[best_idx]",
            "instance_mask[best_mask] = 1"
        ], "cell/sam2.py")
        footer(s)

        # 8. Step 6 postprocess
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "步骤 6：SAM2 raw mask 后处理成 tile 内实例", "功能：把 SAM2 的一个/多个 mask 拆成连通实例，过滤并保存为 tile 级实例 mask")
        io_bar(s, Inches(1.18), "SAM2 raw instance_mask + scores + tile_info", "merged tile mask + tile polygons + debug step5/step7")
        image_panel(s, TILE_DIR / "step5_01_merge_filter_sam2_raw_output.png", Inches(0.75), Inches(2.2), Inches(2.25), Inches(2.25), "SAM2 raw mask", tmp)
        image_panel(s, TILE_DIR / "step5_02_merge_filter_sam2_raw_overlay.png", Inches(3.25), Inches(2.2), Inches(2.5), Inches(2.25), "raw overlay", tmp)
        image_panel(s, TILE_DIR / "step5_05_merge_filter_merged_overlay.png", Inches(6.0), Inches(2.2), Inches(2.5), Inches(2.25), "merged overlay", tmp)
        image_panel(s, TILE_DIR / "step7_sam2_merge_diff.png", Inches(8.75), Inches(2.2), Inches(2.5), Inches(2.25), "raw vs merged diff", tmp)
        metric(s, Inches(11.55), Inches(2.38), Inches(1.0), Inches(0.72), str(merge["total_connected_components"]), "CC", RED, RED_LIGHT)
        metric(s, Inches(11.55), Inches(3.28), Inches(1.0), Inches(0.72), str(merge["final_regions"]), "final", GREEN, GREEN_LIGHT)
        card(s, Inches(0.85), Inches(5.05), Inches(3.55), Inches(1.15), "后处理逻辑", [
            "merge_connected_masks() 对 SAM2 mask 做连通域拆分。",
            "每个连通域成为 tile 内一个实例 id。",
            "当前 min_area=0，所以小碎片也会保留到后面统计。",
        ], WHITE, RED)
        code_box(s, Inches(4.75), Inches(5.05), Inches(3.65), Inches(1.15), [
            "merged = merge_connected_masks(",
            "    sam_mask, scores, min_area)",
            "export_mask = _apply_center_valid_crop(...)",
            "tile_polys = _extract_tile_polygons(...)"
        ], "cell/postprocess.py")
        card(s, Inches(8.75), Inches(5.05), Inches(3.4), Inches(1.15), "为什么要拆分", [
            "SAM2 最佳 mask 可能覆盖多个不相连血管段。",
            "连通域拆分后才能把每个血管段作为独立实例统计和导出。",
        ], WHITE, GOLD)
        footer(s)

        # 9. Step 7 center valid stitching
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "步骤 7：center-valid 拼接和跨 tile 合并", "功能：避免 overlap 重复计数，同时用 overlap 区域判断跨 tile 是否属于同一血管")
        io_bar(s, Inches(1.16), "每个 tile 的 merged mask + row/col/x/y", "全局 polygon groups + stitched debug images")
        image_panel(s, REGION_DIR / "03_region_tile_merged_mosaic.png", Inches(0.65), Inches(2.1), Inches(3.35), Inches(3.05), "tile merged mosaic", tmp)
        image_panel(s, REGION_DIR / "07_region_overlap_matches.png", Inches(4.35), Inches(2.1), Inches(3.35), Inches(3.05), "overlap matches", tmp)
        image_panel(s, REGION_DIR / "06_region_tile_vs_geojson_diff.png", Inches(8.05), Inches(2.1), Inches(3.35), Inches(3.05), "tile vs GeoJSON diff", tmp)
        card(s, Inches(0.85), Inches(5.55), Inches(3.2), Inches(1.0), "center-valid", [
            "每个 tile 只导出中心有效区域。",
            f"overlap={metadata['overlap']}，即四周各裁剪约 {metadata['overlap']//2} px。",
            "边界 tile 没有邻居的一侧不裁剪。",
        ], WHITE, BLUE)
        card(s, Inches(4.35), Inches(5.55), Inches(3.2), Inches(1.0), "跨 tile 合并", [
            "保留未裁剪 mask 作为 overlap 匹配证据。",
            "相邻 tile 的重叠区域有像素对应则 Union-Find 合并。",
            f"test08 接受匹配 {region_summary['accepted_stitch_matches_drawn']} 条。",
        ], WHITE, GREEN)
        code_box(s, Inches(7.85), Inches(5.55), Inches(4.2), Inches(1.0), [
            "_apply_center_valid_crop(...)",
            "export_geojson(..., merge_mode='center-valid')",
            "Pass2: overlap pixel matching",
            "Pass3: unary_union -> GeoJSON"
        ], "postprocess.py / tile_reconstruction.py")
        footer(s)

        # 10. Step 8 GeoJSON stats
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "步骤 8：导出 GeoJSON 和统计结果", "功能：把 tile 内局部坐标转换为 WSI level-0 全局坐标，并输出 QuPath 可读结果")
        io_bar(s, Inches(1.16), "poly_map + masks + crop_origin + level_downsample", "DC2200155 A3 CD34.geojson + summary/statistics.csv + area histogram")
        image_panel(s, REGION_DIR / "05_region_geojson_overlay.png", Inches(0.65), Inches(2.05), Inches(4.65), Inches(3.35), "GeoJSON overlay", tmp)
        image_panel(s, TEST_DIR / "DC2200155 A3 CD34_area_histogram.png", Inches(5.65), Inches(2.05), Inches(6.35), Inches(2.65), "面积分布统计", tmp)
        metric(s, Inches(5.85), Inches(5.25), Inches(1.45), Inches(0.68), str(summary["count"]), "count", RED, RED_LIGHT)
        metric(s, Inches(7.55), Inches(5.25), Inches(1.45), Inches(0.68), f"{summary['mean']:.0f}", "mean px2", GOLD, GOLD_LIGHT)
        metric(s, Inches(9.25), Inches(5.25), Inches(1.45), Inches(0.68), f"{summary['median']:.0f}", "median px2", BLUE, BLUE_LIGHT)
        metric(s, Inches(10.95), Inches(5.25), Inches(1.45), Inches(0.68), f"{summary['max']:.0f}", "max px2", GREEN, GREEN_LIGHT)
        card(s, Inches(0.85), Inches(5.9), Inches(4.45), Inches(0.85), "导出逻辑", [
            "局部坐标先按 tile offset 放回 ROI/crop 坐标，再通过 crop_origin 和 level_downsample 转成 level-0 绝对坐标。",
            "GeoJSON 可直接在 QuPath 中加载查看。",
        ], WHITE, GREEN)
        card(s, Inches(5.65), Inches(6.05), Inches(6.35), Inches(0.7), "结果解释", [
            f"当前 <20 px2 的小区域有 {summary['lt20']} 个，<100 px2 的有 {summary['lt100']} 个，说明后续需要加入更严格的小碎片过滤。",
        ], WHITE, RED)
        footer(s, "输出文件：DC2200155 A3 CD34.geojson / *_summary.csv / *_statistics.csv")

        # 11. Code structure
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "代码结构对应关系", "讲代码时可按这张图定位：每一步的核心函数和输入输出对象")
        rows = [
            ("选 tile", "cell/main.py, cell/utils.py", "ROI/debug region -> tile_info list"),
            ("DeepLIIF", "cell/deepliif.py", "tile PIL batch -> Seg/Marker/DAPI"),
            ("weighted prompt", "cd34_pipeline/sam2_wrapper/weighted_prompt.py", "Seg/Marker/RGB/DAPI -> WeightedPromptResult"),
            ("SAM2", "cell/sam2.py", "BucketItem -> sam_mask, scores"),
            ("后处理", "cell/postprocess.py", "sam_mask -> merged mask + polygons"),
            ("拼接导出", "cd34_pipeline/io/tile_reconstruction.py", "poly_map/masks -> GeoJSON"),
            ("可视化", "cell/debug_vis.py, cell/region_debug.py", "step 图 + region mosaic"),
        ]
        x0 = Inches(0.78)
        y0 = Inches(1.45)
        widths = [Inches(1.8), Inches(4.25), Inches(5.65)]
        headers = ["阶段", "代码位置", "主要输入 -> 输出"]
        for i, h in enumerate(headers):
            box(s, x0 + sum(widths[:i]), y0, widths[i], Inches(0.42), BLUE_LIGHT, LINE)
            textbox(s, x0 + sum(widths[:i]) + Inches(0.08), y0 + Inches(0.1), widths[i] - Inches(0.16), Inches(0.2), h, 10.8, True, BLUE)
        for r, row in enumerate(rows):
            y = y0 + Inches(0.44 + r * 0.62)
            fill = WHITE if r % 2 == 0 else RGBColor(252, 253, 254)
            for c, val in enumerate(row):
                box(s, x0 + sum(widths[:c]), y, widths[c], Inches(0.56), fill, LINE)
                textbox(s, x0 + sum(widths[:c]) + Inches(0.08), y + Inches(0.13), widths[c] - Inches(0.16), Inches(0.25), val, 10.1, c == 0, TEXT)
        code_box(s, Inches(0.95), Inches(6.15), Inches(11.0), Inches(0.72), [
            "Producer._extract_one(): DeepLIIF result -> build_weighted_prompt() -> BucketItem(mask_input, point_coords)",
            "Consumer._run_impl(): bucket items -> SAM2Processor.segment_batch() -> PostProcessor.merge_and_process()",
            "PostProcessor.export_geojson(): in-memory poly_map/masks -> overlap matching -> GeoJSON + statistics"
        ], "主线调用关系")
        footer(s)

        # 12. What matters
        s = prs.slides.add_slide(blank)
        bg(s)
        title(s, "这套 pipeline 的重点逻辑", "汇报时建议把重点放在“怎么把 DeepLIIF 线索变成 SAM2 prompt”和“怎么避免 tile overlap 重复计数”")
        image_panel(s, TILE_DIR / "step3_21_weighted_final_overlay.png", Inches(0.75), Inches(1.55), Inches(3.1), Inches(2.75), "重点 1：weighted prompt", tmp)
        image_panel(s, TILE_DIR / "step5_05_merge_filter_merged_overlay.png", Inches(4.12), Inches(1.55), Inches(3.1), Inches(2.75), "重点 2：实例拆分", tmp)
        image_panel(s, REGION_DIR / "07_region_overlap_matches.png", Inches(7.49), Inches(1.55), Inches(3.1), Inches(2.75), "重点 3：跨 tile 合并", tmp)
        card(s, Inches(0.85), Inches(4.85), Inches(3.3), Inches(1.5), "1. Prompt 是核心", [
            "DeepLIIF 不是最终结果，而是生成 SAM2 prompt 的证据来源。",
            "weighted prompt 使用 -5..5 logits，比二值 mask 更能表达强弱和不确定区域。",
            "DAB/DAPI/lumen 逻辑负责减少误检和补足管腔结构。",
        ], WHITE, RED)
        card(s, Inches(4.35), Inches(4.85), Inches(3.3), Inches(1.5), "2. 后处理决定实例数", [
            "SAM2 best mask 可能包含多个不连续区域。",
            "连通域拆分后形成 tile 内多个实例。",
            "当前小碎片偏多，下一步应加入面积/形态/置信度联合过滤。",
        ], WHITE, GOLD)
        card(s, Inches(7.85), Inches(4.85), Inches(3.3), Inches(1.5), "3. 拼接决定全局可信度", [
            "center-valid 负责避免 overlap 重复导出。",
            "未裁剪 mask 仍用于跨 tile 身份匹配。",
            "最后用 GeoJSON 输出到全局坐标，便于 QuPath 复核。",
        ], WHITE, GREEN)
        footer(s)

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPT)
    return OUT_PPT


if __name__ == "__main__":
    print(build())
